#!/usr/bin/env python3
"""Replace the editable regions of a Yitu Quanjie PPTX template."""

from __future__ import annotations

import argparse
import json
import math
import os
import re
import unicodedata
from pathlib import Path
from typing import Any, Callable

from pptx import Presentation
from pptx.util import Inches


MIN_TABLE_ROW_HEIGHT = int(Inches(0.25))
EMU_PER_INCH = 914400


def clear_runs(text_frame: Any) -> None:
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = ""


def set_first_run(paragraph: Any, text: str) -> None:
    if paragraph.runs:
        paragraph.runs[0].text = text
    else:
        paragraph.add_run().text = text


def replace_text(shape: Any, new_text: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        raise TypeError(f"Shape {getattr(shape, 'name', '<unknown>')} has no text frame")
    clear_runs(shape.text_frame)
    set_first_run(shape.text_frame.paragraphs[0], new_text)


def replace_text_keep_lines(shape: Any, new_text: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        raise TypeError(f"Shape {getattr(shape, 'name', '<unknown>')} has no text frame")
    text_frame = shape.text_frame
    clear_runs(text_frame)
    lines = new_text.split("\n")
    paragraphs = list(text_frame.paragraphs)
    while len(paragraphs) < len(lines):
        paragraphs.append(text_frame.add_paragraph())
    for index, paragraph in enumerate(paragraphs):
        set_first_run(paragraph, lines[index] if index < len(lines) else "")


def replace_cell(cell: Any, new_text: str) -> None:
    clear_runs(cell.text_frame)
    set_first_run(cell.text_frame.paragraphs[0], new_text)


def walk_shapes(shapes: Any, callback: Callable[[Any], None]) -> None:
    for shape in shapes:
        callback(shape)
        if getattr(shape, "shape_type", None) == 6:  # GROUP
            walk_shapes(shape.shapes, callback)


def _legacy_text_capacity(shape: Any) -> float | None:
    original_len = len(getattr(shape, "text", ""))
    if not original_len:
        return None
    return float(max(original_len, int(original_len * 1.05)))


def check_length(shape: Any, new_text: str) -> None:
    capacity = text_capacity(shape)
    required = _weighted_text_length(new_text)
    if capacity is not None and required > capacity:
        raise ValueError(
            f"Replacement for {getattr(shape, 'name', '<unknown>')!r} is too long: "
            f"{len(new_text)} characters ({required:.1f} width units) vs "
            f"{capacity:.1f} capacity"
        )


def _font_size_pt(text_frame: Any, default: float = 14.0) -> float:
    sizes = []
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                sizes.append(float(run.font.size.pt))
    return max(sizes, default=default)


def _weighted_text_length(value: str) -> float:
    total = 0.0
    for char in value:
        if char == "\n":
            continue
        if unicodedata.east_asian_width(char) in {"W", "F"}:
            total += 1.0
        elif char.isspace():
            total += 0.3
        else:
            total += 0.55
    return total


def _text_metrics(
    text_frame: Any,
    value: str,
    width: int | None,
    height: int | None = None,
) -> dict[str, float] | None:
    if not width or not height:
        return None

    font_size = max(_font_size_pt(text_frame), 1.0)
    horizontal_margin = int(text_frame.margin_left or 0) + int(text_frame.margin_right or 0)
    vertical_margin = int(text_frame.margin_top or 0) + int(text_frame.margin_bottom or 0)
    usable_width = max(int(width) - horizontal_margin, int(Inches(0.1)))
    usable_height = max(int(height) - vertical_margin, int(Inches(0.1)))
    line_width_units = max(
        1.0,
        usable_width / EMU_PER_INCH * 72.0 / (font_size * 0.95),
    )
    line_height_pt = font_size * 1.2
    capacity_lines = max(
        1,
        math.floor((usable_height / EMU_PER_INCH * 72.0) / line_height_pt),
    )
    required_lines = sum(
        max(1, math.ceil(_weighted_text_length(line) / line_width_units))
        for line in (str(value or "").splitlines() or [""])
    )
    return {
        "font_size": font_size,
        "line_width_units": line_width_units,
        "capacity_lines": float(capacity_lines),
        "required_lines": float(required_lines),
        "capacity_units": line_width_units * capacity_lines,
    }


def text_capacity(
    shape: Any,
    width: int | None = None,
    height: int | None = None,
) -> float | None:
    """Return physical text capacity; use legacy text length only without geometry."""
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return _legacy_text_capacity(shape)
    if width is None:
        width = getattr(shape, "width", None)
    if height is None:
        height = getattr(shape, "height", None)
    metrics = _text_metrics(text_frame, "", width, height)
    if metrics is not None:
        return metrics["capacity_units"]
    return _legacy_text_capacity(shape)


def _content_row_height(cell: Any, column_width: int, value: str) -> int:
    text_frame = cell.text_frame
    font_size = max(_font_size_pt(text_frame), 1.0)
    horizontal_margin = int(text_frame.margin_left or 0) + int(text_frame.margin_right or 0)
    vertical_margin = int(text_frame.margin_top or 0) + int(text_frame.margin_bottom or 0)
    usable_width = max(int(column_width) - horizontal_margin, int(Inches(0.1)))
    line_width_units = max(
        1.0,
        usable_width / EMU_PER_INCH * 72.0 / (font_size * 0.95),
    )
    lines = sum(
        max(1, math.ceil(_weighted_text_length(line) / line_width_units))
        for line in (str(value or "").splitlines() or [""])
    )
    height_pt = lines * font_size * 1.2 + vertical_margin / 12700
    return max(MIN_TABLE_ROW_HEIGHT, int(Inches(height_pt / 72.0)))


def _row_height(table: Any, row_index: int, replacements: dict[tuple[int, int], str] | None = None) -> int:
    replacements = replacements or {}
    heights = []
    for col_index, column in enumerate(table.columns):
        cell = table.cell(row_index, col_index)
        value = replacements.get((row_index, col_index), cell.text)
        heights.append(_content_row_height(cell, int(column.width), value))
    return max([MIN_TABLE_ROW_HEIGHT, *heights])


def _shape_map(slide: Any) -> dict[str, Any]:
    shapes: dict[str, Any] = {}

    def collect(shape: Any) -> None:
        name = str(getattr(shape, "name", "") or "")
        if name and name not in shapes:
            shapes[name] = shape

    walk_shapes(slide.shapes, collect)
    return shapes


def _table_shapes(slide: Any) -> list[Any]:
    tables: list[Any] = []

    def collect(shape: Any) -> None:
        if getattr(shape, "has_table", False):
            tables.append(shape)

    walk_shapes(slide.shapes, collect)
    return tables


def _parse_cell_key(key: str) -> tuple[int, int] | None:
    match = re.fullmatch(r"cell_(\d+)(?:_(\d+))?", key)
    if not match:
        return None
    if match.group(2) is None:
        digits = match.group(1)
        if len(digits) != 2:
            return None
        return int(digits[0]), int(digits[1])
    return int(match.group(1)), int(match.group(2))


def _output_errors(input_path: str | Path, output_path: str | Path) -> list[str]:
    input_target = Path(input_path)
    output_target = Path(output_path)
    errors: list[str] = []
    if output_target.exists() and output_target.is_dir():
        errors.append("Output path is a directory")
    if output_target.exists() and input_target.exists():
        try:
            if output_target.resolve() == input_target.resolve():
                errors.append("Output path must differ from the input template")
        except OSError:
            pass
    if output_target.parent.exists() and not output_target.parent.is_dir():
        errors.append("Output parent path is not a directory")
    return errors


def _validation_report(input_path: str | Path, output_path: str | Path) -> dict[str, Any]:
    return {
        "ok": False,
        "input": {"path": str(input_path), "errors": []},
        "output": {"path": str(output_path), "errors": _output_errors(input_path, output_path)},
        "missing_shapes": [],
        "missing_mapped_shapes": [],
        "shape_errors": [],
        "table_errors": [],
        "overflows": [],
        "table_layout": [],
        "minimum_row_height": MIN_TABLE_ROW_HEIGHT,
        "errors": [],
    }


def validate_replacement(
    input_path: str | Path,
    output_path: str | Path,
    content_map: dict[str, str],
    table_data: dict[str, str] | None = None,
) -> dict[str, Any]:
    """Validate all mappings without changing or saving the template."""
    report = _validation_report(input_path, output_path)
    input_target = Path(input_path)
    if not input_target.is_file():
        report["input"]["errors"].append("Input template does not exist or is not a file")
        report["errors"].extend(report["input"]["errors"])
        report["errors"].extend(report["output"]["errors"])
        return report

    try:
        prs = Presentation(str(input_target))
    except Exception as exc:
        report["input"]["errors"].append(f"Unable to open template: {exc}")
        report["errors"].extend(report["input"]["errors"])
        report["errors"].extend(report["output"]["errors"])
        return report
    if not prs.slides:
        report["input"]["errors"].append("The template has no slides")
        report["errors"].extend(report["input"]["errors"])
        report["errors"].extend(report["output"]["errors"])
        return report

    slide = prs.slides[0]
    shapes = _shape_map(slide)
    for name, new_text in content_map.items():
        shape = shapes.get(name)
        if shape is None:
            report["missing_shapes"].append(name)
            continue
        if not getattr(shape, "has_text_frame", False):
            report["shape_errors"].append({"shape": name, "error": "Shape has no text frame"})
            continue
        metrics = _text_metrics(shape.text_frame, new_text, shape.width, shape.height)
        capacity = text_capacity(shape)
        required = _weighted_text_length(new_text)
        if capacity is not None and required > capacity:
            report["overflows"].append({
                "target": name,
                "kind": "shape",
                "length": len(new_text),
                "weighted_length": required,
                "capacity": capacity,
                "required_lines": metrics["required_lines"] if metrics else None,
                "capacity_lines": metrics["capacity_lines"] if metrics else None,
            })

    tables = _table_shapes(slide)
    target_table = tables[0] if table_data and tables else None

    replacement_cells: dict[tuple[int, int], str] = {}
    if table_data:
        for key, new_text in table_data.items():
            coordinate = _parse_cell_key(str(key))
            if coordinate is None:
                report["table_errors"].append({"key": key, "error": "Invalid table coordinate"})
                continue
            if target_table is None:
                report["table_errors"].append({"key": key, "error": "No table on the first slide"})
                continue
            row_index, col_index = coordinate
            table = target_table.table
            if row_index >= len(table.rows) or col_index >= len(table.columns):
                report["table_errors"].append({
                    "key": key,
                    "row": row_index,
                    "col": col_index,
                    "error": "Table coordinate is out of bounds",
                })
                continue
            replacement_cells[(row_index, col_index)] = new_text
            cell = table.cell(row_index, col_index)
            column_width = int(table.columns[col_index].width)
            row_height = int(table.rows[row_index].height)
            capacity = text_capacity(cell, column_width, row_height)
            required = _weighted_text_length(new_text)
            metrics = _text_metrics(cell.text_frame, new_text, column_width, row_height)
            if capacity is not None and required > capacity:
                report["overflows"].append({
                    "target": key,
                    "kind": "table_cell",
                    "length": len(new_text),
                    "weighted_length": required,
                    "capacity": capacity,
                    "required_lines": metrics["required_lines"] if metrics else None,
                    "capacity_lines": metrics["capacity_lines"] if metrics else None,
                })

    target_tables = [target_table] if target_table is not None else []
    for table_index, table_shape in enumerate(target_tables):
        table = table_shape.table
        required_height = 0
        for row_index, row in enumerate(table.rows):
            row_values = [
                replacement_cells.get(
                    (row_index, col_index),
                    table.cell(row_index, col_index).text,
                )
                for col_index in range(len(table.columns))
            ]
            replacements = {
                (row_index, col_index): value
                for col_index, value in enumerate(row_values)
            }
            required_height += _row_height(table, row_index, replacements)
        if required_height > int(table_shape.height):
            report["overflows"].append({
                "target": f"table_{table_index}",
                "kind": "table_height",
                "required_height": required_height,
                "capacity": int(table_shape.height),
            })
        report["table_layout"].append({
            "table": table_index,
            "minimum_row_height": MIN_TABLE_ROW_HEIGHT,
            "required_height": required_height,
        })

    report["missing_mapped_shapes"] = list(report["missing_shapes"])
    report["errors"].extend(report["input"]["errors"])
    report["errors"].extend(report["output"]["errors"])
    report["errors"].extend(report["missing_shapes"])
    report["errors"].extend(item["error"] for item in report["shape_errors"])
    report["errors"].extend(item["error"] for item in report["table_errors"])
    report["errors"].extend(item["target"] for item in report["overflows"])
    report["ok"] = not report["errors"]
    return report


def run_replacement(
    input_path: str | Path,
    output_path: str | Path,
    content_map: dict[str, str],
    table_data: dict[str, str] | None = None,
) -> Path:
    """Fill the first slide while preserving the existing run and table styles."""
    report = validate_replacement(input_path, output_path, content_map, table_data)
    if not report["ok"]:
        if report["missing_shapes"]:
            missing = ", ".join(report["missing_shapes"])
            raise KeyError(f"Mapped shapes were not found on the first slide: {missing}")
        if report["overflows"]:
            overflow = report["overflows"][0]
            if "length" in overflow and "capacity" in overflow:
                raise ValueError(
                    f"Replacement for {overflow['target']!r} is too long: "
                    f"{overflow['length']} characters vs {overflow['capacity']} capacity"
                )
            raise ValueError(
                f"Yitu replacement overflow for {overflow['target']!r}: "
                f"required {overflow.get('required_height')} vs capacity {overflow.get('capacity')}"
            )
        raise ValueError("Yitu replacement validation failed: " + "; ".join(report["errors"]))

    prs = Presentation(str(input_path))
    if not prs.slides:
        raise ValueError("The template has no slides")
    slide = prs.slides[0]
    target_tables = _table_shapes(slide)
    target_table = target_tables[0] if table_data and target_tables else None

    def process(shape: Any) -> None:
        name = getattr(shape, "name", "")
        if name in content_map:
            new_text = content_map[name]
            if "\n" in new_text:
                replace_text_keep_lines(shape, new_text)
            else:
                replace_text(shape, new_text)

    walk_shapes(slide.shapes, process)

    if target_table is not None:
        table = target_table.table
        for key, new_text in table_data.items():
            coordinate = _parse_cell_key(str(key))
            if coordinate is None:
                continue
            row_index, col_index = coordinate
            if row_index < len(table.rows) and col_index < len(table.columns):
                replace_cell(table.cell(row_index, col_index), new_text)
        for row_index, row in enumerate(table.rows):
            row.height = _row_height(table, row_index)

    target = Path(output_path)
    target.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(target))
    return target


def load_json(path: str | Path) -> dict[str, str]:
    with Path(path).open("r", encoding="utf-8-sig") as handle:
        value = json.load(handle)
    if not isinstance(value, dict):
        raise ValueError(f"Expected a JSON object: {path}")
    return {str(key): str(item) for key, item in value.items()}


def main() -> int:
    parser = argparse.ArgumentParser(description="Fill a Yitu Quanjie PPTX template.")
    parser.add_argument("--template", required=True, help="Input PPTX template")
    parser.add_argument("--output", required=True, help="Output PPTX path")
    parser.add_argument("--content-map", required=True, help="JSON object mapping shape names to text")
    parser.add_argument("--table-data", help="Optional JSON object for cell_00/cell_01/cell_10/cell_11")
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Validate mappings and print JSON without saving an output",
    )
    args = parser.parse_args()
    content_map = load_json(args.content_map)
    table_data = load_json(args.table_data) if args.table_data else None
    if args.dry_run:
        report = validate_replacement(args.template, args.output, content_map, table_data)
        print(json.dumps(report, ensure_ascii=False, indent=2))
        return 0

    output = run_replacement(args.template, args.output, content_map, table_data)
    print(f"Created: {output} ({os.path.getsize(output):,} bytes)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
