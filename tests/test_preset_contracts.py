import json
import importlib
import importlib.util
import io
import os
import sys
import tempfile
import unittest
import zipfile
from contextlib import redirect_stdout
from pathlib import Path
from argparse import Namespace
from unittest.mock import patch

from pptx import Presentation
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
SCRIPTS_DIR = REPO_ROOT / "ppt-template-batch" / "scripts"
if str(SCRIPTS_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPTS_DIR))

class PresetContractTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        super().setUpClass()
        cls.fill_buyer_briefing_pages = importlib.import_module("fill_buyer_briefing_pages")
        cls.fill_ppt_from_records = importlib.import_module("fill_ppt_from_records")
        cls.generate_layout_config = importlib.import_module("generate_layout_config")

        pipeline_path = REPO_ROOT / "scripts" / "run_ppt_batch_pipeline.py"
        pipeline_spec = importlib.util.spec_from_file_location("run_ppt_batch_pipeline", pipeline_path)
        if not pipeline_spec or not pipeline_spec.loader:
            raise ImportError(f"Unable to load {pipeline_path}")
        cls.run_ppt_batch_pipeline = importlib.util.module_from_spec(pipeline_spec)
        pipeline_spec.loader.exec_module(cls.run_ppt_batch_pipeline)

        builder_path = REPO_ROOT / "scripts" / "build_feishu_agent_skill.py"
        builder_spec = importlib.util.spec_from_file_location("build_feishu_agent_skill", builder_path)
        if not builder_spec or not builder_spec.loader:
            raise ImportError(f"Unable to load {builder_path}")
        cls.build_feishu_agent_skill = importlib.util.module_from_spec(builder_spec)
        builder_spec.loader.exec_module(cls.build_feishu_agent_skill)

        console_path = REPO_ROOT / "ppt-template-batch" / "scripts" / "control_console.py"
        console_spec = importlib.util.spec_from_file_location("control_console", console_path)
        if not console_spec or not console_spec.loader:
            raise ImportError(f"Unable to load {console_path}")
        cls.control_console = importlib.util.module_from_spec(console_spec)
        console_spec.loader.exec_module(cls.control_console)

        yitu_path = REPO_ROOT / "yitu-quanjie" / "scripts" / "yitu_quanjie_replace.py"
        yitu_spec = importlib.util.spec_from_file_location("yitu_quanjie_replace", yitu_path)
        if not yitu_spec or not yitu_spec.loader:
            raise ImportError(f"Unable to load {yitu_path}")
        cls.yitu_quanjie_replace = importlib.util.module_from_spec(yitu_spec)
        yitu_spec.loader.exec_module(cls.yitu_quanjie_replace)

    def test_generic_report_contains_shared_quality_keys(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template = root / "template.pptx"
            records = root / "records.json"
            config = root / "layout-config.json"
            output = root / "output.pptx"
            Presentation().save(template)
            records.write_text(json.dumps({"records": []}), encoding="utf-8")
            config.write_text(json.dumps({"required_fields": [], "slides": []}), encoding="utf-8")

            report = self.fill_ppt_from_records.fill_presentation(template, records, config, output, root / "workspace")

            self.assertIs(type(report["ok"]), bool)
            self.assertIs(report["ok"], True)
            self.assertIsInstance(report["missing_required_fields"], list)
            self.assertEqual(report["missing_required_fields"], [])
            self.assertIsInstance(report["missing_assets"], list)
            self.assertEqual(report["missing_assets"], [])
            self.assertIsInstance(report["warnings"], list)
            self.assertEqual(report["warnings"], [])
            self.assertEqual(report["schema_version"], 2)
            self.assertIsInstance(report["stale_template_text"], list)
            self.assertIsInstance(report["capacity_warnings"], list)
            self.assertEqual(report["expected_slide_count"], 0)
            self.assertIs(report["reopen_ok"], True)
            self.assertEqual(report["reopen_status"]["status"], "ok")

    def test_generic_selector_precedes_numeric_shape_index(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template = root / "template.pptx"
            records = root / "records.json"
            config = root / "layout-config.json"
            output = root / "output.pptx"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            numeric_target = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
            numeric_target.name = "Numeric fallback"
            numeric_target.text = "wrong target"
            stable_target = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
            stable_target.name = "Stable title"
            stable_target.text = "template title"
            presentation.save(template)

            records.write_text(json.dumps({"records": [{"name": "Resolved title"}]}), encoding="utf-8")
            config.write_text(
                json.dumps(
                    {
                        "version": 1,
                        "required_fields": ["name"],
                        "slides": [
                            {
                                "slide_index": 1,
                                "record_index": 1,
                                "texts": [
                                    {
                                        "selector": {"name": "Stable title", "role": "text"},
                                        "shape_index": 1,
                                        "field": "name",
                                    }
                                ],
                            }
                        ],
                    }
                ),
                encoding="utf-8",
            )

            report = self.fill_ppt_from_records.fill_presentation(
                template, records, config, output, root / "workspace"
            )

            self.assertIs(report["ok"], True)
            result = Presentation(output)
            self.assertEqual(result.slides[0].shapes[0].text, "wrong target")
            self.assertEqual(result.slides[0].shapes[1].text, "Resolved title")

    def test_selector_role_text_matches_real_placeholder_shape(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template = root / "template.pptx"
            records = root / "records.json"
            config = root / "layout-config.json"
            output = root / "output.pptx"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            numeric_target = slide.placeholders[0]
            numeric_target.name = "Numeric fallback"
            numeric_target.text = "wrong target"
            placeholder_target = slide.placeholders[1]
            placeholder_target.name = "Real content placeholder"
            placeholder_target.text = "template content"
            presentation.save(template)

            records.write_text(json.dumps({"records": [{"name": "Resolved placeholder"}]}), encoding="utf-8")
            config.write_text(
                json.dumps(
                    {
                        "schema_version": 2,
                        "required_fields": ["name"],
                        "slides": [
                            {
                                "slide_index": 1,
                                "record_index": 1,
                                "texts": [
                                    {
                                        "selector": {
                                            "name": "Real content placeholder",
                                            "role": "text",
                                        },
                                        "shape_index": 1,
                                        "field": "name",
                                    }
                                ],
                            }
                        ],
                    }
                ),
                encoding="utf-8",
            )

            self.fill_ppt_from_records.fill_presentation(
                template, records, config, output, root / "workspace"
            )

            result = Presentation(output)
            self.assertEqual(result.slides[0].shapes[0].text, "wrong target")
            self.assertEqual(result.slides[0].shapes[1].text, "Resolved placeholder")

    def test_generated_layout_config_uses_schema_v2_and_stable_selectors(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template = root / "template.pptx"
            output = root / "layout-config.json"

            presentation = Presentation()
            cover = presentation.slides.add_slide(presentation.slide_layouts[6])
            title = cover.shapes.add_textbox(Inches(1), Inches(0.2), Inches(4), Inches(1))
            title.name = "Cover title"
            title.text = "Title"
            country = cover.shapes.add_textbox(Inches(1), Inches(1.5), Inches(4), Inches(0.5))
            country.name = "Cover country"
            country.text = "Country: Example"
            content = presentation.slides.add_slide(presentation.slide_layouts[6])
            content_title = content.shapes.add_textbox(Inches(1), Inches(0.2), Inches(4), Inches(0.5))
            content_title.name = "Content title"
            content_title.text = "Content"
            table = content.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(2))
            table.name = "Content table"
            presentation.save(template)

            old_argv = sys.argv
            try:
                sys.argv = [
                    "generate_layout_config.py",
                    "--template",
                    str(template),
                    "--output",
                    str(output),
                ]
                self.assertEqual(self.generate_layout_config.main(), 0)
            finally:
                sys.argv = old_argv

            generated = json.loads(output.read_text(encoding="utf-8"))
            self.assertEqual(generated["schema_version"], 2)
            self.assertEqual(generated["cover"]["title_selector"]["name"], "Cover title")
            self.assertEqual(generated["content"]["table_selector"]["name"], "Content table")

    def test_strict_missing_required_fields_rejects_without_writing_output(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template = root / "template.pptx"
            records = root / "records.json"
            config = root / "layout-config.json"
            output = root / "output.pptx"
            Presentation().save(template)
            records.write_text(json.dumps({"records": [{}]}), encoding="utf-8")
            config.write_text(json.dumps({"required_fields": ["name"], "slides": []}), encoding="utf-8")

            report = self.fill_ppt_from_records.fill_presentation(
                template, records, config, output, root / "workspace", strict=True
            )

            self.assertIs(report["ok"], False)
            self.assertEqual(report["missing_required_fields"], [{"record_index": 1, "field": "name"}])
            self.assertEqual(report["reopen_status"]["status"], "not_run")
            self.assertFalse(output.exists())

    def test_generic_report_identifies_stale_text_and_capacity_warnings(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template = root / "template.pptx"
            records = root / "records.json"
            config = root / "layout-config.json"
            output = root / "output.pptx"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            stale = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
            stale.name = "Unmapped placeholder"
            stale.text = "{{record.name}}"
            overflow = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(0.5), Inches(0.2))
            overflow.name = "Small box"
            overflow.text = "x" * 200
            presentation.save(template)
            records.write_text(json.dumps({"records": []}), encoding="utf-8")
            config.write_text(json.dumps({"schema_version": 2, "slides": []}), encoding="utf-8")

            report = self.fill_ppt_from_records.fill_presentation(
                template, records, config, output, root / "workspace"
            )

            self.assertTrue(report["stale_template_text"])
            self.assertTrue(report["capacity_warnings"])
            self.assertEqual(report["slide_count"], 1)
            self.assertIs(report["reopen_ok"], True)

    def test_pipeline_does_not_reuse_stale_report_after_runner_failure(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            stale_report = root / "fill-report.json"
            stale_report.write_text(
                json.dumps({"ok": True, "record_count": 999}), encoding="utf-8"
            )
            job = {
                "template": str(root / "template.pptx"),
                "records": str(root / "records.json"),
                "layout_config": str(root / "layout-config.json"),
                "output": str(root / "output.pptx"),
                "report": str(stale_report),
            }
            defaults = Namespace(
                template=None,
                records=None,
                layout_config=None,
                output=None,
                output_dir=None,
                workspace=None,
                strict=False,
            )

            with patch.object(
                self.run_ppt_batch_pipeline,
                "run",
                side_effect=RuntimeError("simulated filler failure"),
            ):
                with self.assertRaisesRegex(RuntimeError, "simulated filler failure"):
                    self.run_ppt_batch_pipeline.run_single_job(job, defaults, 1)

            self.assertEqual(
                json.loads(stale_report.read_text(encoding="utf-8")),
                {"ok": True, "record_count": 999},
            )

    def test_non_strict_repeat_skips_failed_records_and_reports_them(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template = root / "template.pptx"
            records = root / "records.json"
            config = root / "layout-config.json"
            output = root / "output.pptx"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            name_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
            name_shape.name = "Record name"
            name_shape.text = "template name"
            presentation.save(template)
            records.write_text(json.dumps({"records": [{"name": "Good"}, {}]}), encoding="utf-8")
            config.write_text(
                json.dumps(
                    {
                        "schema_version": 2,
                        "required_fields": ["name"],
                        "repeat": {
                            "source_slide_index": 1,
                            "start_slide_index": 1,
                            "template_slide_count": 1,
                            "texts": [
                                {
                                    "selector": {"name": "Record name", "role": "text"},
                                    "field": "name",
                                }
                            ],
                        },
                    }
                ),
                encoding="utf-8",
            )

            report = self.fill_ppt_from_records.fill_presentation(
                template, records, config, output, root / "workspace"
            )

            self.assertIs(report["ok"], False)
            self.assertEqual(report["failed_records"], [{"record_index": 2, "missing_required_fields": ["name"]}])
            result = Presentation(output)
            self.assertEqual(len(result.slides), 1)
            self.assertEqual(result.slides[0].shapes[0].text, "Good")

    def test_repeat_max_records_truncates_before_filtering_failed_records(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template = root / "template.pptx"
            records = root / "records.json"
            config = root / "layout-config.json"
            output = root / "output.pptx"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            name_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
            name_shape.name = "Record name"
            name_shape.text = "template name"
            presentation.save(template)
            records.write_text(
                json.dumps({"records": [{}, {"name": "B"}, {"name": "C"}]}),
                encoding="utf-8",
            )
            config.write_text(
                json.dumps(
                    {
                        "schema_version": 2,
                        "required_fields": ["name"],
                        "repeat": {
                            "source_slide_index": 1,
                            "start_slide_index": 1,
                            "template_slide_count": 1,
                            "max_records": 2,
                            "texts": [
                                {
                                    "selector": {"name": "Record name", "role": "text"},
                                    "field": "name",
                                }
                            ],
                        },
                    }
                ),
                encoding="utf-8",
            )

            report = self.fill_ppt_from_records.fill_presentation(
                template, records, config, output, root / "workspace"
            )

            self.assertEqual(report["failed_records"], [{"record_index": 1, "missing_required_fields": ["name"]}])
            result = Presentation(output)
            self.assertEqual(len(result.slides), 1)
            self.assertEqual(result.slides[0].shapes[0].text, "B")

    def test_table_cell_capacity_uses_column_and_row_dimensions(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template = root / "template.pptx"
            records = root / "records.json"
            config = root / "layout-config.json"
            output = root / "output.pptx"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            table_shape = slide.shapes.add_table(
                1, 2, Inches(1), Inches(1), Inches(7), Inches(2)
            )
            table_shape.name = "Wide table"
            table_shape.table.columns[0].width = Inches(0.25)
            table_shape.table.columns[1].width = Inches(6.75)
            table_shape.table.rows[0].height = Inches(2)
            table_shape.table.cell(0, 0).text = "narrow cell " + ("x" * 200)
            presentation.save(template)
            records.write_text(json.dumps({"records": []}), encoding="utf-8")
            config.write_text(json.dumps({"schema_version": 2, "slides": []}), encoding="utf-8")

            report = self.fill_ppt_from_records.fill_presentation(
                template, records, config, output, root / "workspace"
            )

            self.assertTrue(
                any(
                    warning.get("row") == 0 and warning.get("col") == 0
                    for warning in report["capacity_warnings"]
                )
            )

    def test_non_strict_failed_slide_is_cleared_instead_of_preserving_template_text(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template = root / "template.pptx"
            records = root / "records.json"
            config = root / "layout-config.json"
            output = root / "output.pptx"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            stale_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
            stale_shape.name = "Failed record content"
            stale_shape.text = "OLD TEMPLATE CONTENT"
            presentation.save(template)
            records.write_text(json.dumps({"records": [{"name": "Good"}, {}]}), encoding="utf-8")
            config.write_text(
                json.dumps(
                    {
                        "schema_version": 2,
                        "required_fields": ["name"],
                        "slides": [
                            {
                                "slide_index": 1,
                                "record_index": 2,
                                "texts": [{"shape_index": 1, "field": "name"}],
                            }
                        ],
                    }
                ),
                encoding="utf-8",
            )

            report = self.fill_ppt_from_records.fill_presentation(
                template, records, config, output, root / "workspace"
            )

            self.assertEqual(report["failed_records"], [{"record_index": 2, "missing_required_fields": ["name"]}])
            result = Presentation(output)
            self.assertFalse(any("OLD TEMPLATE CONTENT" in shape.text for shape in result.slides[0].shapes if getattr(shape, "has_text_frame", False)))

    def test_non_strict_pipeline_preserves_failed_record_details_and_continues(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template = root / "template.pptx"
            config = root / "layout-config.json"
            bad_records = root / "bad-records.json"
            good_records = root / "good-records.json"
            batch = root / "batch.json"
            report_path = root / "batch-report.json"
            bad_output = root / "bad-output.pptx"
            good_output = root / "good-output.pptx"

            Presentation().save(template)
            config.write_text(json.dumps({"required_fields": ["name"], "slides": []}), encoding="utf-8")
            bad_records.write_text(json.dumps({"records": [{}]}), encoding="utf-8")
            good_records.write_text(json.dumps({"records": [{"name": "Good"}]}), encoding="utf-8")
            batch.write_text(
                json.dumps(
                    [
                        {
                            "template": str(template),
                            "records": str(bad_records),
                            "layout_config": str(config),
                            "output": str(bad_output),
                        },
                        {
                            "template": str(template),
                            "records": str(good_records),
                            "layout_config": str(config),
                            "output": str(good_output),
                        },
                    ]
                ),
                encoding="utf-8",
            )

            old_argv = sys.argv
            try:
                sys.argv = [
                    "run_ppt_batch_pipeline.py",
                    "--batch",
                    str(batch),
                    "--report",
                    str(report_path),
                ]
                self.assertEqual(self.run_ppt_batch_pipeline.main(), 2)
            finally:
                sys.argv = old_argv

            batch_report = json.loads(report_path.read_text(encoding="utf-8-sig"))
            self.assertEqual(len(batch_report["jobs"]), 2)
            self.assertEqual(
                batch_report["jobs"][0]["failed_records"],
                [{"record_index": 1, "missing_required_fields": ["name"]}],
            )
            self.assertIs(batch_report["jobs"][1]["ok"], True)
            self.assertTrue(good_output.exists())

    def test_buyer_briefing_clears_unused_slots(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            template = Presentation()
            slide = template.slides.add_slide(template.slide_layouts[6])
            mapping = {
                "title_shape": 1,
                "buyers_per_slide": 6,
                "slots": [],
            }
            for index in range(6):
                summary_shape = slide.shapes.add_textbox(0, 0, 100, 100)
                products_shape = slide.shapes.add_textbox(0, 0, 100, 100)
                summary_shape.text = "stale summary"
                products_shape.text = "stale products"
                mapping["slots"].append({
                    "summary_shape": len(slide.shapes) - 1,
                    "products_shape": len(slide.shapes),
                })
            title_shape = slide.shapes.add_textbox(0, 0, 100, 100)
            title_shape.text = "stale title"
            mapping["title_shape"] = len(slide.shapes)

            page = {
                "title": "Category",
                "buyers": [
                    {"name": f"Buyer {index}", "summary": f"Buyer {index} summary", "products": f"Product {index}"}
                    for index in range(5)
                ] + [{}],
            }
            self.fill_buyer_briefing_pages.fill_slide(slide, page, mapping)

            self.assertEqual(slide.shapes[mapping["slots"][0]["summary_shape"] - 1].text, "Buyer 0 summary")
            self.assertEqual(slide.shapes[mapping["slots"][0]["products_shape"] - 1].text, "采购品类：Product 0")
            self.assertEqual(slide.shapes[mapping["slots"][5]["summary_shape"] - 1].text, "")
            self.assertEqual(slide.shapes[mapping["slots"][5]["products_shape"] - 1].text, "")

    def test_buyer_briefing_accepts_fewer_than_six_and_reports_missing_buyers(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template_path = root / "template.pptx"
            pages_path = root / "pages.json"
            output_path = root / "output.pptx"
            report_path = root / "report.json"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            mapping = {"title_shape": 1, "buyers_per_slide": 6, "slots": []}
            for index in range(6):
                summary = slide.shapes.add_textbox(0, 0, Inches(2), Inches(0.4))
                summary.text = f"stale summary {index}"
                products = slide.shapes.add_textbox(0, 0, Inches(2), Inches(0.3))
                products.text = f"stale products {index}"
                mapping["slots"].append({
                    "summary_shape": len(slide.shapes) - 1,
                    "products_shape": len(slide.shapes),
                })
            title = slide.shapes.add_textbox(0, 0, Inches(2), Inches(0.3))
            title.text = "stale title"
            mapping["title_shape"] = len(slide.shapes)
            presentation.save(template_path)

            pages_path.write_text(
                json.dumps({
                    "pages": [{
                        "title": "Category",
                        "buyers": [{
                            "name": "Buyer 1",
                            "summary": "Buyer 1 is a factual company summary.",
                            "products": "Product 1",
                        }],
                    }],
                    "mapping": mapping,
                }),
                encoding="utf-8",
            )

            report = self.fill_buyer_briefing_pages.fill_presentation(
                template_path, pages_path, output_path, report_path
            )

            self.assertTrue(output_path.exists())
            self.assertEqual(len(report["missing_buyers"]), 5)
            self.assertEqual(report["missing_buyers"][0]["page"], 1)
            self.assertEqual(report["missing_buyers"][0]["slot"], 2)
            result = Presentation(output_path)
            self.assertEqual(result.slides[0].shapes[mapping["slots"][5]["summary_shape"] - 1].text, "")
            self.assertEqual(result.slides[0].shapes[mapping["slots"][5]["products_shape"] - 1].text, "")

    def test_buyer_briefing_reports_overlong_text(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        title = slide.shapes.add_textbox(0, 0, Inches(2), Inches(0.3))
        title.text = "title"
        summary = slide.shapes.add_textbox(0, 0, Inches(0.6), Inches(0.2))
        summary.text = "old"
        products = slide.shapes.add_textbox(0, 0, Inches(0.6), Inches(0.2))
        products.text = "old"
        mapping = {
            "title_shape": 1,
            "buyers_per_slide": 1,
            "slots": [{"summary_shape": 2, "products_shape": 3}],
        }

        report = self.fill_buyer_briefing_pages.fill_slide(
            slide,
            {
                "title": "Category",
                "buyers": [{
                    "name": "Buyer",
                    "summary": "A" * 200,
                    "products": "Product",
                }],
            },
            mapping,
            page_number=1,
        )

        self.assertTrue(report["overlong_text"])
        self.assertEqual(report["overlong_text"][0]["page"], 1)

    def test_buyer_briefing_rejects_more_than_six_buyers(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        mapping = {"title_shape": 1, "buyers_per_slide": 6, "slots": []}
        for index in range(6):
            summary = slide.shapes.add_textbox(0, 0, Inches(2), Inches(0.4))
            products = slide.shapes.add_textbox(0, 0, Inches(2), Inches(0.3))
            mapping["slots"].append({
                "summary_shape": len(slide.shapes) + 1,
                "products_shape": len(slide.shapes) + 2,
            })
        title = slide.shapes.add_textbox(0, 0, Inches(2), Inches(0.3))
        mapping["title_shape"] = len(slide.shapes)

        with self.assertRaisesRegex(ValueError, "at most 6 buyers"):
            self.fill_buyer_briefing_pages.fill_slide(
                slide,
                {"title": "Category", "buyers": [{} for _ in range(7)]},
                mapping,
            )

    def test_buyer_briefing_clears_all_mapped_slots_beyond_configured_capacity(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        mapping = {"title_shape": 1, "buyers_per_slide": 4, "slots": []}
        slot_styles = []
        for index in range(6):
            summary = slide.shapes.add_textbox(0, 0, Inches(2), Inches(0.4))
            summary.text = f"stale summary {index}"
            summary_run = summary.text_frame.paragraphs[0].runs[0]
            summary_run.font.bold = True
            summary_run.font.size = Pt(11)
            products = slide.shapes.add_textbox(0, 0, Inches(2), Inches(0.3))
            products.text = f"stale products {index}"
            slot_styles.append((summary_run.font.bold, summary_run.font.size.pt))
            mapping["slots"].append({
                "summary_shape": len(slide.shapes) - 1,
                "products_shape": len(slide.shapes),
            })
        title = slide.shapes.add_textbox(0, 0, Inches(2), Inches(0.3))
        mapping["title_shape"] = len(slide.shapes)

        self.fill_buyer_briefing_pages.fill_slide(
            slide,
            {
                "title": "Category",
                "buyers": [{
                    "name": "Buyer",
                    "summary": "Buyer summary",
                    "products": "Product",
                }],
            },
            mapping,
        )

        for index, slot in enumerate(mapping["slots"]):
            summary_shape = slide.shapes[slot["summary_shape"] - 1]
            products_shape = slide.shapes[slot["products_shape"] - 1]
            if index == 0:
                self.assertEqual(summary_shape.text, "Buyer summary")
                self.assertEqual(products_shape.text, "采购品类：Product")
            else:
                self.assertEqual(summary_shape.text, "")
                self.assertEqual(products_shape.text, "")
                summary_run = summary_shape.text_frame.paragraphs[0].runs[0]
                self.assertEqual(summary_run.font.bold, slot_styles[index][0])
                self.assertEqual(summary_run.font.size.pt, slot_styles[index][1])

    def test_buyer_briefing_cli_uses_embedded_mapping_without_explicit_config(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template_path = root / "template.pptx"
            pages_path = root / "pages.json"
            explicit_mapping_path = root / "explicit-mapping.json"
            embedded_output = root / "embedded-output.pptx"
            explicit_output = root / "explicit-output.pptx"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            for index in range(6):
                shape = slide.shapes.add_textbox(0, 0, Inches(2), Inches(0.4))
                shape.text = f"stale {index + 1}"
            presentation.save(template_path)

            default_mapping = {
                "title_shape": 1,
                "buyers_per_slide": 1,
                "slots": [{"summary_shape": 2, "products_shape": 3}],
            }
            embedded_mapping = {
                "title_shape": 4,
                "buyers_per_slide": 1,
                "slots": [{"summary_shape": 5, "products_shape": 6}],
            }
            pages_path.write_text(
                json.dumps({
                    "pages": [{
                        "title": "Embedded title",
                        "buyers": [{
                            "name": "Embedded buyer",
                            "summary": "Embedded buyer summary",
                            "products": "Embedded product",
                        }],
                    }],
                    "mapping": embedded_mapping,
                }),
                encoding="utf-8",
            )
            explicit_mapping_path.write_text(json.dumps(default_mapping), encoding="utf-8")

            old_argv = sys.argv
            try:
                with patch.object(self.fill_buyer_briefing_pages, "DEFAULT_MAPPING", default_mapping):
                    sys.argv = [
                        "fill_buyer_briefing_pages.py",
                        str(template_path),
                        str(pages_path),
                        str(embedded_output),
                    ]
                    self.assertEqual(self.fill_buyer_briefing_pages.main(), 0)

                    sys.argv = [
                        "fill_buyer_briefing_pages.py",
                        str(template_path),
                        str(pages_path),
                        str(explicit_output),
                        "--layout-config",
                        str(explicit_mapping_path),
                    ]
                    self.assertEqual(self.fill_buyer_briefing_pages.main(), 0)
            finally:
                sys.argv = old_argv

            embedded_result = Presentation(embedded_output)
            self.assertEqual(embedded_result.slides[0].shapes[3].text, "Embedded title")
            self.assertEqual(embedded_result.slides[0].shapes[4].text, "Embedded buyer summary")
            self.assertEqual(embedded_result.slides[0].shapes[5].text, "采购品类：Embedded product")

            explicit_result = Presentation(explicit_output)
            self.assertEqual(explicit_result.slides[0].shapes[0].text, "Embedded title")
            self.assertEqual(explicit_result.slides[0].shapes[1].text, "Embedded buyer summary")
            self.assertEqual(explicit_result.slides[0].shapes[2].text, "采购品类：Embedded product")
            self.assertEqual(explicit_result.slides[0].shapes[3].text, "stale 4")

    def test_console_briefing_export_allows_fewer_than_six_buyers(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            state = self.control_console.ConsoleState(root / "projects")
            project = state.create_project("briefing", "buyer_briefing")
            paths = state.require(project["slug"])
            self.control_console.write_json(paths["records"], {
                "pages": [{
                    "title": "Briefing",
                    "buyers": [{"name": f"Buyer {index}"} for index in range(5)],
                }],
            })
            job = {
                "output": str(root / "finished.pptx"),
                "report": str(root / "report.json"),
            }

            command = state._briefing_export_command(job, paths, root / "run")

            self.assertEqual(command[0], sys.executable)
            self.assertIn("fill_buyer_briefing_pages.py", command[1])

    def test_console_briefing_export_rejects_empty_or_over_capacity_pages(self):
        for buyer_count in (0, 7):
            with self.subTest(buyer_count=buyer_count), tempfile.TemporaryDirectory() as temp_dir:
                root = Path(temp_dir)
                state = self.control_console.ConsoleState(root / "projects")
                project = state.create_project("briefing", "buyer_briefing")
                paths = state.require(project["slug"])
                self.control_console.write_json(paths["records"], {
                    "pages": [{
                        "title": "Briefing",
                        "buyers": [{"name": f"Buyer {index}"} for index in range(buyer_count)],
                    }],
                })
                job = {
                    "output": str(root / "finished.pptx"),
                    "report": str(root / "report.json"),
                }

                with self.assertRaisesRegex(ValueError, "1-6"):
                    state._briefing_export_command(job, paths, root / "run")

    def test_yitu_missing_shape_fails_without_writing_output(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template_path = root / "template.pptx"
            output_path = root / "output.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            shape = slide.shapes.add_textbox(0, 0, 100, 100)
            shape.name = "Known Shape"
            presentation.save(template_path)

            with self.assertRaisesRegex(KeyError, "Mapped shapes were not found"):
                self.yitu_quanjie_replace.run_replacement(
                    template_path,
                    output_path,
                    {"Missing Shape": "replacement"},
                )

            self.assertFalse(output_path.exists())

    def test_yitu_validate_reports_missing_shape_and_text_capacity(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template_path = root / "template.pptx"
            output_path = root / "output.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            shape = slide.shapes.add_textbox(0, 0, Inches(1), Inches(0.3))
            shape.name = "Known Shape"
            shape.text = "short"
            presentation.save(template_path)

            report = self.yitu_quanjie_replace.validate_replacement(
                template_path,
                output_path,
                {"Known Shape": "x" * 100, "Missing Shape": "replacement"},
            )

            self.assertFalse(report["ok"])
            self.assertEqual(report["missing_shapes"], ["Missing Shape"])
            self.assertTrue(report["overflows"])
            self.assertTrue(report["output"]["path"].endswith("output.pptx"))
            self.assertFalse(output_path.exists())

    def test_yitu_text_capacity_reports_same_length_chinese_overflow_in_narrow_large_font_box(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template_path = root / "template.pptx"
            output_path = root / "output.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            shape = slide.shapes.add_textbox(0, 0, Inches(0.4), Inches(0.3))
            shape.name = "Narrow Shape"
            original_text = "模板文字" * 5
            shape.text = original_text
            shape.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
            presentation.save(template_path)

            replacement = "替换文字" * 5
            self.assertEqual(len(replacement), len(original_text))
            report = self.yitu_quanjie_replace.validate_replacement(
                template_path,
                output_path,
                {"Narrow Shape": replacement},
            )

            self.assertTrue(
                any(
                    item["target"] == "Narrow Shape" and item["kind"] == "shape"
                    for item in report["overflows"]
                )
            )

    def test_yitu_text_capacity_does_not_flag_reasonable_text_in_wide_box(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template_path = root / "template.pptx"
            output_path = root / "output.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            shape = slide.shapes.add_textbox(0, 0, Inches(4), Inches(1))
            shape.name = "Wide Shape"
            shape.text = "x"
            shape.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
            presentation.save(template_path)

            report = self.yitu_quanjie_replace.validate_replacement(
                template_path,
                output_path,
                {"Wide Shape": "这是一段可以放入宽文本框的合理中文文本"},
            )

            self.assertTrue(report["ok"])
            self.assertFalse(any(item["target"] == "Wide Shape" for item in report["overflows"]))

    def test_yitu_text_capacity_rejects_explicit_newline_beyond_vertical_capacity(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template_path = root / "template.pptx"
            output_path = root / "output.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            shape = slide.shapes.add_textbox(0, 0, Inches(3), Inches(0.3))
            shape.name = "Single Line Shape"
            shape.text = "template"
            shape.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
            presentation.save(template_path)

            report = self.yitu_quanjie_replace.validate_replacement(
                template_path,
                output_path,
                {"Single Line Shape": "A\nB"},
            )

            self.assertTrue(
                any(
                    item["target"] == "Single Line Shape"
                    and item["kind"] == "shape"
                    and item["required_lines"] > item["capacity_lines"]
                    for item in report["overflows"]
                )
            )

    def test_yitu_capacity_preserves_trailing_empty_lines_for_lf_and_crlf(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template_path = root / "template.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            shape = slide.shapes.add_textbox(0, 0, Inches(3), Inches(0.3))
            shape.name = "Trailing Lines Shape"
            shape.text = "template"
            shape.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
            presentation.save(template_path)

            for replacement in ("A\n", "\n", "A\r\n"):
                with self.subTest(replacement=repr(replacement)):
                    report = self.yitu_quanjie_replace.validate_replacement(
                        template_path,
                        root / "output.pptx",
                        {"Trailing Lines Shape": replacement},
                    )

                    overflows = [
                        item for item in report["overflows"]
                        if item["target"] == "Trailing Lines Shape"
                    ]
                    self.assertTrue(overflows)
                    self.assertEqual(overflows[0]["required_lines"], 2.0)

    def test_yitu_validate_reports_invalid_table_coordinate_and_output_path(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template_path = root / "template.pptx"
            output_dir = root / "output-dir"
            output_dir.mkdir()
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            table_shape = slide.shapes.add_table(1, 1, 0, 0, Inches(2), Inches(1))
            table_shape.name = "Data Table"
            presentation.save(template_path)

            report = self.yitu_quanjie_replace.validate_replacement(
                template_path,
                output_dir,
                {},
                {"cell_99": "invalid", "cell_bad": "invalid"},
            )

            self.assertFalse(report["ok"])
            self.assertGreaterEqual(len(report["table_errors"]), 2)
            self.assertTrue(report["output"]["errors"])
            self.assertFalse((root / "output.pptx").exists())

    def test_yitu_replaces_only_the_table_validated_by_table_data(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template_path = root / "template.pptx"
            output_path = root / "output.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            first_table = slide.shapes.add_table(1, 1, 0, 0, Inches(2), Inches(1))
            second_table = slide.shapes.add_table(1, 1, Inches(3), 0, Inches(2), Inches(1))
            first_table.table.cell(0, 0).text = "first table"
            second_table.table.cell(0, 0).text = "second table"
            presentation.save(template_path)

            self.yitu_quanjie_replace.run_replacement(
                template_path,
                output_path,
                {},
                {"cell_00": "updated"},
            )

            result = Presentation(output_path)
            self.assertEqual(result.slides[0].shapes[0].table.cell(0, 0).text, "updated")
            self.assertEqual(result.slides[0].shapes[1].table.cell(0, 0).text, "second table")

    def test_yitu_table_height_uses_each_column_width_for_narrow_long_text(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template_path = root / "template.pptx"
            output_path = root / "output.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            table_shape = slide.shapes.add_table(1, 2, 0, 0, Inches(3), Inches(0.9))
            table_shape.table.columns[0].width = Inches(2.5)
            table_shape.table.columns[1].width = Inches(0.5)
            table_shape.table.cell(0, 0).text = "wide label"
            table_shape.table.cell(0, 1).text = ""
            presentation.save(template_path)

            report = self.yitu_quanjie_replace.validate_replacement(
                template_path,
                output_path,
                {},
                {"cell_01": "窄列长文本" * 6},
            )

            self.assertTrue(
                any(
                    item["kind"] == "table_height" and item["target"] == "table_0"
                    for item in report["overflows"]
                )
            )

    def test_yitu_table_height_overflow_raises_stable_error_without_length_fields(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template_path = root / "template.pptx"
            output_path = root / "output.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            table_shape = slide.shapes.add_table(1, 1, 0, 0, Inches(1), Inches(0.05))
            table_shape.table.cell(0, 0).text = "template cell text with enough capacity"
            presentation.save(template_path)

            report = self.yitu_quanjie_replace.validate_replacement(
                template_path,
                output_path,
                {},
                {"cell_00": "short"},
            )
            height_overflows = [
                item for item in report["overflows"] if item["kind"] == "table_height"
            ]
            self.assertTrue(height_overflows)
            self.assertNotIn("length", height_overflows[0])

            try:
                self.yitu_quanjie_replace.run_replacement(
                    template_path,
                    output_path,
                    {},
                    {"cell_00": "short"},
                )
            except Exception as exc:
                self.assertIsInstance(exc, ValueError)
                self.assertIn("table_0", str(exc))
            else:
                self.fail("Expected a stable ValueError for table height overflow")
            self.assertFalse(output_path.exists())

    def test_yitu_dry_run_prints_json_and_does_not_save(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template_path = root / "template.pptx"
            output_path = root / "output.pptx"
            content_path = root / "content.json"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            shape = slide.shapes.add_textbox(0, 0, Inches(2), Inches(0.3))
            shape.name = "Known Shape"
            shape.text = "template"
            presentation.save(template_path)
            content_path.write_text(json.dumps({"Missing Shape": "replacement"}), encoding="utf-8")

            old_argv = sys.argv
            stdout = io.StringIO()
            try:
                sys.argv = [
                    "yitu_quanjie_replace.py",
                    "--template", str(template_path),
                    "--output", str(output_path),
                    "--content-map", str(content_path),
                    "--dry-run",
                ]
                with redirect_stdout(stdout):
                    self.assertEqual(self.yitu_quanjie_replace.main(), 0)
            finally:
                sys.argv = old_argv

            report = json.loads(stdout.getvalue())
            self.assertFalse(report["ok"])
            self.assertEqual(report["missing_shapes"], ["Missing Shape"])
            self.assertFalse(output_path.exists())

    def test_feishu_zip_contains_only_portable_skill_files(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            output = self.build_feishu_agent_skill.build(Path(temp_dir) / "skill.zip")

            with zipfile.ZipFile(output) as archive:
                names = archive.namelist()

            self.assertIn("SKILL.md", names)
            self.assertTrue(names)
            self.assertTrue(all(name == "SKILL.md" or name.startswith("references/") for name in names))
            self.assertIn("references/agent-runtime.md", names)
            self.assertIn("references/input-schema.json", names)

    def test_feishu_builder_rejects_missing_contract_files(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            package_root = Path(temp_dir) / "package"
            package_root.mkdir()
            (package_root / "SKILL.md").write_text("# skill\n", encoding="utf-8")

            with patch.object(self.build_feishu_agent_skill, "PACKAGE_ROOT", package_root):
                with self.assertRaisesRegex(ValueError, "required"):
                    self.build_feishu_agent_skill.build(Path(temp_dir) / "skill.zip")

    def test_feishu_builder_rejects_engine_and_extra_paths(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            package_root = Path(temp_dir) / "package"
            references = package_root / "references"
            references.mkdir(parents=True)
            (package_root / "SKILL.md").write_text("# skill\n", encoding="utf-8")
            (references / "agent-runtime.md").write_text("runtime\n", encoding="utf-8")
            (references / "input-schema.json").write_text("{}\n", encoding="utf-8")
            (package_root / "engine").mkdir()
            (package_root / "engine" / "runner.py").write_text("print('no')\n", encoding="utf-8")
            (package_root / "README.md").write_text("no\n", encoding="utf-8")

            with patch.object(self.build_feishu_agent_skill, "PACKAGE_ROOT", package_root):
                with self.assertRaisesRegex(ValueError, "Unsupported package paths"):
                    self.build_feishu_agent_skill.build(Path(temp_dir) / "skill.zip")

    def test_feishu_builder_rejects_source_path_resolving_outside_package_root(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            package_root = Path(temp_dir) / "package"
            references = package_root / "references"
            references.mkdir(parents=True)
            (package_root / "SKILL.md").write_text("# skill\n", encoding="utf-8")
            (references / "agent-runtime.md").write_text("runtime\n", encoding="utf-8")
            (references / "input-schema.json").write_text("{}\n", encoding="utf-8")
            escaped = Path(temp_dir) / "outside.md"
            escaped.write_text("outside\n", encoding="utf-8")

            with patch.object(Path, "rglob", return_value=[escaped]):
                with self.assertRaisesRegex(ValueError, "escapes package root"):
                    self.build_feishu_agent_skill.validate_package(package_root)

    def test_feishu_builder_rejects_root_and_reference_symlinks(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            package_root = Path(temp_dir) / "package"
            references = package_root / "references"
            references.mkdir(parents=True)
            (package_root / "SKILL.md").write_text("# skill\n", encoding="utf-8")
            (references / "agent-runtime.md").write_text("runtime\n", encoding="utf-8")
            (references / "input-schema.json").write_text("{}\n", encoding="utf-8")
            outside = Path(temp_dir) / "outside.md"
            outside.write_text("outside\n", encoding="utf-8")

            try:
                os.symlink(outside, package_root / "root-link.md")
                os.symlink(outside, references / "reference-link.md")
            except (OSError, NotImplementedError) as exc:
                self.skipTest(f"symlink creation is unavailable: {exc}")

            with patch.object(self.build_feishu_agent_skill, "PACKAGE_ROOT", package_root):
                with self.assertRaisesRegex(ValueError, "Symlinks are not allowed"):
                    self.build_feishu_agent_skill.build(Path(temp_dir) / "skill.zip")

    def test_feishu_builder_uses_same_directory_atomic_replace(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "skill.zip"
            with patch.object(self.build_feishu_agent_skill, "os", os, create=True):
                with patch.object(os, "replace", wraps=os.replace) as replace:
                    self.build_feishu_agent_skill.build(output)

            replace.assert_called_once()
            temporary_path, replaced_path = replace.call_args.args
            self.assertEqual(Path(replaced_path), output.resolve())
            self.assertEqual(Path(temporary_path).parent, output.parent.resolve())
            self.assertFalse(Path(temporary_path).exists())

    def test_feishu_builder_preserves_existing_output_when_zip_validation_fails(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "skill.zip"
            original = b"keep this output"
            output.write_bytes(original)

            with patch.object(
                self.build_feishu_agent_skill.zipfile.ZipFile,
                "testzip",
                side_effect=RuntimeError("zip validation failed"),
            ):
                with self.assertRaisesRegex(RuntimeError, "zip validation failed"):
                    self.build_feishu_agent_skill.build(output)

            self.assertEqual(output.read_bytes(), original)
            self.assertEqual(list(Path(temp_dir).glob(".skill.zip.*.tmp")), [])

    def test_feishu_contract_documents_stable_metadata(self):
        skill_text = (REPO_ROOT / "feishu-agent-skill" / "SKILL.md").read_text(encoding="utf-8")
        runtime_text = (REPO_ROOT / "feishu-agent-skill" / "references" / "agent-runtime.md").read_text(encoding="utf-8")
        schema = json.loads((REPO_ROOT / "feishu-agent-skill" / "references" / "input-schema.json").read_text(encoding="utf-8"))

        for field in ("contract_version", "sources", "verification_status", "warnings"):
            self.assertIn(field, skill_text)
            self.assertIn(field, runtime_text)
            self.assertIn(field, schema["properties"])

    def test_console_defaults_to_loopback_and_requires_opt_in_for_public_bind(self):
        old_argv = sys.argv
        try:
            sys.argv = ["run_control_console.py"]
            args = self.control_console.parse_args()
        finally:
            sys.argv = old_argv
        self.assertEqual(args.host, "127.0.0.1")
        self.assertFalse(args.allow_non_loopback)

        with tempfile.TemporaryDirectory() as temp_dir:
            state = self.control_console.ConsoleState(Path(temp_dir))
            with self.assertRaisesRegex(ValueError, "loopback"):
                self.control_console.build_server("0.0.0.0", 0, state)
            server = self.control_console.build_server(
                "0.0.0.0", 0, state, allow_non_loopback=True
            )
            server.server_close()

    def test_console_model_payload_redacts_credentials(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            state = self.control_console.ConsoleState(Path(temp_dir))
            state.model_settings["unified_key"] = "super-secret-key"
            state.model_settings["unified_base_url"] = "https://user:super-secret-key@example.invalid/v1?api_key=super-secret-key"

            payload = state.model_settings_payload()
            serialized = json.dumps(payload, ensure_ascii=False)

        self.assertNotIn("super-secret-key", serialized)
        self.assertIn("redacted", serialized.lower())

    def test_console_json_writes_are_atomic(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            target = Path(temp_dir) / "state.json"
            with patch.object(self.control_console.os, "replace", wraps=self.control_console.os.replace) as replace:
                self.control_console.write_json(target, {"ok": True})

            self.assertTrue(replace.called)
            self.assertEqual(json.loads(target.read_text(encoding="utf-8")), {"ok": True})
            self.assertEqual(list(Path(temp_dir).glob("*.tmp")), [])

    def test_console_subprocesses_receive_bounded_timeouts(self):
        completed = Namespace(returncode=0, stdout=b"", stderr=b"")
        with patch.object(self.control_console.subprocess, "run", return_value=completed) as run:
            self.assertEqual(self.control_console.run_command(["helper"]), (0, "", ""))

        self.assertIn("timeout", run.call_args.kwargs)
        self.assertGreater(run.call_args.kwargs["timeout"], 0)

    def test_console_diagnostic_child_has_explicit_timeout(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            state = self.control_console.ConsoleState(Path(temp_dir))
            with patch.object(self.control_console, "probe_base_url", return_value={
                "ok": False, "reachable": False, "status": 0, "error": "offline"
            }):
                with patch.object(
                    self.control_console,
                    "run_command",
                    return_value=(0, json.dumps({"reachable": False}), ""),
                ) as run:
                    state.diagnose_model_connection({"role": "research"})

        self.assertEqual(
            run.call_args.kwargs["timeout"],
            self.control_console.DIAGNOSTIC_SUBPROCESS_TIMEOUT_SECONDS,
        )
