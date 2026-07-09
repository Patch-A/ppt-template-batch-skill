from __future__ import annotations

import argparse
import json
import tempfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

from image_layout_utils import prepare_logo_image, prepare_site_image


def load_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8-sig"))


def remove_shape(shape) -> None:
    sp = shape._element
    sp.getparent().remove(sp)


def find_shape(slide, left_pt: float, top_pt: float):
    target_left = Pt(left_pt)
    target_top = Pt(top_pt)
    tolerance = Pt(3)
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        if abs(shape.left - target_left) <= tolerance and abs(shape.top - target_top) <= tolerance:
            return shape
    raise ValueError(f"Target picture shape not found near ({left_pt}, {top_pt})")


def fit_into_box(shape, left, top, width, height, align_left: bool = False) -> None:
    ratio_x = width / shape.width
    ratio_y = height / shape.height
    ratio = min(ratio_x, ratio_y)
    shape.width = int(shape.width * ratio)
    shape.height = int(shape.height * ratio)
    shape.left = int(left if align_left else left + ((width - shape.width) / 2))
    shape.top = int(top + ((height - shape.height) / 2))


def replace_picture(slide, target, image_path: Path, fill: bool) -> None:
    left, top, width, height = target.left, target.top, target.width, target.height
    remove_shape(target)
    if fill:
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
        return
    new_shape = slide.shapes.add_picture(str(image_path), left, top)
    fit_into_box(new_shape, left, top, width, height, align_left=True)


def clear_picture_target(slide, left_pt: float, top_pt: float) -> None:
    try:
        target = find_shape(slide, left_pt, top_pt)
    except ValueError:
        return
    remove_shape(target)


def clear_region(slide, region: dict[str, float]) -> None:
    left = Pt(region["left"])
    top = Pt(region["top"])
    right = Pt(region["right"])
    bottom = Pt(region["bottom"])
    to_delete = []
    for shape in slide.shapes:
        if shape.left >= left and shape.left <= right and shape.top >= top and shape.top <= bottom:
            if shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.PLACEHOLDER}:
                to_delete.append(shape)
    for shape in to_delete:
        remove_shape(shape)


def add_logo(slide, image_path: Path, logo_cfg: dict[str, float]) -> None:
    left = Pt(logo_cfg["left"])
    top = Pt(logo_cfg["top"])
    width = Pt(logo_cfg["width"])
    height = Pt(logo_cfg["height"])
    new_shape = slide.shapes.add_picture(str(image_path), left, top)
    fit_into_box(new_shape, left, top, width, height, align_left=True)


def main() -> int:
    parser = argparse.ArgumentParser(description="Fallback image application without PowerPoint COM.")
    parser.add_argument("--input-ppt", required=True)
    parser.add_argument("--buyers-json", required=True)
    parser.add_argument("--layout-config", required=True)
    parser.add_argument("--output-ppt", required=True)
    parser.add_argument("--preview-dir")
    args = parser.parse_args()

    input_ppt = Path(args.input_ppt)
    buyers = load_json(Path(args.buyers_json))
    layout = load_json(Path(args.layout_config))
    prs = Presentation(str(input_ppt))

    temp_dir = Path(tempfile.mkdtemp(prefix="buyer-board-images-"))
    start_index = int(layout["content"]["start_slide_index"]) - 1
    slots = {int(item["slide_offset"]): item for item in layout["images"]["slides"]}

    for idx, buyer in enumerate(buyers):
        slide = prs.slides[start_index + idx]
        slot = slots.get(idx)
        if not slot:
            continue

        if slot.get("site") and buyer.get("site_image_path"):
            site_cfg = slot["site"]
            target = find_shape(slide, float(site_cfg["target_left"]), float(site_cfg["target_top"]))
            prepared = prepare_site_image(
                Path(buyer["site_image_path"]),
                temp_dir,
                int(target.width),
                int(target.height),
            )
            replace_picture(slide, target, prepared.output_path, True)
        elif slot.get("site"):
            clear_picture_target(slide, float(slot["site"]["target_left"]), float(slot["site"]["target_top"]))

        if slot.get("logo") and buyer.get("logo_path"):
            logo_cfg = slot["logo"]
            try:
                if logo_cfg["mode"] == "add":
                    target_width = int(Pt(logo_cfg["width"]))
                    target_height = int(Pt(logo_cfg["height"]))
                else:
                    target = find_shape(slide, float(logo_cfg["target_left"]), float(logo_cfg["target_top"]))
                    target_width = int(target.width)
                    target_height = int(target.height)
                prepared_logo = prepare_logo_image(Path(buyer["logo_path"]), temp_dir, target_width, target_height)
            except RuntimeError:
                prepared_logo = None
            if logo_cfg["mode"] == "add":
                if logo_cfg.get("clear_region"):
                    clear_region(slide, logo_cfg["clear_region"])
                if prepared_logo is not None:
                    add_logo(slide, prepared_logo.output_path, logo_cfg)
            else:
                if prepared_logo is not None:
                    target = find_shape(slide, float(logo_cfg["target_left"]), float(logo_cfg["target_top"]))
                    replace_picture(slide, target, prepared_logo.output_path, False)
                else:
                    clear_picture_target(slide, float(logo_cfg["target_left"]), float(logo_cfg["target_top"]))
        elif slot.get("logo") and slot["logo"]["mode"] == "replace":
            clear_picture_target(slide, float(slot["logo"]["target_left"]), float(slot["logo"]["target_top"]))
        elif slot.get("logo") and slot["logo"]["mode"] == "add" and slot["logo"].get("clear_region"):
            clear_region(slide, slot["logo"]["clear_region"])

    output_ppt = Path(args.output_ppt)
    output_ppt.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_ppt))
    if args.preview_dir:
        preview_dir = Path(args.preview_dir)
        preview_dir.mkdir(parents=True, exist_ok=True)
        (preview_dir / "README.txt").write_text(
            "Slide preview export was skipped because PowerPoint COM was unavailable. "
            "The PPT was generated with the Python fallback image pipeline.",
            encoding="utf-8",
        )
    print(output_ppt)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
