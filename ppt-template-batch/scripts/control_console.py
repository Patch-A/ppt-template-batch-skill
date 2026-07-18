from __future__ import annotations

import argparse
import ipaddress
import json
import mimetypes
import os
import re
import shutil
import subprocess
import sys
import threading
import tempfile
import uuid
import webbrowser
import zipfile
from datetime import datetime, timezone
from http import HTTPStatus
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from io import BytesIO
from pathlib import Path
from typing import Any
from urllib.error import HTTPError
from urllib.parse import parse_qsl, quote, unquote, urlencode, urlparse, urlsplit, urlunsplit
from urllib.request import Request, urlopen

from discover_buyer_profiles import buyer_evidence_context, normalize_products, pad_or_trim_bio, refine_buyer_products

MAX_UPLOAD_BYTES = 250 * 1024 * 1024
DEFAULT_SUBPROCESS_TIMEOUT_SECONDS = 300
DIAGNOSTIC_SUBPROCESS_TIMEOUT_SECONDS = 20
EXPORT_SUBPROCESS_TIMEOUT_SECONDS = 300
REDACTED = "[redacted]"
SENSITIVE_KEY_NAMES = {
    "access_token", "api_key", "apikey", "authorization", "client_secret", "password",
    "secret", "token", "unified_key", "research_key", "visual_key", "layout_key",
}
PROJECT_SLUG_RE = re.compile(r"^[a-z0-9][a-z0-9-]{0,62}$")
BUYER_FIELDS = (
    "name", "country", "website", "products", "bio", "logo_path", "site_image_path", "research_notes",
    "buyer_type", "demand_scenarios", "local_presence", "import_signal", "evidence", "confidence", "risks",
)
BUYER_SCORE_FIELDS = ("fit_score", "demand_score", "import_score", "verification_score", "total_score")

BUYER_BRIEFING_DEFAULT_MAPPING: dict[str, Any] = {
    "title_shape": 5,
    "buyers_per_slide": 6,
    "slots": [
        {"summary_shape": 15, "products_shape": 23},
        {"summary_group": 16, "summary_child": 2, "products_group": 16, "products_child": 3},
        {"summary_group": 17, "summary_child": 2, "products_group": 17, "products_child": 3},
        {"summary_shape": 19, "products_shape": 26},
        {"summary_shape": 20, "products_shape": 25},
        {"summary_shape": 22, "products_shape": 24},
    ],
}

MODEL_PROVIDER_DEFAULTS: dict[str, dict[str, Any]] = {
    "openai": {"label": "OpenAI", "base_url": "https://api.openai.com/v1", "research_model": "gpt-4.1", "visual_model": "gpt-image-1", "layout_model": "gpt-4.1", "models": ["gpt-4.1", "gpt-4.1-mini", "gpt-4o", "gpt-image-1"]},
    "deepseek": {"label": "DeepSeek", "base_url": "https://api.deepseek.com", "research_model": "deepseek-chat", "visual_model": "", "layout_model": "deepseek-chat", "models": ["deepseek-chat", "deepseek-reasoner"]},
    "qwen": {"label": "通义千问 / Qwen", "base_url": "https://dashscope.aliyuncs.com/compatible-mode/v1", "research_model": "qwen-plus", "visual_model": "", "layout_model": "qwen-plus", "models": ["qwen-plus", "qwen-turbo", "qwen-max", "qwen-long"]},
    "zhipu": {"label": "智谱 GLM", "base_url": "https://open.bigmodel.cn/api/paas/v4", "research_model": "glm-4-plus", "visual_model": "", "layout_model": "glm-4-plus", "models": ["glm-4-plus", "glm-4-air", "glm-4-flash"]},
    "kimi": {"label": "Kimi / Moonshot", "base_url": "https://api.moonshot.cn/v1", "research_model": "moonshot-v1-8k", "visual_model": "", "layout_model": "moonshot-v1-8k", "models": ["moonshot-v1-8k", "moonshot-v1-32k", "moonshot-v1-128k"]},
    "doubao": {"label": "豆包 / 火山方舟", "base_url": "https://ark.cn-beijing.volces.com/api/v3", "research_model": "", "visual_model": "", "layout_model": "", "models": []},
    "baidu": {"label": "百度千帆 / 文心", "base_url": "", "research_model": "", "visual_model": "", "layout_model": "", "models": []},
    "minimax": {"label": "MiniMax", "base_url": "https://api.minimax.chat/v1", "research_model": "", "visual_model": "", "layout_model": "", "models": []},
    "siliconflow": {"label": "硅基流动 SiliconFlow", "base_url": "https://api.siliconflow.cn/v1", "research_model": "Qwen/Qwen2.5-72B-Instruct", "visual_model": "", "layout_model": "Qwen/Qwen2.5-72B-Instruct", "models": ["Qwen/Qwen2.5-72B-Instruct", "deepseek-ai/DeepSeek-V3", "deepseek-ai/DeepSeek-R1"]},
    "openrouter": {"label": "OpenRouter", "base_url": "https://openrouter.ai/api/v1", "research_model": "", "visual_model": "", "layout_model": "", "models": []},
    "ollama": {"label": "本地 Ollama", "base_url": "http://127.0.0.1:11434/v1", "research_model": "qwen2.5", "visual_model": "", "layout_model": "qwen2.5", "models": ["qwen2.5", "llama3.1", "deepseek-r1"]},
    "lmstudio": {"label": "本地 LM Studio", "base_url": "http://127.0.0.1:1234/v1", "research_model": "local-model", "visual_model": "", "layout_model": "local-model", "models": ["local-model"]},
    "compatible": {"label": "自定义 OpenAI 兼容接口", "base_url": "", "research_model": "", "visual_model": "", "layout_model": "", "models": []},
}
PROVIDER_CHOICES = tuple(MODEL_PROVIDER_DEFAULTS)
ROLE_CHOICES = ("research", "visual", "layout")


def clean_provider(value: Any) -> str:
    provider = str(value or "deepseek").strip().lower()
    return provider if provider in MODEL_PROVIDER_DEFAULTS else "compatible"


def clean_base_url(value: Any, provider: str) -> str:
    base_url = str(value or "").strip().rstrip("/")
    if base_url:
        return base_url
    return MODEL_PROVIDER_DEFAULTS.get(provider, MODEL_PROVIDER_DEFAULTS["compatible"])["base_url"]


def is_loopback_host(host: str) -> bool:
    normalized = str(host or "").strip().strip("[]").lower()
    if normalized == "localhost":
        return True
    try:
        return ipaddress.ip_address(normalized).is_loopback
    except ValueError:
        return False


def validate_bind_host(host: str, allow_non_loopback: bool = False) -> str:
    normalized = str(host or "").strip()
    if not normalized:
        raise ValueError("Console host must be explicit; default to the loopback address.")
    if not allow_non_loopback and not is_loopback_host(normalized):
        raise ValueError("Non-loopback console binding requires --allow-non-loopback.")
    return normalized


def model_list_endpoint(base_url: str) -> str:
    base = base_url.rstrip("/")
    return f"{base}/models"


def probe_base_url(url: str, timeout: int = 8) -> dict[str, Any]:
    if not url:
        return {"ok": False, "reachable": False, "status": 0, "error": "Base URL is empty."}
    request = Request(url.rstrip("/"), headers={"Accept": "application/json"}, method="GET")
    try:
        with urlopen(request, timeout=timeout) as response:
            return {"ok": True, "reachable": True, "status": int(response.status), "error": ""}
    except HTTPError as exc:
        return {"ok": True, "reachable": True, "status": int(exc.code), "error": str(exc.reason or exc)}
    except Exception as exc:
        return {"ok": False, "reachable": False, "status": 0, "error": str(exc)}


def clean_research_mode(value: Any) -> str:
    mode = str(value or "model_only").strip().lower()
    return mode if mode in {"model_only", "openai_web_search"} else "model_only"


def infer_provider_from_model(model: Any, provider: str, explicit_provider: bool = False) -> str:
    model_name = str(model or "").strip().lower()
    if not explicit_provider and provider in {"openai", "compatible", ""} and model_name.startswith("deepseek"):
        return "deepseek"
    return provider


def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds")


def repo_root() -> Path:
    return Path(__file__).resolve().parents[2]


def skill_root() -> Path:
    return Path(__file__).resolve().parents[1]


def static_root() -> Path:
    return skill_root() / "assets" / "control-console"


def skill_script(name: str) -> Path:
    return Path(__file__).resolve().parent / name


def generic_pipeline_script() -> Path:
    repo_runner = repo_root() / "scripts" / "run_ppt_batch_pipeline.py"
    return repo_runner if repo_runner.is_file() else skill_script("fill_ppt_from_records.py")


def buyer_pipeline_script() -> Path | None:
    candidate = repo_root() / "scripts" / "run_buyer_board_pipeline.py"
    return candidate if candidate.is_file() else None


def read_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8-sig"))


def write_json(path: Path, payload: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    fd, temporary_name = tempfile.mkstemp(prefix=f".{path.name}.", suffix=".tmp", dir=str(path.parent))
    temporary_path = Path(temporary_name)
    try:
        with os.fdopen(fd, "w", encoding="utf-8", newline="\n") as handle:
            json.dump(payload, handle, ensure_ascii=False, indent=2)
            handle.write("\n")
            handle.flush()
            os.fsync(handle.fileno())
        os.replace(temporary_path, path)
    except Exception:
        try:
            temporary_path.unlink()
        except FileNotFoundError:
            pass
        raise


def _sensitive_key(value: Any) -> bool:
    normalized = str(value or "").strip().lower().replace("-", "_")
    return normalized in SENSITIVE_KEY_NAMES or normalized.endswith(("_api_key", "_token", "_secret"))


def _redact_url(value: str) -> str:
    try:
        parsed = urlsplit(value)
        if not parsed.scheme or not parsed.netloc:
            return value
        hostname = parsed.hostname or ""
        if ":" in hostname and not hostname.startswith("["):
            hostname = f"[{hostname}]"
        try:
            port = f":{parsed.port}" if parsed.port is not None else ""
        except ValueError:
            port = ""
        userinfo = ""
        query = []
        for key, item in parse_qsl(parsed.query, keep_blank_values=True):
            query.append((key, REDACTED if _sensitive_key(key) else item))
        return urlunsplit((parsed.scheme, f"{userinfo}{hostname}{port}", parsed.path, urlencode(query), parsed.fragment))
    except ValueError:
        return value


def redact_credentials(value: Any) -> Any:
    if isinstance(value, dict):
        return {
            key: REDACTED if _sensitive_key(key) else redact_credentials(item)
            for key, item in value.items()
        }
    if isinstance(value, list):
        return [redact_credentials(item) for item in value]
    if isinstance(value, tuple):
        return tuple(redact_credentials(item) for item in value)
    if isinstance(value, str):
        return _redact_url(value)
    return value


def _redact_output(value: str, env: dict[str, str]) -> str:
    for key, secret in env.items():
        if _sensitive_key(key) and secret:
            value = value.replace(secret, REDACTED)
    return value


def decode_output(value: bytes | None) -> str:
    if not value:
        return ""
    if isinstance(value, str):
        return value
    for encoding in ("utf-8", "gbk", sys.getdefaultencoding()):
        try:
            return value.decode(encoding)
        except UnicodeDecodeError:
            continue
    return value.decode("utf-8", errors="replace")


def run_command(
    command: list[str],
    cwd: Path | None = None,
    extra_env: dict[str, str] | None = None,
    timeout: float = DEFAULT_SUBPROCESS_TIMEOUT_SECONDS,
) -> tuple[int, str, str]:
    env = os.environ.copy()
    env.setdefault("PYTHONUTF8", "1")
    if extra_env:
        env.update(extra_env)
    try:
        result = subprocess.run(
            command,
            cwd=cwd or repo_root(),
            capture_output=True,
            text=False,
            env=env,
            timeout=max(0.1, float(timeout)),
        )
    except subprocess.TimeoutExpired as exc:
        stderr = decode_output(exc.stderr)
        suffix = f"subprocess timed out after {timeout:g}s"
        return 124, _redact_output(decode_output(exc.stdout), env), _redact_output(
            f"{suffix}: {stderr}" if stderr else suffix, env
        )
    return result.returncode, _redact_output(decode_output(result.stdout), env), _redact_output(
        decode_output(result.stderr), env
    )


def slugify(value: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "-", value.strip().lower()).strip("-")
    return (slug or f"project-{datetime.now().strftime('%Y%m%d-%H%M%S')}")[:63].rstrip("-")


def safe_filename(value: str) -> str:
    name = Path(value or "finished.pptx").name
    stem = re.sub(r'[<>:"/\\|?*\x00-\x1f]+', "-", Path(name).stem).strip("-. ") or "finished"
    return f"{stem[:96]}.pptx"


def natural_sort_key(value: str) -> list[Any]:
    return [int(part) if part.isdigit() else part.lower() for part in re.split(r"(\d+)", value)]


def ensure_pptx(data: bytes) -> None:
    try:
        with zipfile.ZipFile(BytesIO(data)) as archive:
            if "ppt/presentation.xml" not in archive.namelist():
                raise ValueError("上传文件不是有效的 PowerPoint PPTX 文件。")
    except zipfile.BadZipFile as exc:
        raise ValueError("上传文件不是有效的 PowerPoint PPTX 文件。") from exc


def within(root: Path, candidate: Path) -> Path:
    root, candidate = root.resolve(), candidate.resolve()
    try:
        candidate.relative_to(root)
    except ValueError as exc:
        raise ValueError("路径超出允许的项目目录。") from exc
    return candidate


def buyer_records(data: Any) -> list[dict[str, Any]]:
    if isinstance(data, list):
        return [dict(item) for item in data if isinstance(item, dict)]
    if isinstance(data, dict) and isinstance(data.get("records"), list):
        return [dict(item) for item in data["records"] if isinstance(item, dict)]
    return []


def record_globals(data: Any) -> dict[str, Any]:
    if isinstance(data, dict) and isinstance(data.get("globals"), dict):
        return dict(data["globals"])
    return {}


def record_count(data: Any) -> int:
    if isinstance(data, dict) and isinstance(data.get("pages"), list):
        return sum(len(page.get("buyers") or []) for page in data["pages"] if isinstance(page, dict))
    return len(buyer_records(data))


def briefing_pages_from_buyers(buyers: list[dict[str, Any]], title: str) -> list[dict[str, Any]]:
    pages: list[dict[str, Any]] = []
    for offset in range(0, len(buyers), 6):
        entries = []
        for buyer in buyers[offset : offset + 6]:
            name = str(buyer.get("name", "") or "").strip()
            bio = str(buyer.get("bio", "") or "").strip()
            products = str(buyer.get("products", "") or "").strip()
            entries.append({
                "name": name,
                "summary": bio if bio.startswith(name) else f"{name}{bio}",
                "products": products if products.startswith("采购产品") else f"采购产品：{products}",
            })
        pages.append({"title": title or "买家商情", "buyers": entries})
    return pages


def presentation_engine(value: Any) -> str:
    engine = str(value or "auto").strip().lower()
    return engine if engine in {"auto", "office", "wps", "compatible"} else "auto"


def inspect_output_ppt(path: Path) -> dict[str, Any]:
    if not path.is_file() or path.stat().st_size <= 0:
        raise RuntimeError("生成进程已结束，但没有找到有效的PPTX输出文件。")
    try:
        from pptx import Presentation
        slide_count = len(Presentation(path).slides)
    except Exception as exc:
        raise RuntimeError(f"PPTX输出文件无法读取：{exc}") from exc
    return {
        "output": str(path),
        "output_size": path.stat().st_size,
        "slide_count": slide_count,
    }


def is_buyer_layout(config: Any) -> bool:
    return isinstance(config, dict) and isinstance(config.get("cover"), dict) and isinstance(config.get("content"), dict)


def is_buyer_briefing_layout(config: Any) -> bool:
    return (
        isinstance(config, dict)
        and "title_shape" in config
        and isinstance(config.get("slots"), list)
        and len(config.get("slots") or []) > 0
    )


def has_buyer_image_slots(config: Any) -> bool:
    if not isinstance(config, dict):
        return False
    slides = config.get("images", {}).get("slides", []) if isinstance(config.get("images"), dict) else []
    return any(isinstance(slot, dict) and (slot.get("logo") or slot.get("site")) for slot in slides)


def default_research_strategy(procurement_need: str) -> dict[str, Any]:
    need = procurement_need.strip().lower()
    defaults: dict[str, Any] = {
        "preferred_industries": "当地制造企业、工程项目承包商、设备维护服务商、进口商和区域经销商",
        "excluded_company_types": "仅销售同类产品且无进口、代理或自用场景的直接竞争制造商；无当地实体或无官网企业",
        "custom_requirements": "企业需在目标国家有工厂、项目、服务网点或明确经销业务；优先有进口、代理、跨境采购或持续使用需求的公开证据。",
        "prefer_import_evidence": True,
        "candidate_multiplier": 3,
    }
    if any(word in need for word in ("电机", "马达", "motor")):
        defaults.update({
            "preferred_industries": "食品机械、包装机械、泵阀、风机、输送设备、暖通设备、矿山设备、物流仓储、工业自动化、工程承包和设备维护",
            "excluded_company_types": "纯电机、发电机或减速电机制造商；仅销售电机且无进口代理或自用场景的品牌商；无当地实体企业",
            "custom_requirements": "优先选择生产线、设备或项目必然使用电机的当地企业，以及进口商、代理商和维修服务商；必须说明电机的具体使用环节，并优先有进口或持续采购证据。",
        })
    elif any(word in need for word in ("五轴", "数控", "cnc", "机床", "加工中心")):
        defaults.update({
            "preferred_industries": "航空航天零部件、汽车零部件、模具制造、医疗器械加工、精密机械、能源设备零部件、合同制造和CNC加工服务商",
            "excluded_company_types": "五轴数控机床制造商、纯机床品牌商；无本地工厂、加工服务或经销网点的海外企业",
            "custom_requirements": "企业必须在当地有工厂、机加工服务、维修网点或进口代理业务；优先有复杂零件加工、模具加工、五轴加工能力或进口机床代理的公开证据。",
        })
    return defaults


def merge_research_strategy(procurement_need: str, supplied: dict[str, Any] | None) -> dict[str, Any]:
    merged = default_research_strategy(procurement_need)
    for key, value in (supplied or {}).items():
        if key in {"preferred_industries", "excluded_company_types", "custom_requirements"} and str(value or "").strip():
            merged[key] = str(value).strip()
        elif key == "prefer_import_evidence":
            merged[key] = bool(value)
        elif key == "candidate_multiplier":
            try:
                merged[key] = max(2, min(5, int(value)))
            except (TypeError, ValueError):
                pass
    return merged


def _u(value: str) -> str:
    return value.encode("ascii").decode("unicode_escape")


def default_briefing_records(name: str) -> dict[str, Any]:
    sample_buyers = [
        {
            "name": _u(r"\u793a\u4f8b\u4f01\u4e1aA"),
            "summary": _u(r"\u793a\u4f8b\u4f01\u4e1aA\u662f\u5f53\u5730\u5927\u578b\u98df\u54c1\u5236\u9020\u4f01\u4e1a\uff0c\u8986\u76d6\u591a\u7c7b\u52a0\u5de5\u573a\u666f\uff0c\u5177\u5907\u7a33\u5b9a\u8bbe\u5907\u91c7\u8d2d\u9700\u6c42\u3002"),
            "products": _u(r"\u91c7\u8d2d\u54c1\u7c7b\uff1a\u793a\u4f8b\u8bbe\u5907\u3001\u914d\u5957\u4ea7\u7ebf\u3001\u5907\u4ef6\u7b49"),
        },
        {
            "name": _u(r"\u793a\u4f8b\u4f01\u4e1aB"),
            "summary": _u(r"\u793a\u4f8b\u4f01\u4e1aB\u662f\u533a\u57df\u6027\u5de5\u4e1a\u5ba2\u6237\uff0c\u62e5\u6709\u591a\u5904\u751f\u4ea7\u6216\u670d\u52a1\u7f51\u70b9\uff0c\u91cd\u89c6\u8bbe\u5907\u6548\u7387\u4e0e\u552e\u540e\u4fdd\u969c\u3002"),
            "products": _u(r"\u91c7\u8d2d\u54c1\u7c7b\uff1a\u793a\u4f8b\u8bbe\u5907\u3001\u81ea\u52a8\u5316\u7cfb\u7edf\u3001\u8017\u6750\u7b49"),
        },
        {
            "name": _u(r"\u793a\u4f8b\u4f01\u4e1aC"),
            "summary": _u(r"\u793a\u4f8b\u4f01\u4e1aC\u662f\u7ec6\u5206\u884c\u4e1a\u91c7\u8d2d\u5546\uff0c\u670d\u52a1\u591a\u4e2a\u7ec8\u7aef\u5e02\u573a\uff0c\u6301\u7eed\u5173\u6ce8\u53ef\u9760\u8bbe\u5907\u4e0e\u672c\u5730\u670d\u52a1\u3002"),
            "products": _u(r"\u91c7\u8d2d\u54c1\u7c7b\uff1a\u793a\u4f8b\u8bbe\u5907\u3001\u63a7\u5236\u7cfb\u7edf\u3001\u5907\u4ef6\u7b49"),
        },
        {
            "name": _u(r"\u793a\u4f8b\u4f01\u4e1aD"),
            "summary": _u(r"\u793a\u4f8b\u4f01\u4e1aD\u662f\u672c\u5730\u9886\u5148\u8fd0\u8425\u5546\u4e4b\u4e00\uff0c\u4e1a\u52a1\u8986\u76d6\u751f\u4ea7\u3001\u4ed3\u50a8\u6216\u5de5\u7a0b\u573a\u666f\uff0c\u91c7\u8d2d\u54c1\u7c7b\u8f83\u96c6\u4e2d\u3002"),
            "products": _u(r"\u91c7\u8d2d\u54c1\u7c7b\uff1a\u793a\u4f8b\u8bbe\u5907\u3001\u8f85\u52a9\u8bbe\u5907\u3001\u7ef4\u62a4\u670d\u52a1\u7b49"),
        },
        {
            "name": _u(r"\u793a\u4f8b\u4f01\u4e1aE"),
            "summary": _u(r"\u793a\u4f8b\u4f01\u4e1aE\u662f\u591a\u5143\u5316\u96c6\u56e2\u4e0b\u5c5e\u4e1a\u52a1\u5355\u4f4d\uff0c\u5177\u5907\u89c4\u6a21\u5316\u91c7\u8d2d\u80fd\u529b\uff0c\u5173\u6ce8\u8bbe\u5907\u7a33\u5b9a\u6027\u3002"),
            "products": _u(r"\u91c7\u8d2d\u54c1\u7c7b\uff1a\u793a\u4f8b\u8bbe\u5907\u3001\u4ea7\u7ebf\u7ec4\u4ef6\u3001\u6613\u635f\u4ef6\u7b49"),
        },
        {
            "name": _u(r"\u793a\u4f8b\u4f01\u4e1aF"),
            "summary": _u(r"\u793a\u4f8b\u4f01\u4e1aF\u662f\u884c\u4e1a\u6e20\u9053\u6216\u7ec8\u7aef\u5ba2\u6237\uff0c\u957f\u671f\u670d\u52a1\u672c\u5730\u5e02\u573a\uff0c\u5bf9\u4ea4\u4ed8\u5468\u671f\u548c\u6280\u672f\u652f\u6301\u8981\u6c42\u8f83\u9ad8\u3002"),
            "products": _u(r"\u91c7\u8d2d\u54c1\u7c7b\uff1a\u793a\u4f8b\u8bbe\u5907\u3001\u914d\u5957\u4ea7\u54c1\u3001\u5907\u4ef6\u7b49"),
        },
    ]
    return {"globals": {"deck_title": name, "country": "", "procurement_need": ""}, "pages": [{"title": _u(r"\u54c1\u7c7b\u540d\u79f0"), "buyers": sample_buyers}]}

class ConsoleState:
    def __init__(self, projects_root: Path):
        self.projects_root = projects_root.resolve()
        self.projects_root.mkdir(parents=True, exist_ok=True)
        self.jobs: dict[str, dict[str, Any]] = {}
        self.jobs_lock = threading.Lock()
        self.model_settings: dict[str, Any] = {
            "unified": True,
            "unified_provider": "deepseek",
            "unified_base_url": MODEL_PROVIDER_DEFAULTS["deepseek"]["base_url"],
            "unified_key": "",
            "research_provider": "deepseek",
            "research_base_url": MODEL_PROVIDER_DEFAULTS["deepseek"]["base_url"],
            "research_key": "",
            "visual_provider": "openai",
            "visual_base_url": MODEL_PROVIDER_DEFAULTS["openai"]["base_url"],
            "visual_key": "",
            "layout_provider": "deepseek",
            "layout_base_url": MODEL_PROVIDER_DEFAULTS["deepseek"]["base_url"],
            "layout_key": "",
            "research_model": "deepseek-chat",
            "visual_model": "gpt-image-1",
            "layout_model": "deepseek-chat",
            "research_enabled": True,
            "research_mode": "model_only",
            "visual_enabled": False,
            "layout_enabled": False,
        }

    def project_dir(self, slug: str) -> Path:
        if not PROJECT_SLUG_RE.fullmatch(slug):
            raise ValueError("项目标识不合法。")
        return within(self.projects_root, self.projects_root / slug)

    def paths(self, slug: str) -> dict[str, Path]:
        root = self.project_dir(slug)
        return {
            "root": root,
            "project": root / "project.json",
            "template": root / "template.pptx",
            "records": root / "records.json",
            "layout": root / "layout-config.json",
            "assets": root / "assets",
            "output": root / "output",
            "workspace": root / "workspace",
            "imports": root / "imports",
        }

    def require(self, slug: str) -> dict[str, Path]:
        paths = self.paths(slug)
        if not paths["project"].is_file():
            raise FileNotFoundError("项目不存在。")
        return paths

    def output_files(self, paths: dict[str, Path]) -> list[dict[str, Any]]:
        if not paths["output"].exists():
            return []
        result = []
        items = sorted(paths["output"].glob("*.pptx"), key=lambda item: item.stat().st_mtime, reverse=True)
        for path in items:
            stat = path.stat()
            preview_dir = paths["output"] / ".previews" / path.stem
            preview_count = len(list(preview_dir.glob("*.png"))) if preview_dir.is_dir() else 0
            result.append({
                "name": path.name,
                "size": stat.st_size,
                "modified_at": datetime.fromtimestamp(stat.st_mtime, timezone.utc).isoformat(timespec="seconds"),
                "preview_count": preview_count,
            })
        return result

    def unique_output_filename(self, paths: dict[str, Path], requested: str) -> str:
        filename = safe_filename(requested)
        if not (paths["output"] / filename).exists():
            return filename
        stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
        source = Path(filename)
        candidate = f"{source.stem}-{stamp}{source.suffix}"
        counter = 2
        while (paths["output"] / candidate).exists():
            candidate = f"{source.stem}-{stamp}-{counter}{source.suffix}"
            counter += 1
        return candidate

    def list_projects(self) -> list[dict[str, Any]]:
        result = []
        for project_file in self.projects_root.glob("*/project.json"):
            try:
                item = read_json(project_file)
                paths = self.paths(str(item["slug"]))
                item["template_ready"] = paths["template"].is_file()
                item["outputs"] = self.output_files(paths)
                result.append(item)
            except (OSError, ValueError, KeyError, json.JSONDecodeError):
                continue
        return sorted(result, key=lambda item: str(item.get("updated_at", "")), reverse=True)

    def create_project(self, name: str, mode: str = "generic") -> dict[str, Any]:
        name = name.strip()
        mode = mode if mode in {"generic", "buyer_board", "buyer_briefing"} else "generic"
        if not name:
            raise ValueError("项目名称不能为空。")
        base, slug, suffix = slugify(name), slugify(name), 2
        while self.project_dir(slug).exists():
            tail = f"-{suffix}"
            slug = f"{base[:63 - len(tail)]}{tail}"
            suffix += 1
        paths = self.paths(slug)
        for key in ("root", "assets", "output", "workspace"):
            paths[key].mkdir(parents=True, exist_ok=True)
        now = utc_now()
        inferred_mode = mode
        project = {
            "version": 1,
            "name": name,
            "slug": slug,
            "mode": inferred_mode,
            "created_at": now,
            "updated_at": now,
            "paths": {
                "template": "template.pptx",
                "records": "records.json",
                "layout_config": "layout-config.json",
                "assets": "assets",
                "output": "output",
                "workspace": "workspace",
            },
            "export": {"filename": "finished.pptx", "strict": False, "presentation_engine": "auto"},
            "last_runs": [],
            "research_strategy": {
                "preferred_industries": "",
                "excluded_company_types": "",
                "custom_requirements": "",
                "prefer_import_evidence": True,
                "candidate_multiplier": 3,
            },
            "layout_instruction": "",
        }
        write_json(paths["project"], project)
        if inferred_mode == "buyer_briefing":
            write_json(paths["records"], default_briefing_records(name))
            write_json(paths["layout"], BUYER_BRIEFING_DEFAULT_MAPPING)
        else:
            write_json(paths["records"], {"globals": {"deck_title": name, "country": "", "procurement_need": ""}, "records": []})
            write_json(paths["layout"], {"version": 1, "record_key": "records", "required_fields": [], "slides": []})
        return project

    def inspect_template(self, path: Path) -> dict[str, Any]:
        if not path.is_file():
            return {"ready": False, "slides": [], "slide_count": 0}
        try:
            from pptx import Presentation
        except ImportError:
            return {"ready": True, "error": "未安装 python-pptx，暂时无法读取模板结构。", "slides": []}
        presentation = Presentation(path)
        slides = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            shapes = []
            for shape_index, shape in enumerate(slide.shapes, start=1):
                text = str(shape.text).replace("\n", " ").strip()[:80] if getattr(shape, "has_text_frame", False) else ""
                shapes.append({
                    "index": shape_index,
                    "id": int(shape.shape_id),
                    "name": shape.name,
                    "type": str(shape.shape_type),
                    "text": text,
                    "has_text_frame": bool(getattr(shape, "has_text_frame", False)),
                    "has_table": bool(getattr(shape, "has_table", False)),
                    "left": round(shape.left / 914400, 3),
                    "top": round(shape.top / 914400, 3),
                    "width": round(shape.width / 914400, 3),
                    "height": round(shape.height / 914400, 3),
                })
            slides.append({"index": slide_index, "shape_count": len(shapes), "shapes": shapes})
        return {
            "ready": True,
            "slide_count": len(slides),
            "width": round(presentation.slide_width / 914400, 3),
            "height": round(presentation.slide_height / 914400, 3),
            "slides": slides,
        }

    def role_model_config(self, role: str, overrides: dict[str, Any] | None = None) -> dict[str, str]:
        if role not in ROLE_CHOICES:
            raise ValueError("Unsupported model role.")
        settings = self.model_settings
        overrides = overrides or {}
        if bool(settings.get("unified", True)) and not overrides:
            provider = clean_provider(settings.get("unified_provider"))
            key = str(settings.get("unified_key") or os.environ.get("OPENAI_API_KEY", "")).strip()
            base_url = clean_base_url(settings.get("unified_base_url"), provider)
        else:
            provider = clean_provider(overrides.get("provider", settings.get(f"{role}_provider")))
            key = str(overrides.get("api_key") or "").strip()
            if not key:
                key = str(settings.get("unified_key") if settings.get("unified") else settings.get(f"{role}_key", "") or "").strip()
            key = key or os.environ.get("OPENAI_API_KEY", "")
            base_url = clean_base_url(overrides.get("base_url", settings.get(f"{role}_base_url")), provider)
        model = str(overrides.get("model", settings.get(f"{role}_model", "")) or "").strip()
        if not model:
            model = MODEL_PROVIDER_DEFAULTS.get(provider, MODEL_PROVIDER_DEFAULTS["compatible"]).get(f"{role}_model", "")
        inferred_provider = infer_provider_from_model(model, provider, bool(overrides.get("provider")))
        if inferred_provider != provider:
            provider = inferred_provider
            base_url = clean_base_url("", provider)
        return {"provider": provider, "api_key": key, "base_url": base_url, "model": model}

    def model_settings_payload(self) -> dict[str, Any]:
        settings = self.model_settings
        payload = {
            "unified": bool(settings["unified"]),
            "providers": {key: dict(value) for key, value in MODEL_PROVIDER_DEFAULTS.items()},
            "unified_provider": clean_provider(settings.get("unified_provider")),
            "unified_base_url": clean_base_url(settings.get("unified_base_url"), clean_provider(settings.get("unified_provider"))),
            "unified_configured": bool(settings["unified_key"] or os.environ.get("OPENAI_API_KEY")),
            "research_enabled": bool(settings.get("research_enabled", True)),
            "research_mode": clean_research_mode(settings.get("research_mode")),
            "visual_enabled": bool(settings.get("visual_enabled", False)),
            "layout_enabled": bool(settings.get("layout_enabled", False)),
        }
        for role in ROLE_CHOICES:
            cfg = self.role_model_config(role)
            configured = bool(cfg["api_key"]) if not settings.get("unified") else bool(settings["unified_key"] or os.environ.get("OPENAI_API_KEY"))
            payload.update({
                f"{role}_provider": cfg["provider"],
                f"{role}_base_url": cfg["base_url"],
                f"{role}_configured": configured,
                f"{role}_model": cfg["model"],
            })
        return redact_credentials(payload)

    def configure_models(self, payload: dict[str, Any]) -> dict[str, Any]:
        self.model_settings["unified"] = bool(payload.get("unified", True))
        if "research_enabled" in payload:
            self.model_settings["research_enabled"] = bool(payload.get("research_enabled"))
        if "visual_enabled" in payload:
            self.model_settings["visual_enabled"] = bool(payload.get("visual_enabled"))
        if "layout_enabled" in payload:
            self.model_settings["layout_enabled"] = bool(payload.get("layout_enabled"))
        if "research_mode" in payload:
            self.model_settings["research_mode"] = clean_research_mode(payload.get("research_mode"))
        for prefix in ("unified", "research", "visual", "layout"):
            provider_key = f"{prefix}_provider"
            base_url_key = f"{prefix}_base_url"
            if provider_key in payload:
                provider = clean_provider(payload.get(provider_key))
                self.model_settings[provider_key] = provider
                self.model_settings[base_url_key] = clean_base_url(payload.get(base_url_key), provider)
            elif base_url_key in payload:
                provider = clean_provider(self.model_settings.get(provider_key))
                self.model_settings[base_url_key] = clean_base_url(payload.get(base_url_key), provider)
        for key in ("unified_key", "research_key", "visual_key", "layout_key"):
            value = str(payload.get(key, "") or "").strip()
            if value:
                self.model_settings[key] = value
        for key in ("research_model", "visual_model", "layout_model"):
            value = str(payload.get(key, "") or "").strip()
            if value:
                self.model_settings[key] = value
        inferred = infer_provider_from_model(self.model_settings.get("research_model"), clean_provider(self.model_settings.get("unified_provider")), bool(payload.get("unified_provider")))
        if self.model_settings.get("unified") and inferred != self.model_settings.get("unified_provider"):
            self.model_settings["unified_provider"] = inferred
            self.model_settings["unified_base_url"] = clean_base_url("", inferred)
        return self.model_settings_payload()

    def diagnose_model_connection(self, payload: dict[str, Any]) -> dict[str, Any]:
        role = str(payload.get("role") or "research").strip()
        config = self.role_model_config(role, {
            "provider": payload.get("provider"),
            "base_url": payload.get("base_url"),
            "api_key": payload.get("api_key"),
        })
        base_url = config["base_url"].rstrip("/")
        parent_probe = probe_base_url(base_url)
        child_code = (
            "import json, urllib.request, urllib.error\n"
            "url = " + repr(base_url) + "\n"
            "req = urllib.request.Request(url, headers={'Accept':'application/json'}, method='GET')\n"
            "try:\n"
            "    r = urllib.request.urlopen(req, timeout=8)\n"
            "    print(json.dumps({'ok': True, 'reachable': True, 'status': getattr(r, 'status', 0), 'error': ''}, ensure_ascii=False))\n"
            "except urllib.error.HTTPError as e:\n"
            "    print(json.dumps({'ok': True, 'reachable': True, 'status': e.code, 'error': str(e.reason)}, ensure_ascii=False))\n"
            "except Exception as e:\n"
            "    print(json.dumps({'ok': False, 'reachable': False, 'status': 0, 'error': str(e)}, ensure_ascii=False))\n"
        )
        returncode, stdout, stderr = run_command(
            [sys.executable, "-c", child_code],
            extra_env=self.model_env(role),
            timeout=DIAGNOSTIC_SUBPROCESS_TIMEOUT_SECONDS,
        )
        try:
            child_probe = json.loads(stdout.strip() or "{}")
        except json.JSONDecodeError:
            child_probe = {"ok": False, "reachable": False, "status": 0, "error": (stderr or stdout).strip()}
        return redact_credentials({
            "ok": bool(parent_probe.get("reachable") or child_probe.get("reachable")),
            "provider": config["provider"],
            "base_url": base_url,
            "model": config["model"],
            "python": sys.executable,
            "parent_process_probe": parent_probe,
            "child_process_probe": child_probe,
            "child_returncode": returncode,
            "proxy_env": {
                "HTTP_PROXY": os.environ.get("HTTP_PROXY", ""),
                "HTTPS_PROXY": os.environ.get("HTTPS_PROXY", ""),
                "NO_PROXY": os.environ.get("NO_PROXY", ""),
            },
            "hint": "HTTP 401/403 also means the network path is reachable; it only means the probe did not send a valid API request.",
        })

    def list_models(self, payload: dict[str, Any]) -> dict[str, Any]:
        role = str(payload.get("role") or "research").strip()
        config = self.role_model_config(role, {
            "provider": payload.get("provider"),
            "base_url": payload.get("base_url"),
            "api_key": payload.get("api_key"),
        })
        provider = config["provider"]
        defaults = MODEL_PROVIDER_DEFAULTS.get(provider, MODEL_PROVIDER_DEFAULTS["compatible"])
        fallback = list(defaults.get("models") or [])
        if not fallback:
            fallback = [value for key, value in defaults.items() if key.endswith("_model") and value]
        if not config["api_key"]:
            return {"ok": False, "provider": provider, "models": fallback, "error": "API Key is required to fetch upstream models."}
        request = Request(
            model_list_endpoint(config["base_url"]),
            headers={"Authorization": f"Bearer {config['api_key']}", "Accept": "application/json"},
            method="GET",
        )
        try:
            with urlopen(request, timeout=18) as response:
                data = json.loads(response.read().decode("utf-8"))
        except Exception as exc:
            return redact_credentials({"ok": False, "provider": provider, "models": fallback, "error": str(exc)})
        raw_models = data.get("data", data if isinstance(data, list) else [])
        models: list[str] = []
        if isinstance(raw_models, list):
            for item in raw_models:
                if isinstance(item, dict) and item.get("id"):
                    models.append(str(item["id"]))
                elif isinstance(item, str):
                    models.append(item)
        models = sorted(dict.fromkeys(models))
        return redact_credentials({"ok": True, "provider": provider, "base_url": config["base_url"], "models": models or fallback})

    def model_env(self, role: str) -> dict[str, str]:
        config = self.role_model_config(role)
        env: dict[str, str] = {}
        if config["api_key"]:
            env["OPENAI_API_KEY"] = config["api_key"]
        if config["base_url"]:
            env["OPENAI_BASE_URL"] = config["base_url"]
        if role == "research":
            env["BUYER_RESEARCH_PROVIDER"] = config["provider"]
            env["BUYER_RESEARCH_BASE_URL"] = config["base_url"]
            env["BUYER_RESEARCH_MODEL"] = config["model"]
            env["BUYER_RESEARCH_MODE"] = clean_research_mode(self.model_settings.get("research_mode"))
        elif role == "visual":
            env["BUYER_VISUAL_PROVIDER"] = config["provider"]
            env["BUYER_VISUAL_BASE_URL"] = config["base_url"]
            env["BUYER_VISUAL_MODEL"] = config["model"]
        return env

    def payload(self, slug: str) -> dict[str, Any]:
        paths = self.require(slug)
        project = read_json(paths["project"])
        changed = False
        if "mode" not in project:
            project["mode"] = "buyer_board" if re.search(r"买家|采购商|采购", str(project.get("name", ""))) else "generic"
            changed = True
        if "research" not in project:
            project["research"] = {"country": "", "procurement_need": "", "buyer_count": 10}
            changed = True
        if "research_strategy" not in project:
            project["research_strategy"] = {
                "preferred_industries": "",
                "excluded_company_types": "",
                "custom_requirements": "",
                "prefer_import_evidence": True,
                "candidate_multiplier": 3,
            }
            changed = True
        if "layout_instruction" not in project:
            project["layout_instruction"] = ""
            changed = True
        if "layout_recipes" not in project:
            project["layout_recipes"] = []
            changed = True
        if changed:
            write_json(paths["project"], project)
        return {
            "project": project,
            "records": read_json(paths["records"]),
            "layout_config": read_json(paths["layout"]),
            "template": self.inspect_template(paths["template"]),
            "outputs": self.output_files(paths),
            "model_settings": self.model_settings_payload(),
        }

    def touch(self, paths: dict[str, Path], updates: dict[str, Any] | None = None) -> dict[str, Any]:
        project = read_json(paths["project"])
        if updates:
            project.update(updates)
        project["updated_at"] = utc_now()
        write_json(paths["project"], project)
        return project

    def save_document(self, slug: str, kind: str, payload: Any) -> None:
        paths = self.require(slug)
        write_json(paths["records"] if kind == "records" else paths["layout"], payload)
        self.touch(paths)

    def mapping_preview(self, slug: str, config: dict[str, Any] | None = None) -> dict[str, Any]:
        paths = self.require(slug)
        config = config if isinstance(config, dict) else read_json(paths["layout"])
        inspection = self.inspect_template(paths["template"])
        data = read_json(paths["records"])
        records = buyer_records(data)
        globals_data = record_globals(data)
        shapes_by_slide = {
            int(slide.get("index", 0)): {int(shape.get("index", 0)): shape for shape in slide.get("shapes", [])}
            for slide in inspection.get("slides", [])
        }
        entries: list[dict[str, Any]] = []
        warnings: list[str] = []

        def sample_value(field: str) -> str:
            source: Any = globals_data if field.startswith("globals.") else (records[0] if records else {})
            key = field.split(".", 1)[1] if "." in field else field
            if not isinstance(source, dict):
                return ""
            for part in key.split("."):
                if not isinstance(source, dict) or part not in source:
                    return ""
                source = source[part]
            return str(source or "")

        def collect(slide_index: int, specs: Any, kind: str) -> None:
            if not isinstance(specs, list):
                return
            for spec in specs:
                if not isinstance(spec, dict):
                    continue
                field = str(spec.get("field") or spec.get("template") or spec.get("value") or "固定内容")
                shape_index = spec.get("shape_index")
                shape = shapes_by_slide.get(slide_index, {}).get(int(shape_index or 0))
                target = f"第 {slide_index} 页 · 元素 {shape_index}" if shape_index else f"第 {slide_index} 页 · 自定义区域"
                if shape_index and shape is None:
                    warnings.append(f"{target} 不存在，请重新选择模板元素。")
                entries.append({
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "shape_name": str((shape or {}).get("name") or "自定义区域"),
                    "kind": kind,
                    "field": field,
                    "sample": sample_value(field)[:80],
                    "target": target,
                })

        for slide in config.get("slides", []) if isinstance(config.get("slides"), list) else []:
            if isinstance(slide, dict):
                index = int(slide.get("slide_index") or 0)
                collect(index, slide.get("texts"), "文本")
                collect(index, slide.get("images"), "图片")
        repeat = config.get("repeat") if isinstance(config.get("repeat"), dict) else {}
        if repeat:
            index = int(repeat.get("source_slide_index") or repeat.get("start_slide_index") or 0)
            collect(index, repeat.get("texts"), "批量文本")
            collect(index, repeat.get("images"), "批量图片")
        if not entries:
            warnings.append("尚未找到可预览的映射。请先用自然语言生成映射或选择一个已保存方案。")
        return {
            "entries": entries,
            "warnings": list(dict.fromkeys(warnings)),
            "record_count": len(records),
            "template_ready": bool(inspection.get("ready")),
        }

    def export_preflight(self, slug: str) -> dict[str, Any]:
        paths = self.require(slug)
        project = read_json(paths["project"])
        config = read_json(paths["layout"])
        data = read_json(paths["records"])
        records = buyer_records(data)
        mode = str(project.get("mode", "generic"))
        if mode == "buyer_board":
            return {
                "ok": bool(paths["template"].is_file() and records),
                "errors": ([] if paths["template"].is_file() else ["尚未上传 PPTX 模板。"])
                + ([] if records else ["尚未填写买家资料。"]),
                "warnings": [],
                "mapping_count": 0,
                "record_count": len(records),
                "expected_slide_count": 1 + len(records),
            }
        if mode == "buyer_briefing":
            pages = data.get("pages", []) if isinstance(data, dict) else []
            return {
                "ok": bool(paths["template"].is_file() and pages),
                "errors": ([] if paths["template"].is_file() else ["尚未上传 PPTX 模板。"])
                + ([] if pages else ["尚未填写买家商情页面。"]),
                "warnings": [],
                "mapping_count": 0,
                "record_count": record_count(data),
                "expected_slide_count": len(pages),
            }
        preview = self.mapping_preview(slug, config)
        errors: list[str] = []
        warnings = list(preview["warnings"])
        if not paths["template"].is_file():
            errors.append("尚未上传 PPTX 模板。")
        if project.get("mode") == "generic" and not records:
            errors.append("尚未导入可批量填充的资料。")
        if project.get("mode") == "generic" and not preview["entries"]:
            errors.append("尚未配置字段映射。")
        required = [str(value) for value in config.get("required_fields", []) if str(value).strip()]
        for record_index, record in enumerate(records, start=1):
            for field in required:
                key = field.split(".")[-1]
                if not str(record.get(key, "") or "").strip():
                    errors.append(f"第 {record_index} 条资料缺少必填字段：{field}。")
        for entry in preview["entries"]:
            if entry["kind"].endswith("图片") and entry["sample"] and not Path(entry["sample"]).is_file():
                warnings.append(f"{entry['target']} 的示例图片路径不可用，导出时会保留空位。")
        mapped_fields = self.mapped_record_fields(config)
        for field in sorted(mapped_fields):
            missing_count = 0
            for record in records:
                value: Any = record
                for part in field.split("."):
                    value = value.get(part) if isinstance(value, dict) else None
                if value is None or not str(value).strip():
                    missing_count += 1
            if missing_count:
                warnings.append(f"已映射字段“{field}”有 {missing_count}/{len(records)} 条资料为空。")
        for record_index, record in enumerate(records, start=1):
            for field, value in record.items():
                if not isinstance(value, str):
                    continue
                if "{{" in value or "}}" in value:
                    warnings.append(f"第 {record_index} 条资料的“{field}”仍含占位符标记。")
                if len(value) > 1200:
                    warnings.append(f"第 {record_index} 条资料的“{field}”超过 1200 个字符，建议确认版面容量。")
        repeat = config.get("repeat") if isinstance(config.get("repeat"), dict) else {}
        template_slide_count = len(self.inspect_template(paths["template"]).get("slides") or [])
        expected_slide_count = template_slide_count + max(len(records) - 1, 0) if repeat else max(template_slide_count, 1)
        return {
            "ok": not errors,
            "errors": list(dict.fromkeys(errors)),
            "warnings": list(dict.fromkeys(warnings)),
            "mapping_count": len(preview["entries"]),
            "record_count": len(records),
            "expected_slide_count": expected_slide_count,
        }

    def save_layout_recipe(self, slug: str, payload: dict[str, Any]) -> dict[str, Any]:
        paths = self.require(slug)
        name = str(payload.get("name", "") or "").strip()[:80]
        if not name:
            raise ValueError("请为版式方案填写名称。")
        project = read_json(paths["project"])
        recipes = [item for item in project.get("layout_recipes", []) if isinstance(item, dict)]
        recipe = {
            "id": uuid.uuid4().hex[:10],
            "name": name,
            "created_at": utc_now(),
            "instruction": str(payload.get("instruction", project.get("layout_instruction", "")) or ""),
            "layout_config": read_json(paths["layout"]),
        }
        recipes.insert(0, recipe)
        project["layout_recipes"] = recipes[:20]
        project["updated_at"] = utc_now()
        write_json(paths["project"], project)
        return {"recipe": recipe, "recipes": project["layout_recipes"]}

    def apply_layout_recipe(self, slug: str, recipe_id: str) -> dict[str, Any]:
        paths = self.require(slug)
        project = read_json(paths["project"])
        recipe = next((item for item in project.get("layout_recipes", []) if isinstance(item, dict) and item.get("id") == recipe_id), None)
        if not recipe or not isinstance(recipe.get("layout_config"), dict):
            raise FileNotFoundError("版式方案不存在。")
        write_json(paths["layout"], recipe["layout_config"])
        self.touch(paths, {"layout_instruction": str(recipe.get("instruction", "") or "")})
        return {"layout_config": recipe["layout_config"], "recipe": recipe, "preview": self.mapping_preview(slug, recipe["layout_config"])}

    def import_records_document(self, slug: str, data: bytes, filename: str, instruction: str = "") -> dict[str, Any]:
        paths = self.require(slug)
        source_name = safe_filename(filename).replace(".pptx", Path(filename or "document.txt").suffix or ".txt")
        suffix = Path(source_name).suffix.lower()
        if suffix not in {".txt", ".md", ".csv", ".json", ".docx", ".xlsx"}:
            raise ValueError("资料导入支持 TXT、Markdown、CSV、JSON、DOCX、XLSX。")
        paths["imports"].mkdir(parents=True, exist_ok=True)
        source_path = paths["imports"] / source_name
        source_path.write_bytes(data)
        project = read_json(paths["project"])
        if project.get("mode") in {"buyer_board", "buyer_briefing"}:
            return self.import_buyer_document(paths, project, source_path)
        command = [
            sys.executable,
            str(skill_script("import_content_document.py")),
            "--input", str(source_path),
            "--output", str(paths["records"]),
            "--project-name", str(project.get("name", "PPT项目")),
            "--instruction", instruction,
        ]
        returncode, stdout, stderr = run_command(command)
        if returncode != 0:
            raise RuntimeError(stderr or stdout or "资料导入失败。")
        self.touch(paths, {"mode": "generic" if project.get("mode") == "generic" else project.get("mode", "generic")})
        records = read_json(paths["records"])
        return {"records": records, "record_count": len(buyer_records(records)), "source": str(source_path.name)}

    def rename_generic_fields(self, slug: str, mapping: dict[str, Any]) -> dict[str, Any]:
        paths = self.require(slug)
        project = read_json(paths["project"])
        if project.get("mode") != "generic":
            raise ValueError("Field mapping is available for generic PPT projects only.")
        data = read_json(paths["records"])
        records = buyer_records(data)
        normalized_mapping = {
            str(source).strip(): re.sub(r"[^\w\u4e00-\u9fff-]+", "_", str(target or "").strip())[:80]
            for source, target in mapping.items()
            if str(source).strip() and str(target or "").strip() and str(source).strip() != str(target).strip()
        }
        if not normalized_mapping:
            return {"records": data, "fields": self.field_candidates_from_records(data)}
        for record in records:
            for source, target in normalized_mapping.items():
                if source not in record:
                    continue
                value = record.pop(source)
                if target not in record or not str(record.get(target) or "").strip():
                    record[target] = value
        if isinstance(data, dict):
            data["records"] = records
        else:
            data = {"globals": {}, "records": records}
        write_json(paths["records"], data)
        self.touch(paths)
        return {"records": data, "fields": self.field_candidates_from_records(data)}

    def import_buyer_document(self, paths: dict[str, Path], project: dict[str, Any], source_path: Path) -> dict[str, Any]:
        current = read_json(paths["records"])
        globals_data = record_globals(current)
        imported_path = paths["workspace"] / f"imported-buyers-{uuid.uuid4().hex[:8]}.json"
        command = [
            sys.executable,
            str(skill_script("import_buyer_records.py")),
            "--input", str(source_path),
            "--output", str(imported_path),
            "--country", str(globals_data.get("country", "")),
            "--procurement-need", str(globals_data.get("procurement_need", "")),
        ]
        returncode, stdout, stderr = run_command(command)
        if returncode != 0:
            raise RuntimeError(stderr or stdout or "Buyer data import failed.")

        imported = read_json(imported_path)
        buyers = buyer_records(imported)
        imported_globals = record_globals(imported)
        country = str(imported_globals.get("country") or globals_data.get("country") or "")
        procurement_need = str(imported_globals.get("procurement_need") or globals_data.get("procurement_need") or "")
        if project.get("mode") == "buyer_board":
            result = self.save_buyer_data(
                str(project["slug"]),
                {"country": country, "procurement_need": procurement_need, "buyers": buyers},
            )
            result.update({"record_count": len(buyer_records(result["records"])), "source": str(source_path.name)})
            return result

        normalized, warnings = self.normalize_buyers(buyers, country, procurement_need)
        records = {
            "globals": {**globals_data, "country": country, "procurement_need": procurement_need},
            "pages": briefing_pages_from_buyers(normalized, procurement_need),
        }
        write_json(paths["records"], records)
        self.touch(paths, {"mode": "buyer_briefing"})
        return {"records": records, "record_count": len(normalized), "source": str(source_path.name), "warnings": warnings}

    def field_candidates_from_records(self, data: Any) -> list[str]:
        records = buyer_records(data)
        seen: list[str] = []
        for record in records:
            for key, value in record.items():
                if key in seen or isinstance(value, (dict, list)):
                    continue
                if str(value or "").strip():
                    seen.append(key)
        return seen or ["title", "content"]

    def mapped_record_fields(self, config: dict[str, Any]) -> set[str]:
        fields: set[str] = set()

        def collect_specs(specs: Any) -> None:
            if not isinstance(specs, list):
                return
            for spec in specs:
                if not isinstance(spec, dict):
                    continue
                for value in (spec.get("field"), spec.get("template")):
                    if not isinstance(value, str):
                        continue
                    if value.startswith("globals."):
                        continue
                    if value.startswith("record."):
                        fields.add(value.split(".", 1)[1])
                    elif value and "{" not in value:
                        fields.add(value)
                    for token in re.findall(r"\{record\.([^}]+)\}", value):
                        fields.add(token)

        for slide in config.get("slides", []) if isinstance(config.get("slides"), list) else []:
            if isinstance(slide, dict):
                collect_specs(slide.get("texts"))
                collect_specs(slide.get("images"))
                collect_specs(slide.get("tables"))
        repeat = config.get("repeat") if isinstance(config.get("repeat"), dict) else {}
        collect_specs(repeat.get("texts"))
        collect_specs(repeat.get("images"))
        collect_specs(repeat.get("tables"))
        return fields

    def generate_layout_from_instruction(self, slug: str, payload: dict[str, Any]) -> dict[str, Any]:
        paths = self.require(slug)
        instruction = str(payload.get("instruction", "") or "").strip()
        if not instruction:
            raise ValueError("请先填写自然语言版式要求。")
        records = read_json(paths["records"])
        project = read_json(paths["project"])
        fields = self.field_candidates_from_records(records)
        config: dict[str, Any] = {
            "version": 1,
            "record_key": "records",
            "required_fields": fields[:2],
            "notes": [
                "Generated from natural-language mapping instructions in the local console.",
                instruction,
                "This is a starter mapping. Check shape_index values in Template Structure before production export.",
            ],
        }
        inspection = self.inspect_template(paths["template"])
        if inspection.get("ready") and inspection.get("slides"):
            slides = inspection["slides"]
            first_slide_texts = [
                shape for shape in slides[0].get("shapes", [])
                if shape.get("has_text_frame") and shape.get("text") and not shape.get("has_table")
            ]
            if first_slide_texts:
                title_shape = max(first_slide_texts, key=lambda item: float(item.get("height") or 0))
                config["slides"] = [{
                    "slide_index": 1,
                    "texts": [{"shape_index": title_shape["index"], "field": "globals.deck_title"}],
                }]
            else:
                config["slides"] = []
            try:
                source_slide_index = int(payload.get("source_slide_index") or 0)
            except (TypeError, ValueError):
                source_slide_index = 0
            if source_slide_index < 1 or source_slide_index > len(slides):
                source_slide_index = 2 if len(slides) > 1 else 1
            repeat_slide = slides[source_slide_index - 1]
            text_shapes = [
                shape for shape in repeat_slide.get("shapes", [])
                if shape.get("has_text_frame") and not shape.get("has_table")
            ]
            text_shapes.sort(key=lambda item: (float(item.get("top") or 0), float(item.get("left") or 0)))
            repeat_texts = []
            for shape, field in zip(text_shapes, fields):
                repeat_texts.append({"shape_index": shape["index"], "field": field, "mode": "clear"})
            if repeat_texts:
                config["repeat"] = {
                    "source_slide_index": repeat_slide["index"],
                    "start_slide_index": repeat_slide["index"],
                    "template_slide_count": 1,
                    "trim_extra_template_slides": True,
                    "texts": repeat_texts,
                }
                image_fields = [field for field in fields if re.search(r"image|photo|picture|logo|图片|配图|照片", field, re.I)]
                picture_shapes = [shape for shape in repeat_slide.get("shapes", []) if "PICTURE" in str(shape.get("type", "")).upper()]
                if image_fields and picture_shapes:
                    config["repeat"]["images"] = [
                        {
                            "shape_index": shape["index"],
                            "field": field,
                            "fit": "contain" if "logo" in field.lower() else "cover",
                            "clear_if_missing": True,
                        }
                        for shape, field in zip(picture_shapes, image_fields)
                    ]
        else:
            config["slides"] = []
        write_json(paths["layout"], config)
        self.touch(paths, {"layout_instruction": instruction, "mode": project.get("mode", "generic")})
        return {"layout_config": config, "field_candidates": fields, "source_slide_index": repeat_slide["index"] if inspection.get("slides") else 0}

    def prepare_buyer_layout(self, paths: dict[str, Path], country: str, procurement_need: str) -> dict[str, Any]:
        if not paths["template"].is_file():
            return {"ready": False, "warning": "请先上传PPTX模板，再自动生成买家看板版式映射。"}
        title = f"{country}{procurement_need}买家".strip() or read_json(paths["project"]).get("name", "买家看板")
        country_line = f"国家：{country}" if country else ""
        current = read_json(paths["layout"]) if paths["layout"].is_file() else {}
        if is_buyer_layout(current):
            current.setdefault("content", {})
            # Existing buyer-board configs predate the fixed-header rule.
            # Preserve the template's content-page header unless a user
            # explicitly opts into replacing it in layout-config.json.
            current["content"].setdefault("preserve_title", True)
            current["content"].setdefault("dynamic_row_height", True)
            current.setdefault("defaults", {})
            current["defaults"].update({
                "cover_title": title,
                "cover_country": country_line,
                "content_title": title,
            })
            notes = current.get("notes") if isinstance(current.get("notes"), list) else []
            generated_scaffold = "This file is a scaffold generated from a reference PPT." in notes
            image_layout_version = int((current.get("images") or {}).get("layout_version", 0) or 0)
            if not has_buyer_image_slots(current) or (generated_scaffold and image_layout_version < 2):
                generated_path = paths["workspace"] / "buyer-layout-images.generated.json"
                command = [
                    sys.executable,
                    str(skill_script("generate_layout_config.py")),
                    "--template", str(paths["template"]),
                    "--output", str(generated_path),
                    "--cover-title", title,
                    "--cover-country", country_line,
                    "--content-title", title,
                ]
                returncode, stdout, stderr = run_command(command)
                if returncode == 0 and generated_path.is_file():
                    generated = read_json(generated_path)
                    if has_buyer_image_slots(generated):
                        current["images"] = generated["images"]
                        image_note = "Image insertion regions were generated from the buyer-board table layout."
                        if image_note not in current.setdefault("notes", []):
                            current["notes"].append(image_note)
                else:
                    return {"ready": False, "warning": "图片区域自动识别失败，请重新上传模板后重试。", "details": stderr or stdout}
            write_json(paths["layout"], current)
            return {"ready": True, "generated": False, "layout_config": current}

        command = [
            sys.executable,
            str(skill_script("generate_layout_config.py")),
            "--template", str(paths["template"]),
            "--output", str(paths["layout"]),
            "--cover-title", title,
            "--cover-country", country_line,
            "--content-title", title,
        ]
        returncode, stdout, stderr = run_command(command)
        if returncode != 0:
            return {
                "ready": False,
                "warning": "模板未能自动识别为买家看板版式，请在“版式映射”中手动配置。",
                "details": stderr or stdout,
            }
        return {"ready": True, "generated": True, "layout_config": read_json(paths["layout"])}

    def normalize_buyers(
        self,
        buyers: list[Any],
        default_country: str,
        procurement_need: str = "",
        enforce_research_copy_rules: bool = False,
    ) -> tuple[list[dict[str, Any]], list[str]]:
        normalized = []
        warnings = []
        for index, raw in enumerate(buyers, start=1):
            if not isinstance(raw, dict):
                continue
            item = {key: str(raw.get(key, "") or "").strip() for key in BUYER_FIELDS}
            source_urls = raw.get("source_urls")
            item["source_urls"] = [str(url).strip() for url in source_urls if str(url).strip()] if isinstance(source_urls, list) else []
            for score_key in BUYER_SCORE_FIELDS:
                try:
                    item[score_key] = max(0, min(100, int(raw.get(score_key, 0) or 0)))
                except (TypeError, ValueError):
                    item[score_key] = 0
            item["country"] = item["country"] or default_country
            if enforce_research_copy_rules:
                item["products"] = refine_buyer_products(
                    item["products"],
                    procurement_need,
                    buyer_evidence_context(item),
                )
                item["bio"] = pad_or_trim_bio(item["bio"])
                if item["products"] == "需核实具体设备":
                    warnings.append(f"{item['name'] or f'第{index}家买家'}采购品类目前只能确认到大类，建议人工补充具体设备。")
            if not item["name"]:
                warnings.append(f"第{index}家买家缺少企业名称。")
            chinese_count = sum(1 for char in item["bio"] if "\u4e00" <= char <= "\u9fff")
            if enforce_research_copy_rules and item["bio"] and not 120 <= chinese_count <= 130:
                label = item["name"] or f"第{index}家买家"
                warnings.append(f"{label}简介目前为{chinese_count}个中文字符，建议调整到120-130字。")
            normalized.append(item)
        return normalized, warnings

    def save_buyer_data(self, slug: str, payload: dict[str, Any]) -> dict[str, Any]:
        paths = self.require(slug)
        country = str(payload.get("country", "") or "").strip()
        procurement_need = str(payload.get("procurement_need", "") or "").strip()
        raw_buyers = payload.get("buyers")
        if not isinstance(raw_buyers, list):
            raise ValueError("买家资料必须是列表。")
        buyers, warnings = self.normalize_buyers(
            raw_buyers,
            country,
            procurement_need,
            bool(payload.get("enforce_research_copy_rules", False)),
        )
        current_data = read_json(paths["records"])
        globals_data = record_globals(current_data)
        project = read_json(paths["project"])
        globals_data.update({
            "deck_title": f"{country}{procurement_need}买家".strip() or project.get("name", "买家看板"),
            "country": country,
            "procurement_need": procurement_need,
        })
        records_payload = {"globals": globals_data, "records": buyers}
        write_json(paths["records"], records_payload)
        research = dict(project.get("research") or {})
        research.update({
            "country": country,
            "procurement_need": procurement_need,
            "buyer_count": max(len(buyers), int(research.get("buyer_count") or 10)),
        })
        self.touch(paths, {"mode": "buyer_board", "research": research})
        layout_result = self.prepare_buyer_layout(paths, country, procurement_need)
        if layout_result.get("warning"):
            warnings.append(str(layout_result["warning"]))
        return {
            "records": records_payload,
            "layout_config": read_json(paths["layout"]),
            "layout_ready": bool(layout_result.get("ready")),
            "warnings": warnings,
        }
    def save_template(self, slug: str, data: bytes) -> dict[str, Any]:
        if len(data) > MAX_UPLOAD_BYTES:
            raise ValueError("模板超过250MB限制。")
        ensure_pptx(data)
        paths = self.require(slug)
        paths["template"].write_bytes(data)
        project = self.touch(paths)
        inspection = self.inspect_template(paths["template"])
        if project.get("mode") == "buyer_board":
            globals_data = record_globals(read_json(paths["records"]))
            layout_result = self.prepare_buyer_layout(
                paths,
                str(globals_data.get("country", "")),
                str(globals_data.get("procurement_need", "")),
            )
            inspection["layout_ready"] = bool(layout_result.get("ready"))
            inspection["layout_warning"] = layout_result.get("warning", "")
        return inspection

    def create_research_job(self, slug: str, options: dict[str, Any]) -> dict[str, Any]:
        self.require(slug)
        if not self.model_settings.get("research_enabled", True):
            raise ValueError("当前处于手动/导入资料模式，未启用买家资料生成模型。请在模型设置中开启“买家资料生成”，或手动添加买家。")
        country = str(options.get("country", "") or "").strip()
        procurement_need = str(options.get("procurement_need", "") or "").strip()
        if not country or not procurement_need:
            raise ValueError("请同时填写国家和采购需求。")
        buyer_count = int(options.get("buyer_count", 10))
        if buyer_count < 1 or buyer_count > 30:
            raise ValueError("买家数量必须在1到30之间。")
        strategy = merge_research_strategy(procurement_need, dict(options.get("strategy") or {}))
        job_id = uuid.uuid4().hex[:12]
        job = {
            "id": job_id,
            "type": "buyer_research",
            "project": slug,
            "status": "queued",
            "created_at": utc_now(),
            "country": country,
            "procurement_need": procurement_need,
            "buyer_count": buyer_count,
            "fetch_assets": bool(options.get("fetch_assets", False)),
            "asset_mode": str(options.get("asset_mode", "light") or "light"),
            "enable_ai_visual_fallback": bool(options.get("enable_ai_visual_fallback", False)) and bool(self.model_settings.get("visual_enabled", False)),
            "strategy": strategy,
            "stdout": "",
            "stderr": "",
        }
        with self.jobs_lock:
            self.jobs[job_id] = job
        threading.Thread(target=self._run_research_job, args=(job_id,), daemon=True).start()
        return dict(job)
    def create_job(self, slug: str, options: dict[str, Any]) -> dict[str, Any]:
        paths = self.require(slug)
        for key in ("template", "records", "layout"):
            if not paths[key].is_file():
                raise ValueError(f"缺少必需文件：{paths[key].name}")
        project = read_json(paths["project"])
        config = read_json(paths["layout"])
        globals_data = record_globals(read_json(paths["records"]))
        if project.get("mode") == "buyer_board" and not is_buyer_layout(config):
            self.prepare_buyer_layout(
                paths,
                str(globals_data.get("country", "")),
                str(globals_data.get("procurement_need", "")),
            )
            config = read_json(paths["layout"])
        if project.get("mode") == "buyer_board" and is_buyer_layout(config):
            pipeline_mode = "buyer_board"
        elif project.get("mode") == "buyer_briefing" and is_buyer_briefing_layout(config):
            pipeline_mode = "buyer_briefing"
        else:
            pipeline_mode = "generic"
        if project.get("mode") == "buyer_board" and pipeline_mode != "buyer_board":
            raise ValueError(_u(r"\u5f53\u524d\u9879\u76ee\u662f\u4e70\u5bb6\u770b\u677f\uff0c\u4f46\u6a21\u677f\u7248\u5f0f\u5c1a\u672a\u8bc6\u522b\u6210\u529f\uff0c\u8bf7\u5148\u68c0\u67e5\u201c\u7248\u5f0f\u6620\u5c04\u201d\u3002"))
        if project.get("mode") == "buyer_briefing" and pipeline_mode != "buyer_briefing":
            raise ValueError(_u(r"\u5f53\u524d\u9879\u76ee\u662f\u4e70\u5bb6\u5546\u60c5\uff0c\u4f46\u7248\u5f0f\u6620\u5c04\u4e0d\u662f6\u4e70\u5bb6\u5546\u60c5\u7ed3\u6784\uff0c\u8bf7\u5148\u4f7f\u7528\u793a\u4f8b\u6620\u5c04\u6216\u68c0\u67e5\u201c\u7248\u5f0f\u6620\u5c04\u201d\u3002"))
        job_id = uuid.uuid4().hex[:12]
        requested_filename = safe_filename(str(options.get("filename", "finished.pptx")))
        filename = self.unique_output_filename(paths, requested_filename)
        job = {
            "id": job_id,
            "type": "export",
            "pipeline_mode": pipeline_mode,
            "project": slug,
            "status": "queued",
            "created_at": utc_now(),
            "filename": filename,
            "requested_filename": requested_filename,
            "output": str(paths["output"] / filename),
            "report": str(paths["workspace"] / f"run-{job_id}" / "console_export_report.json"),
            "strict": bool(options.get("strict", False)),
            "presentation_engine": presentation_engine(options.get("presentation_engine")),
            "stdout": "",
            "stderr": "",
        }
        with self.jobs_lock:
            self.jobs[job_id] = job
        threading.Thread(target=self._run_export_job, args=(job_id,), daemon=True).start()
        return dict(job)
    def get_job(self, job_id: str) -> dict[str, Any]:
        with self.jobs_lock:
            if job_id not in self.jobs:
                raise FileNotFoundError("任务不存在。")
            return dict(self.jobs[job_id])

    def set_job(self, job_id: str, **updates: Any) -> None:
        with self.jobs_lock:
            self.jobs[job_id].update(updates)

    def prune_jobs(self, keep_finished: int = 100) -> None:
        with self.jobs_lock:
            finished = sorted(
                (job for job in self.jobs.values() if job.get("status") in {"completed", "failed"}),
                key=lambda item: str(item.get("finished_at") or item.get("created_at") or ""),
                reverse=True,
            )
            for job in finished[keep_finished:]:
                self.jobs.pop(str(job["id"]), None)

    def record_run(self, paths: dict[str, Path], job: dict[str, Any]) -> None:
        project = read_json(paths["project"])
        runs = list(project.get("last_runs") or [])
        run = {
            "id": job["id"],
            "type": job.get("type"),
            "status": job["status"],
            "finished_at": job.get("finished_at"),
        }
        if job.get("filename"):
            run["filename"] = job["filename"]
        if job.get("country"):
            run["country"] = job["country"]
            run["procurement_need"] = job.get("procurement_need", "")
        runs.insert(0, run)
        project["last_runs"] = runs[:20]
        if job.get("type") == "export":
            project["export"] = {
                "filename": job.get("requested_filename", job["filename"]),
                "strict": job["strict"],
                "presentation_engine": job.get("presentation_engine", "auto"),
            }
        project["updated_at"] = utc_now()
        write_json(paths["project"], project)

    def _run_research_job(self, job_id: str) -> None:
        job = self.get_job(job_id)
        paths = self.require(str(job["project"]))
        run_dir = paths["workspace"] / f"research-{job_id}"
        source_json = run_dir / "buyers.generated.json"
        result_json = source_json
        self.set_job(job_id, status="running", started_at=utc_now(), stage="正在搜索并核实买家资料")
        command = [
            sys.executable,
            str(skill_script("discover_buyer_profiles.py")),
            "--country", str(job["country"]),
            "--procurement-need", str(job["procurement_need"]),
            "--output", str(source_json),
            "--workspace", str(run_dir / "research"),
            "--buyer-count", str(job["buyer_count"]),
            "--research-mode", clean_research_mode(self.model_settings.get("research_mode")),
            "--preferred-industries", str((job.get("strategy") or {}).get("preferred_industries", "")),
            "--excluded-company-types", str((job.get("strategy") or {}).get("excluded_company_types", "")),
            "--custom-requirements", str((job.get("strategy") or {}).get("custom_requirements", "")),
            "--candidate-multiplier", str((job.get("strategy") or {}).get("candidate_multiplier", 3)),
        ]
        if (job.get("strategy") or {}).get("prefer_import_evidence", True):
            command.append("--prefer-import-evidence")
        stdout_parts: list[str] = []
        stderr_parts: list[str] = []
        try:
            returncode, stdout, stderr = run_command(command, extra_env=self.model_env("research"))
            stdout_parts.append(stdout)
            stderr_parts.append(stderr)
            if returncode != 0:
                raise RuntimeError(stderr or stdout or "买家搜索失败。")

            if job.get("fetch_assets"):
                self.set_job(job_id, stage="正在获取企业Logo和官网产品图")
                result_json = run_dir / "buyers.with-assets.json"
                asset_command = [
                    sys.executable,
                    str(skill_script("fetch_buyer_assets.py")),
                    "--buyers", str(source_json),
                    "--output", str(result_json),
                    "--assets-dir", str(paths["assets"]),
                    "--cache-file", str(paths["workspace"] / "asset-cache.json"),
                    "--report-file", str(run_dir / "asset_fetch_report.json"),
                    "--asset-mode", str(job.get("asset_mode", "light")),
                ]
                if job.get("enable_ai_visual_fallback"):
                    asset_command.append("--enable-ai-visual-fallback")
                asset_code, asset_stdout, asset_stderr = run_command(
                    asset_command,
                    extra_env=self.model_env("visual"),
                )
                stdout_parts.append(asset_stdout)
                stderr_parts.append(asset_stderr)
                if asset_code != 0:
                    result_json = source_json
                    stderr_parts.append("素材抓取失败，已保留买家文字资料。")

            buyers = read_json(result_json)
            if not isinstance(buyers, list) or not buyers:
                raise ValueError(
                    "本次搜索没有生成可回填的买家。请放宽排除企业类型，补充优先行业/应用场景，"
                    "或切换到支持联网搜索的模型后重试。系统已保留现有买家表单，不会用空结果覆盖。"
                )
            saved = self.save_buyer_data(str(job["project"]), {
                "country": job["country"],
                "procurement_need": job["procurement_need"],
                "buyers": buyers,
                "enforce_research_copy_rules": True,
            })
            report_data = {
                "buyer_count": len(saved["records"]["records"]),
                "layout_ready": saved["layout_ready"],
                "warnings": saved["warnings"],
                "assets_requested": bool(job.get("fetch_assets")),
                "strategy": job.get("strategy") or {},
            }
            project = read_json(paths["project"])
            project["research_strategy"] = job.get("strategy") or {}
            write_json(paths["project"], project)
            self.set_job(
                job_id,
                status="completed",
                stage="已完成",
                returncode=0,
                stdout="\n".join(part for part in stdout_parts if part),
                stderr="\n".join(part for part in stderr_parts if part),
                report_data=report_data,
                finished_at=utc_now(),
            )
        except Exception as exc:
            message = str(exc)
            if "OPENAI_API_KEY" in message or "API key" in message or "API Key" in message:
                message = "未检测到可用的检索模型 API Key。请先打开“模型设置”配置供应商、API Key 和模型。"
            elif "PowerShell fallback also failed" in message:
                message = "Python urllib \u88ab\u7cfb\u7edf\u7b56\u7565\u62e6\u622a\uff0c\u5df2\u5c1d\u8bd5 PowerShell \u515c\u5e95\u8bf7\u6c42\uff0c\u4f46\u515c\u5e95\u4e5f\u5931\u8d25\u3002\u8be6\u7ec6\u539f\u56e0\uff1a" + message
            elif "WinError 10013" in message or "api.openai.com" in message:
                message = "\u8bbf\u95ee\u6a21\u578b\u63a5\u53e3\u88ab\u7cfb\u7edf\u6216\u7f51\u7edc\u7b56\u7565\u62e6\u622a\u3002\u82e5\u4f7f\u7528 DeepSeek\u3001Qwen \u7b49\u56fd\u5185\u6216\u517c\u5bb9\u670d\u52a1\u5546\uff0c\u8bf7\u5728\u201c\u6a21\u578b\u8bbe\u7f6e\u201d\u4e2d\u9009\u62e9\u5bf9\u5e94\u670d\u52a1\u5546\uff0c\u5e76\u4f7f\u7528 model_only \u6a21\u5f0f\uff1bDeepSeek \u9ed8\u8ba4\u6a21\u578b\u5efa\u8bae\u4f7f\u7528 deepseek-chat\u3002\u53ea\u6709\u663e\u5f0f\u9009\u62e9 OpenAI \u5185\u7f6e\u8054\u7f51\u641c\u7d22\u65f6\u624d\u4f1a\u8bbf\u95ee OpenAI Responses \u63a5\u53e3\u3002"
            self.set_job(
                job_id,
                status="failed",
                stage="搜索失败",
                returncode=-1,
                stdout="\n".join(part for part in stdout_parts if part),
                stderr=message,
                finished_at=utc_now(),
            )
        self.record_run(paths, self.get_job(job_id))
        self.prune_jobs()
    def _buyer_export_command(self, job: dict[str, Any], paths: dict[str, Path], run_dir: Path) -> list[str]:
        data = read_json(paths["records"])
        buyers = buyer_records(data)
        if not buyers:
            raise ValueError("还没有买家资料，请先手动录入或使用AI搜索生成。")
        if job.get("strict") and any(not item.get("name") for item in buyers):
            raise ValueError("严格模式下，每家买家都必须填写企业名称。")
        buyers_path = run_dir / "buyers.input.json"
        write_json(buyers_path, buyers)
        runner = buyer_pipeline_script()
        if runner is not None:
            return [
                sys.executable,
                str(runner),
                "--template", str(paths["template"]),
                "--buyers", str(buyers_path),
                "--layout-config", str(paths["layout"]),
                "--output", str(job["output"]),
                "--preview-dir", str(run_dir / "previews"),
                "--workspace", str(run_dir / "buyer-pipeline"),
            ]

        draft = run_dir / "text-draft.pptx"
        code, stdout, stderr = run_command([
            sys.executable,
            str(skill_script("fill_buyer_board_text.py")),
            str(paths["template"]),
            str(buyers_path),
            str(paths["layout"]),
            str(draft),
        ])
        if code != 0:
            raise RuntimeError(stderr or stdout)
        return [
            sys.executable,
            str(skill_script("apply_buyer_board_images_fallback.py")),
            "--input-ppt", str(draft),
            "--buyers-json", str(buyers_path),
            "--layout-config", str(paths["layout"]),
            "--output-ppt", str(job["output"]),
        ]

    def _briefing_export_command(self, job: dict[str, Any], paths: dict[str, Path], run_dir: Path) -> list[str]:
        data = read_json(paths["records"])
        pages = data.get("pages") if isinstance(data, dict) else data
        if not isinstance(pages, list) or not pages:
            raise ValueError(_u(r"\u8fd8\u6ca1\u6709\u4e70\u5bb6\u5546\u60c5\u9875\u9762\uff0c\u8bf7\u5148\u5728\u201c\u4e70\u5bb6\u5546\u60c5\u8868\u5355\u201d\u5f55\u5165\u81f3\u5c111\u9875\u3001\u6bcf\u98756\u5bb6\u4e70\u5bb6\u3002"))
        for page_index, page in enumerate(pages, start=1):
            if not isinstance(page, dict):
                raise ValueError(_u(r"\u7b2c") + str(page_index) + _u(r"\u9875\u4e70\u5bb6\u5546\u60c5\u6570\u636e\u683c\u5f0f\u4e0d\u6b63\u786e\u3002"))
            buyers = page.get("buyers")
            if not isinstance(buyers, list) or not 1 <= len(buyers) <= 6:
                raise ValueError(_u(r"\u7b2c") + str(page_index) + _u(r"\u9875\u5fc5\u987b\u586b\u51991-6\u5bb6\u4e70\u5bb6\u3002"))
            for buyer_index, buyer in enumerate(buyers, start=1):
                if not isinstance(buyer, dict) or not str(buyer.get("name", "")).strip():
                    raise ValueError(_u(r"\u7b2c") + str(page_index) + _u(r"\u9875\u7b2c") + str(buyer_index) + _u(r"\u5bb6\u4e70\u5bb6\u7f3a\u5c11\u4f01\u4e1a\u540d\u79f0\u3002"))
        return [
            sys.executable,
            str(skill_script("fill_buyer_briefing_pages.py")),
            str(paths["template"]),
            str(paths["records"]),
            str(job["output"]),
            "--layout-config", str(paths["layout"]),
            "--report", str(job["report"]),
        ]

    def _generic_export_command(self, job: dict[str, Any], paths: dict[str, Path], run_dir: Path) -> list[str]:
        config = read_json(paths["layout"])
        data = read_json(paths["records"])
        records = buyer_records(data)
        template = paths["template"]
        layout = paths["layout"]
        repeat = config.get("repeat") if isinstance(config.get("repeat"), dict) else {}
        if os.name == "nt" and job.get("presentation_engine") != "compatible" and repeat and records:
            source_index = int(repeat.get("source_slide_index") or repeat.get("start_slide_index") or 0)
            template_count = int(repeat.get("template_slide_count") or max(len(self.inspect_template(template).get("slides") or []) - source_index + 1, 1))
            copy_count = max(len(records) - template_count, 0)
            if source_index > 0 and copy_count:
                template = run_dir / "generic-template-expanded.pptx"
                code, stdout, stderr = run_command(
                    [
                        "powershell", "-NoProfile", "-NonInteractive", "-ExecutionPolicy", "Bypass",
                        "-File", str(skill_script("duplicate_ppt_slide.ps1")),
                        "-InputPpt", str(paths["template"]),
                        "-OutputPpt", str(template),
                        "-SourceSlideIndex", str(source_index),
                        "-CopyCount", str(copy_count),
                        "-Engine", str(job.get("presentation_engine", "auto")),
                    ],
                    timeout=EXPORT_SUBPROCESS_TIMEOUT_SECONDS,
                )
                if code != 0:
                    raise RuntimeError(stderr or stdout or "PowerPoint 无法复制通用模板页。")
                expanded_config = json.loads(json.dumps(config))
                expanded_config["repeat"]["template_slide_count"] = len(records)
                layout = run_dir / "generic-layout-expanded.json"
                write_json(layout, expanded_config)
        command = [
            sys.executable, str(generic_pipeline_script()),
            "--template", str(template),
            "--records", str(paths["records"]),
            "--layout-config", str(layout),
            "--output", str(job["output"]),
            "--workspace", str(run_dir),
            "--report", str(job["report"]),
        ]
        if job["strict"]:
            command.append("--strict")
        return command

    def _run_export_job(self, job_id: str) -> None:
        job = self.get_job(job_id)
        paths = self.require(str(job["project"]))
        run_dir = Path(job["report"]).parent
        self.set_job(job_id, status="running", started_at=utc_now(), stage="正在生成PPTX")
        try:
            if job["pipeline_mode"] == "buyer_board":
                command = self._buyer_export_command(job, paths, run_dir)
            elif job["pipeline_mode"] == "buyer_briefing":
                command = self._briefing_export_command(job, paths, run_dir)
            else:
                command = self._generic_export_command(job, paths, run_dir)
            returncode, stdout, stderr = run_command(command, timeout=EXPORT_SUBPROCESS_TIMEOUT_SECONDS)
            if returncode != 0:
                raise RuntimeError(stderr or stdout or "PPTX生成失败。")
            preview_source = run_dir / "previews"
            if job.get("presentation_engine") != "compatible" and not any(preview_source.glob("*.png")):
                preview_code, preview_stdout, preview_stderr = run_command(
                    [
                        "powershell",
                        "-NoProfile",
                        "-NonInteractive",
                        "-ExecutionPolicy", "Bypass",
                        "-File", str(skill_script("export_ppt_previews.ps1")),
                        "-InputPpt", str(job["output"]),
                        "-PreviewDir", str(preview_source),
                        "-Engine", str(job.get("presentation_engine", "auto")),
                    ],
                    timeout=EXPORT_SUBPROCESS_TIMEOUT_SECONDS,
                )
                if preview_code != 0:
                    stderr = "\n".join(part for part in (stderr, preview_stderr or preview_stdout) if part)
            preview_target = paths["output"] / ".previews" / Path(job["filename"]).stem
            preview_files = sorted(preview_source.glob("*.png"))
            if preview_files:
                if preview_target.exists():
                    shutil.rmtree(preview_target)
                shutil.copytree(preview_source, preview_target)
            output_metadata = inspect_output_ppt(Path(job["output"]))
            report_data = read_json(Path(job["report"])) if Path(job["report"]).is_file() else {
                "ok": True,
                "pipeline_mode": job["pipeline_mode"],
            }
            report_data.update(output_metadata)
            report_data["pipeline_mode"] = job["pipeline_mode"]
            report_data["record_count"] = record_count(read_json(paths["records"]))
            report_data["preview_count"] = len(preview_files)
            if job["pipeline_mode"] in {"buyer_board", "buyer_briefing", "generic"}:
                write_json(Path(job["report"]), report_data)
            self.set_job(
                job_id,
                status="completed",
                stage="已完成",
                returncode=returncode,
                stdout=stdout,
                stderr=stderr,
                report_data=report_data,
                finished_at=utc_now(),
            )
        except Exception as exc:
            self.set_job(
                job_id,
                status="failed",
                stage="生成失败",
                returncode=-1,
                stderr=str(exc),
                finished_at=utc_now(),
            )
        self.record_run(paths, self.get_job(job_id))
        self.prune_jobs()


class ConsoleHandler(BaseHTTPRequestHandler):
    server_version = "PPTTemplateConsole/0.2"

    @property
    def state(self) -> ConsoleState:
        return self.server.state  # type: ignore[attr-defined]

    def log_message(self, format: str, *args: Any) -> None:
        sys.stdout.write(f"[{self.log_date_time_string()}] {format % args}\n")

    def send_json(self, payload: Any, status: int = HTTPStatus.OK) -> None:
        data = json.dumps(payload, ensure_ascii=False).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(data)))
        self.send_header("Cache-Control", "no-store")
        self.end_headers()
        self.wfile.write(data)

    def send_error_json(self, exc: Exception) -> None:
        if isinstance(exc, FileNotFoundError):
            status = HTTPStatus.NOT_FOUND
        elif isinstance(exc, (ValueError, json.JSONDecodeError)):
            status = HTTPStatus.BAD_REQUEST
        else:
            status = HTTPStatus.INTERNAL_SERVER_ERROR
        self.send_json({"ok": False, "error": redact_credentials(str(exc)), "error_type": exc.__class__.__name__}, status)

    def read_body(self, max_bytes: int = 5 * 1024 * 1024) -> bytes:
        length = int(self.headers.get("Content-Length", "0"))
        if length < 0 or length > max_bytes:
            raise ValueError("请求内容超过大小限制。")
        return self.rfile.read(length)

    def read_json_body(self) -> Any:
        data = self.read_body()
        return json.loads(data.decode("utf-8-sig")) if data else {}

    def parts(self) -> list[str]:
        return [unquote(part) for part in urlparse(self.path).path.strip("/").split("/") if part]

    def do_GET(self) -> None:
        try:
            parts = self.parts()
            if parts == ["api", "health"]:
                self.send_json({
                    "ok": True,
                    "version": 5,
                    "features": [
                        "network-probe",
                        "powershell-fallback",
                        "qualified-buyer-research",
                        "document-import",
                        "natural-language-layout",
                        "layout-preflight",
                        "layout-recipes",
                        "json-recovery",
                    ],
                })
            elif parts == ["api", "model-settings"]:
                self.send_json(self.state.model_settings_payload())
            elif parts == ["api", "projects"]:
                self.send_json({"projects": self.state.list_projects()})
            elif len(parts) == 3 and parts[:2] == ["api", "projects"]:
                self.send_json(self.state.payload(parts[2]))
            elif len(parts) == 3 and parts[:2] == ["api", "jobs"]:
                self.send_json(self.state.get_job(parts[2]))
            elif len(parts) == 5 and parts[:2] == ["api", "projects"] and parts[3] == "output":
                self.send_output(parts[2], parts[4])
            elif len(parts) == 6 and parts[:2] == ["api", "projects"] and parts[3] == "preview":
                self.send_preview(parts[2], parts[4], parts[5])
            elif urlparse(self.path).path.startswith("/api/"):
                raise FileNotFoundError("接口不存在。")
            else:
                self.send_static(urlparse(self.path).path)
        except Exception as exc:
            self.send_error_json(exc)

    def do_POST(self) -> None:
        try:
            parts = self.parts()
            if parts == ["api", "projects"]:
                payload = self.read_json_body()
                self.send_json(self.state.create_project(str(payload.get("name", "")), str(payload.get("mode", "generic"))), HTTPStatus.CREATED)
            elif parts == ["api", "models"]:
                self.send_json(self.state.list_models(self.read_json_body()))
            elif parts == ["api", "network-probe"]:
                self.send_json(self.state.diagnose_model_connection(self.read_json_body()))
            elif len(parts) == 4 and parts[:2] == ["api", "projects"] and parts[3] == "template":
                self.send_json(self.state.save_template(parts[2], self.read_body(MAX_UPLOAD_BYTES)))
            elif len(parts) == 4 and parts[:2] == ["api", "projects"] and parts[3] == "import-records":
                query = urlparse(self.path).query
                params = dict(item.split("=", 1) if "=" in item else (item, "") for item in query.split("&") if item)
                filename = unquote(params.get("filename", "") or self.headers.get("X-Filename", "document.txt"))
                instruction = unquote(params.get("instruction", "") or self.headers.get("X-Instruction", ""))
                self.send_json(self.state.import_records_document(parts[2], self.read_body(MAX_UPLOAD_BYTES), filename, instruction))
            elif len(parts) == 4 and parts[:2] == ["api", "projects"] and parts[3] == "field-map":
                payload = self.read_json_body()
                mapping = payload.get("mapping") if isinstance(payload, dict) else None
                if not isinstance(mapping, dict):
                    raise ValueError("Field mapping must be an object.")
                self.send_json(self.state.rename_generic_fields(parts[2], mapping))
            elif len(parts) == 4 and parts[:2] == ["api", "projects"] and parts[3] == "layout-from-instruction":
                self.send_json(self.state.generate_layout_from_instruction(parts[2], self.read_json_body()))
            elif len(parts) == 4 and parts[:2] == ["api", "projects"] and parts[3] == "layout-preview":
                payload = self.read_json_body()
                self.send_json(self.state.mapping_preview(parts[2], payload.get("layout_config")))
            elif len(parts) == 4 and parts[:2] == ["api", "projects"] and parts[3] == "preflight":
                self.send_json(self.state.export_preflight(parts[2]))
            elif len(parts) == 4 and parts[:2] == ["api", "projects"] and parts[3] == "layout-recipes":
                self.send_json(self.state.save_layout_recipe(parts[2], self.read_json_body()), HTTPStatus.CREATED)
            elif len(parts) == 6 and parts[:2] == ["api", "projects"] and parts[3] == "layout-recipes" and parts[5] == "apply":
                self.send_json(self.state.apply_layout_recipe(parts[2], parts[4]))
            elif len(parts) == 4 and parts[:2] == ["api", "projects"] and parts[3] == "research-buyers":
                self.send_json(self.state.create_research_job(parts[2], self.read_json_body()), HTTPStatus.ACCEPTED)
            elif len(parts) == 4 and parts[:2] == ["api", "projects"] and parts[3] == "run":
                self.send_json(self.state.create_job(parts[2], self.read_json_body()), HTTPStatus.ACCEPTED)
            else:
                raise FileNotFoundError("接口不存在。")
        except Exception as exc:
            self.send_error_json(exc)

    def do_PUT(self) -> None:
        try:
            parts = self.parts()
            if parts == ["api", "model-settings"]:
                self.send_json(self.state.configure_models(self.read_json_body()))
            elif len(parts) == 4 and parts[:2] == ["api", "projects"] and parts[3] == "buyer-data":
                self.send_json(self.state.save_buyer_data(parts[2], self.read_json_body()))
            elif len(parts) == 5 and parts[:2] == ["api", "projects"] and parts[3] == "document":
                if parts[4] not in {"records", "layout"}:
                    raise ValueError("不支持的文档类型。")
                self.state.save_document(parts[2], parts[4], self.read_json_body())
                self.send_json({"ok": True})
            else:
                raise FileNotFoundError("接口不存在。")
        except Exception as exc:
            self.send_error_json(exc)

    def send_output(self, slug: str, filename: str) -> None:
        paths = self.state.require(slug)
        target = within(paths["output"], paths["output"] / Path(filename).name)
        if not target.is_file() or target.suffix.lower() != ".pptx":
            raise FileNotFoundError("输出文件不存在。")
        data = target.read_bytes()
        self.send_response(HTTPStatus.OK)
        self.send_header("Content-Type", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        ascii_name = re.sub(r"[^A-Za-z0-9._-]+", "-", target.name).strip("-") or "download.pptx"
        encoded_name = quote(target.name)
        disposition = f"attachment; filename=\"{ascii_name}\"; filename*=UTF-8''{encoded_name}"
        self.send_header("Content-Disposition", disposition)
        self.send_header("Content-Length", str(len(data)))
        self.end_headers()
        self.wfile.write(data)

    def send_preview(self, slug: str, filename: str, index_value: str) -> None:
        paths = self.state.require(slug)
        stem = Path(filename).stem
        try:
            index = int(index_value)
        except ValueError as exc:
            raise FileNotFoundError("预览页不存在。") from exc
        preview_dir = within(paths["output"], paths["output"] / ".previews" / stem)
        images = sorted(preview_dir.glob("*.png"), key=lambda item: natural_sort_key(item.name))
        if index < 1 or index > len(images):
            raise FileNotFoundError("预览页不存在。")
        target = within(preview_dir, images[index - 1])
        data = target.read_bytes()
        self.send_response(HTTPStatus.OK)
        self.send_header("Content-Type", "image/png")
        self.send_header("Content-Length", str(len(data)))
        self.send_header("Cache-Control", "private, max-age=300")
        self.end_headers()
        self.wfile.write(data)

    def send_static(self, path_value: str) -> None:
        relative = "index.html" if path_value in {"", "/"} else unquote(path_value).lstrip("/")
        target = within(static_root(), static_root() / relative)
        if not target.is_file():
            target = static_root() / "index.html"
        data = target.read_bytes()
        content_type = mimetypes.guess_type(target.name)[0] or "application/octet-stream"
        self.send_response(HTTPStatus.OK)
        self.send_header("Content-Type", f"{content_type}; charset=utf-8" if content_type.startswith("text/") else content_type)
        self.send_header("Content-Length", str(len(data)))
        self.send_header("Cache-Control", "no-cache")
        self.end_headers()
        self.wfile.write(data)


class ConsoleServer(ThreadingHTTPServer):
    daemon_threads = True

    def __init__(self, address: tuple[str, int], state: ConsoleState):
        super().__init__(address, ConsoleHandler)
        self.state = state


def build_server(
    host: str,
    port: int,
    state: ConsoleState,
    allow_non_loopback: bool = False,
) -> ConsoleServer:
    host = validate_bind_host(host, allow_non_loopback=allow_non_loopback)
    if port == 0:
        return ConsoleServer((host, 0), state)
    last_error: OSError | None = None
    for candidate in range(port, port + 20):
        try:
            return ConsoleServer((host, candidate), state)
        except OSError as exc:
            last_error = exc
    raise OSError(f"端口 {port}-{port + 19} 均不可用。") from last_error


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Run the local PPT template batch control console.")
    parser.add_argument("--host", default="127.0.0.1")
    parser.add_argument("--port", type=int, default=5310, help="Preferred port; searches the next 19 ports if busy")
    parser.add_argument("--projects-root", default="console-projects")
    parser.add_argument("--no-open", action="store_true")
    parser.add_argument(
        "--allow-non-loopback",
        action="store_true",
        help="Explicitly allow binding beyond the local machine.",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    if not static_root().is_dir():
        raise FileNotFoundError(f"Control console assets are missing: {static_root()}")
    state = ConsoleState(Path(args.projects_root))
    server = build_server(args.host, args.port, state, allow_non_loopback=args.allow_non_loopback)
    url = f"http://{args.host}:{server.server_address[1]}/"
    print(f"PPT Template Batch Console: {url}")
    print(f"Projects: {state.projects_root}")
    if not args.no_open:
        threading.Timer(0.5, lambda: webbrowser.open(url)).start()
    try:
        server.serve_forever(poll_interval=0.25)
    except KeyboardInterrupt:
        pass
    finally:
        server.server_close()
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
