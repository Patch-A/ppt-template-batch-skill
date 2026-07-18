from __future__ import annotations

import argparse
import copy
import json
import os
import re
import sys
import tempfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Inches

try:
    from PIL import Image
except Exception:  # pragma: no cover - dependency is optional at import time
    Image = None


TOKEN_RE = re.compile(r"\{([^{}]+)\}")
UNRESOLVED_TOKEN_RE = re.compile(r"\{\{[^{}]+\}\}|\{(?:record|globals|data)\.[^{}]+\}")
STALE_TEXT_HINTS = (
    "placeholder",
    "sample text",
    "replace me",
    "stale template",
    "lorem ipsum",
    "todo",
    "tbd",
    "请替换",
    "待替换",
    "示例文本",
)
LAYOUT_SCHEMA_VERSION = 2


def load_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8-sig"))


def write_json(path: Path, payload: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8-sig")


def normalize_config(config: dict[str, Any]) -> dict[str, Any]:
    """Normalize legacy configs without mutating the user's JSON input."""
    if not isinstance(config, dict):
        raise ValueError("layout config must be a JSON object.")
    normalized = copy.deepcopy(config)
    version = normalized.get("schema_version", normalized.get("version", 1))
    try:
        version = int(version)
    except (TypeError, ValueError) as exc:
        raise ValueError("layout config schema_version must be an integer.") from exc
    if version > LAYOUT_SCHEMA_VERSION:
        raise ValueError(f"Unsupported layout config schema_version: {version}.")
    normalized["schema_version"] = LAYOUT_SCHEMA_VERSION
    return normalized


def resolve_path(path: str | Path, base_dir: Path) -> Path:
    candidate = Path(path)
    if candidate.is_absolute():
        return candidate
    return (base_dir / candidate).resolve()


def get_records(data: Any, config: dict[str, Any]) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    if isinstance(data, list):
        return [dict(item) for item in data], {}
    if not isinstance(data, dict):
        raise ValueError("records JSON must be either a list or an object.")

    record_key = str(config.get("record_key") or config.get("records_key") or "records")
    records = data.get(record_key)
    if records is None:
        records = data.get("items")
    if records is None:
        records = [data]
    if not isinstance(records, list):
        raise ValueError(f"records field '{record_key}' must be a list.")

    globals_data = {}
    for key in ("globals", "global", "meta", "defaults"):
        if isinstance(data.get(key), dict):
            globals_data.update(data[key])
    return [dict(item) for item in records], globals_data


def resolve_field(field: str, context: dict[str, Any], report: dict[str, Any]) -> Any:
    roots = {
        "record": context.get("record", {}),
        "globals": context.get("globals", {}),
        "data": context.get("data", {}),
    }
    if "." in field:
        root_name, remainder = field.split(".", 1)
        if root_name in roots:
            return resolve_dotted(roots[root_name], remainder, field, report)

    for root_name in ("record", "globals", "data"):
        value = resolve_dotted(roots[root_name], field, field, report, mark_missing=False)
        if value is not None:
            return value

    report.setdefault("missing_fields", []).append(field)
    return ""


def resolve_dotted(source: Any, field: str, display: str, report: dict[str, Any], mark_missing: bool = True) -> Any:
    current = source
    for part in field.split("."):
        if isinstance(current, dict) and part in current:
            current = current[part]
            continue
        if mark_missing:
            report.setdefault("missing_fields", []).append(display)
        return None
    return current


def render_value(spec: dict[str, Any], context: dict[str, Any], report: dict[str, Any]) -> str:
    if "value" in spec:
        return str(spec.get("value") or "")
    if "field" in spec:
        value = resolve_field(str(spec["field"]), context, report)
        if isinstance(value, list):
            return str(spec.get("separator", "、")).join(str(item) for item in value if item is not None)
        return str(value or "")
    if "template" in spec:
        template = str(spec["template"])

        def replace(match: re.Match[str]) -> str:
            return str(resolve_field(match.group(1).strip(), context, report) or "")

        return TOKEN_RE.sub(replace, template)
    return ""


def _shape_role(shape) -> str:
    explicit_role = getattr(shape, "role", None)
    if explicit_role:
        return str(explicit_role).strip().lower()
    element_role = shape.element.get("role") or shape.element.get("data-role")
    if element_role:
        return str(element_role).strip().lower()
    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "is_placeholder", False):
        return "placeholder"
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if getattr(shape, "has_text_frame", False):
        return "text"
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        return "group"
    return "shape"


def _matches_stable_selector(shape, stable_selector: dict[str, Any]) -> bool:
    name = stable_selector.get("name")
    if name is not None and shape.name != str(name):
        return False
    role = stable_selector.get("role")
    if role is None:
        return True
    expected = str(role).strip().lower()
    actual = _shape_role(shape)
    aliases = {
        "textbox": "text",
        "text_box": "text",
        "picture": "image",
        "pic": "image",
    }
    normalized_expected = aliases.get(expected, expected)
    accepted_roles = {
        "text": {"text", "placeholder"},
        "image": {"image"},
        "table": {"table"},
        "group": {"group"},
        "placeholder": {"placeholder"},
        "shape": {"shape"},
    }
    return actual in accepted_roles.get(normalized_expected, {normalized_expected})


def get_shape(slide, selector: dict[str, Any]):
    stable_selector = selector.get("selector")
    if not isinstance(stable_selector, dict):
        direct_stable_keys = {key: selector[key] for key in ("name", "role") if key in selector}
        stable_selector = direct_stable_keys or None

    if stable_selector:
        for shape in slide.shapes:
            if _matches_stable_selector(shape, stable_selector):
                return shape

    if "shape_id" in selector:
        shape_id = int(selector["shape_id"])
        for shape in slide.shapes:
            if shape.shape_id == shape_id:
                return shape
        raise KeyError(f"shape_id {shape_id} was not found.")

    if "shape_name" in selector:
        shape_name = str(selector["shape_name"])
        for shape in slide.shapes:
            if shape.name == shape_name:
                return shape
        raise KeyError(f"shape_name {shape_name!r} was not found.")

    if "shape_index" in selector:
        index = int(selector["shape_index"]) - 1
        if index < 0 or index >= len(slide.shapes):
            raise IndexError(f"shape_index {selector['shape_index']} is outside slide shape range.")
        return slide.shapes[index]

    raise ValueError("A mapping must include selector, shape_index, shape_id, or shape_name.")


def copy_font(source_run, target_run) -> None:
    if source_run is None:
        return
    target_run.font.name = source_run.font.name
    target_run.font.size = source_run.font.size
    target_run.font.bold = source_run.font.bold
    target_run.font.italic = source_run.font.italic
    target_run.font.underline = source_run.font.underline
    try:
        target_run.font.color.rgb = source_run.font.color.rgb
    except Exception:
        try:
            target_run.font.color.theme_color = source_run.font.color.theme_color
        except Exception:
            pass


def apply_font_overrides(text_frame, spec: dict[str, Any]) -> None:
    font = spec.get("font") or {}
    if not font:
        return
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if "name" in font:
                run.font.name = font["name"]
            if "size_pt" in font:
                run.font.size = Inches(float(font["size_pt"]) / 72)
            if "bold" in font:
                run.font.bold = bool(font["bold"])
            if "italic" in font:
                run.font.italic = bool(font["italic"])
            if "color" in font:
                run.font.color.rgb = RGBColor.from_string(str(font["color"]).lstrip("#"))


def replace_text_frame(text_frame, value: str, mode: str = "clear") -> None:
    source_paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else None
    source_run = source_paragraph.runs[0] if source_paragraph and source_paragraph.runs else None
    alignment = source_paragraph.alignment if source_paragraph else None
    level = source_paragraph.level if source_paragraph else 0

    if mode == "in_place":
        paragraph = text_frame.paragraphs[0]
        while len(paragraph.runs) < 1:
            paragraph.add_run()
        paragraph.runs[0].text = value
        for run in paragraph.runs[1:]:
            run.text = ""
        for extra_paragraph in text_frame.paragraphs[1:]:
            for run in extra_paragraph.runs:
                run.text = ""
        return

    text_frame.clear()
    lines = value.splitlines() or [""]
    for line_index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if line_index == 0 else text_frame.add_paragraph()
        paragraph.alignment = alignment
        paragraph.level = level
        run = paragraph.add_run()
        run.text = line
        copy_font(source_run, run)


def set_shape_text(slide, spec: dict[str, Any], context: dict[str, Any], report: dict[str, Any]) -> None:
    shape = get_shape(slide, spec)
    if not getattr(shape, "has_text_frame", False):
        raise ValueError(f"Target shape {shape.name!r} does not have a text frame.")
    value = render_value(spec, context, report)
    replace_text_frame(shape.text_frame, value, str(spec.get("mode", "clear")))
    apply_font_overrides(shape.text_frame, spec)


def set_table_cell(slide, table_spec: dict[str, Any], cell_spec: dict[str, Any], context: dict[str, Any], report: dict[str, Any]) -> None:
    shape = get_shape(slide, table_spec)
    if not getattr(shape, "has_table", False):
        raise ValueError(f"Target shape {shape.name!r} does not have a table.")
    row = int(cell_spec["row"])
    col = int(cell_spec["col"])
    cell = shape.table.cell(row, col)
    value = render_value(cell_spec, context, report)
    replace_text_frame(cell.text_frame, value, str(cell_spec.get("mode", table_spec.get("mode", "clear"))))
    apply_font_overrides(cell.text_frame, cell_spec)


def iter_text_frames(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            yield shape.text_frame
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame


def iter_text_targets(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            yield shape, shape.text_frame, None
        if getattr(shape, "has_table", False):
            for row_index, row in enumerate(shape.table.rows):
                for col_index, cell in enumerate(row.cells):
                    yield shape, cell.text_frame, {
                        "row": row_index,
                        "col": col_index,
                        "width": int(shape.table.columns[col_index].width),
                        "height": int(shape.table.rows[row_index].height),
                    }


def _text_capacity_warning(shape, text_frame, slide_index: int, cell: dict[str, int] | None):
    text = text_frame.text or ""
    if not text.strip() or not getattr(shape, "width", 0) or not getattr(shape, "height", 0):
        return None

    font_size_pt = 12.0
    for paragraph in text_frame.paragraphs:
        if paragraph.runs and paragraph.runs[0].font.size:
            font_size_pt = max(float(paragraph.runs[0].font.size.pt), 1.0)
            break
    target_width = cell["width"] if cell is not None else shape.width
    target_height = cell["height"] if cell is not None else shape.height
    width_pt = max(float(target_width) / 12700, 1.0)
    height_pt = max(float(target_height) / 12700, 1.0)
    chars_per_line = max(int(width_pt / (font_size_pt * 0.6)), 1)
    line_capacity = max(int(height_pt / (font_size_pt * 1.2)), 1)
    estimated_lines = sum(
        max((len(line) + chars_per_line - 1) // chars_per_line, 1)
        for line in (text.splitlines() or [""])
    )
    if estimated_lines <= line_capacity:
        return None

    item = {
        "slide_index": slide_index,
        "shape_name": shape.name,
        "estimated_lines": estimated_lines,
        "capacity_lines": line_capacity,
        "text_length": len(text),
    }
    if cell is not None:
        item["row"] = cell["row"]
        item["col"] = cell["col"]
    return item


def scan_shared_preflight(presentation: Presentation, report: dict[str, Any]) -> None:
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape, text_frame, cell in iter_text_targets(slide):
            text = text_frame.text or ""
            if not text.strip():
                continue
            lowered = text.casefold()
            unresolved = UNRESOLVED_TOKEN_RE.search(text)
            hint = next((marker for marker in STALE_TEXT_HINTS if marker in lowered), None)
            if unresolved or hint:
                item = {
                    "slide_index": slide_index,
                    "shape_name": shape.name,
                    "text": text,
                    "reason": "unresolved_placeholder" if unresolved else "template_marker",
                }
                if cell is not None:
                    item["row"] = cell["row"]
                    item["col"] = cell["col"]
                report.setdefault("stale_template_text", []).append(item)
                report.setdefault("warnings", []).append(
                    f"Stale template text remains on slide {slide_index} in {shape.name!r}."
                )

            capacity_warning = _text_capacity_warning(shape, text_frame, slide_index, cell)
            if capacity_warning:
                report.setdefault("capacity_warnings", []).append(capacity_warning)
                report.setdefault("warnings", []).append(
                    f"Text may exceed capacity on slide {slide_index} in {shape.name!r}."
                )


def replace_placeholders(slide, placeholders: dict[str, Any], context: dict[str, Any], report: dict[str, Any]) -> None:
    rendered = {}
    for token, spec in placeholders.items():
        if isinstance(spec, str):
            rendered[token] = str(resolve_field(spec, context, report) or "")
        elif isinstance(spec, dict):
            rendered[token] = render_value(spec, context, report)
        else:
            rendered[token] = str(spec)

    for text_frame in iter_text_frames(slide):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text
                for token, value in rendered.items():
                    text = text.replace(token, value)
                run.text = text


def remove_shape(shape) -> None:
    shape.element.getparent().remove(shape.element)


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        remove_shape(shape)


def length_to_emu(value: Any, unit: str) -> int:
    if value is None:
        raise ValueError("Image bounds require left, top, width, and height.")
    if unit == "emu":
        return int(value)
    if unit == "in":
        return int(Inches(float(value)))
    if unit == "pt":
        return int(Emu(float(value) * 12700))
    if unit == "px":
        return int(Emu(float(value) * 9525))
    raise ValueError(f"Unsupported length unit: {unit}")


def get_image_bounds(slide, spec: dict[str, Any]) -> tuple[int, int, int, int, Any | None]:
    if any(key in spec for key in ("selector", "shape_index", "shape_id", "shape_name")):
        shape = get_shape(slide, spec)
        return int(shape.left), int(shape.top), int(shape.width), int(shape.height), shape

    unit = str(spec.get("unit", "in"))
    left = length_to_emu(spec.get("left"), unit)
    top = length_to_emu(spec.get("top"), unit)
    width = length_to_emu(spec.get("width"), unit)
    height = length_to_emu(spec.get("height"), unit)
    return left, top, width, height, None


def normalize_image_path(image_path: Path, workspace: Path, report: dict[str, Any]) -> Path:
    if image_path.suffix.lower() != ".svg":
        return image_path
    try:
        import cairosvg
    except Exception:
        report.setdefault("warnings", []).append(f"SVG image could not be converted because cairosvg is unavailable: {image_path}")
        return image_path

    workspace.mkdir(parents=True, exist_ok=True)
    output = workspace / f"{image_path.stem}.converted.png"
    cairosvg.svg2png(url=str(image_path), write_to=str(output))
    return output


def crop_cover_image(image_path: Path, target_width: int, target_height: int, workspace: Path, report: dict[str, Any]) -> Path:
    if Image is None:
        report.setdefault("warnings", []).append("Pillow is unavailable; cover images will be inserted without pre-crop.")
        return image_path

    with Image.open(image_path) as image:
        image = image.convert("RGB")
        src_w, src_h = image.size
        target_ratio = target_width / target_height
        src_ratio = src_w / src_h
        if src_ratio > target_ratio:
            new_w = int(src_h * target_ratio)
            left = max((src_w - new_w) // 2, 0)
            box = (left, 0, left + new_w, src_h)
        else:
            new_h = int(src_w / target_ratio)
            top = max((src_h - new_h) // 2, 0)
            box = (0, top, src_w, top + new_h)
        cropped = image.crop(box)
        workspace.mkdir(parents=True, exist_ok=True)
        output = workspace / f"{image_path.stem}.cover.jpg"
        cropped.save(output, quality=92)
        return output


def add_image(slide, image_path: Path, left: int, top: int, width: int, height: int, fit: str, workspace: Path, report: dict[str, Any]) -> None:
    prepared = normalize_image_path(image_path, workspace, report)
    if fit == "cover":
        prepared = crop_cover_image(prepared, width, height, workspace, report)
        slide.shapes.add_picture(str(prepared), left, top, width=width, height=height)
        return

    if Image is None:
        slide.shapes.add_picture(str(prepared), left, top, width=width)
        return

    with Image.open(prepared) as image:
        src_w, src_h = image.size
    scale = min(width / src_w, height / src_h)
    fitted_w = int(src_w * scale)
    fitted_h = int(src_h * scale)
    fitted_left = left + int((width - fitted_w) / 2)
    fitted_top = top + int((height - fitted_h) / 2)
    slide.shapes.add_picture(str(prepared), fitted_left, fitted_top, width=fitted_w, height=fitted_h)


def set_image(slide, spec: dict[str, Any], context: dict[str, Any], paths_base: Path, workspace: Path, report: dict[str, Any]) -> None:
    value = render_value(spec, context, report)
    left, top, width, height, placeholder_shape = get_image_bounds(slide, spec)
    if not value:
        if spec.get("clear_if_missing", True) and placeholder_shape is not None:
            remove_shape(placeholder_shape)
        return

    image_path = resolve_path(value, paths_base)
    if not image_path.exists():
        report.setdefault("missing_assets", []).append(str(image_path))
        if spec.get("clear_if_missing", True) and placeholder_shape is not None:
            remove_shape(placeholder_shape)
        return

    if spec.get("replace_existing", True) and placeholder_shape is not None:
        remove_shape(placeholder_shape)
    add_image(slide, image_path, left, top, width, height, str(spec.get("fit", "contain")), workspace, report)


def duplicate_slide(presentation: Presentation, source_index: int):
    source_slide = presentation.slides[source_index]
    blank_layout = presentation.slide_layouts[6]
    new_slide = presentation.slides.add_slide(blank_layout)

    # Copied XML keeps the source rIds. Create equivalent relationships on the
    # new slide and rewrite those rIds so PowerPoint can open the result.
    relationship_ids: dict[str, str] = {}
    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
            continue
        try:
            relationship_ids[rel.rId] = new_slide.part.relate_to(
                rel._target,
                rel.reltype,
                is_external=bool(rel.is_external),
            )
        except Exception:
            continue

    for shape in source_slide.shapes:
        copied_shape = copy.deepcopy(shape.element)
        for element in copied_shape.iter():
            for attribute in (qn("r:embed"), qn("r:link"), qn("r:id")):
                source_rid = element.get(attribute)
                if source_rid in relationship_ids:
                    element.set(attribute, relationship_ids[source_rid])
        new_slide.shapes._spTree.insert_element_before(copied_shape, "p:extLst")
    return new_slide


def move_slide(presentation: Presentation, old_index: int, new_index: int) -> None:
    slide_ids = presentation.slides._sldIdLst
    slide_id = list(slide_ids)[old_index]
    slide_ids.remove(slide_id)
    slide_ids.insert(new_index, slide_id)


def delete_slide(presentation: Presentation, slide_index: int) -> None:
    slide_ids = presentation.slides._sldIdLst
    slide_id = list(slide_ids)[slide_index]
    rel_id = slide_id.rId
    slide_ids.remove(slide_id)
    presentation.part.drop_rel(rel_id)


def prepare_repeated_slides(presentation: Presentation, config: dict[str, Any], records: list[dict[str, Any]], report: dict[str, Any]) -> None:
    repeat = config.get("repeat")
    if not repeat:
        return

    start_index = int(repeat.get("start_slide_index") or repeat.get("source_slide_index")) - 1
    source_index = int(repeat.get("source_slide_index") or repeat.get("start_slide_index")) - 1
    record_count = min(len(records), int(repeat.get("max_records", len(records))))
    template_count = int(repeat.get("template_slide_count") or max(len(presentation.slides) - start_index, 1))

    if record_count > template_count:
        for offset in range(record_count - template_count):
            duplicate_slide(presentation, source_index)
            move_slide(presentation, len(presentation.slides) - 1, start_index + template_count + offset)
        report.setdefault("warnings", []).append(
            f"Duplicated {record_count - template_count} repeated slide(s) from slide {source_index + 1}."
        )

    if repeat.get("trim_extra_template_slides", True) and template_count > record_count:
        for _ in range(template_count - record_count):
            delete_slide(presentation, start_index + record_count)
        report.setdefault("warnings", []).append(f"Trimmed {template_count - record_count} extra repeated slide(s).")


def validate_required_fields(
    records: list[dict[str, Any]], required_fields: list[str], report: dict[str, Any]
) -> set[int]:
    failed_indexes: set[int] = set()
    for record_index, record in enumerate(records, start=1):
        missing_for_record: list[str] = []
        for field in required_fields:
            value = resolve_dotted(record, field, field, report, mark_missing=False)
            if value in (None, ""):
                report.setdefault("missing_required_fields", []).append(
                    {"record_index": record_index, "field": field}
                )
                missing_for_record.append(field)
        if missing_for_record:
            failed_indexes.add(record_index)
            report.setdefault("failed_records", []).append(
                {"record_index": record_index, "missing_required_fields": missing_for_record}
            )
    return failed_indexes


def apply_slide_mapping(
    slide,
    mapping: dict[str, Any],
    context: dict[str, Any],
    paths_base: Path,
    workspace: Path,
    report: dict[str, Any],
) -> None:
    placeholders = mapping.get("placeholders") or {}
    if placeholders:
        replace_placeholders(slide, placeholders, context, report)

    for spec in mapping.get("texts", []):
        set_shape_text(slide, spec, context, report)

    for table_spec in mapping.get("tables", []):
        for cell_spec in table_spec.get("cells", []):
            set_table_cell(slide, table_spec, cell_spec, context, report)

    for spec in mapping.get("images", []):
        set_image(slide, spec, context, paths_base, workspace, report)

    for spec in mapping.get("clear_shapes", []):
        remove_shape(get_shape(slide, spec))


def save_presentation_atomically(
    presentation: Presentation, output_path: Path, report: dict[str, Any]
) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    temp_path: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(
            mode="wb",
            prefix=f".{output_path.stem}.",
            suffix=".tmp.pptx",
            dir=output_path.parent,
            delete=False,
        ) as handle:
            temp_path = Path(handle.name)

        presentation.save(temp_path)
        reopened = Presentation(temp_path)
        report["slide_count"] = len(reopened.slides)
        report["reopen_ok"] = True
        report["reopen_status"] = {
            "status": "ok",
            "ok": True,
            "slide_count": len(reopened.slides),
            "replaced": False,
        }
        os.replace(temp_path, output_path)
        report["reopen_status"]["replaced"] = True
        temp_path = None
    except Exception as exc:
        report["reopen_ok"] = False
        report["reopen_status"] = {
            "status": "failed",
            "ok": False,
            "error_type": exc.__class__.__name__,
            "error": str(exc),
        }
    finally:
        if temp_path is not None:
            try:
                temp_path.unlink()
            except FileNotFoundError:
                pass


def fill_presentation(
    template: Path,
    records_path: Path,
    config_path: Path,
    output_path: Path,
    workspace: Path,
    strict: bool = False,
) -> dict[str, Any]:
    config = normalize_config(load_json(config_path))
    data = load_json(records_path)
    records, globals_data = get_records(data, config)
    report: dict[str, Any] = {
        "schema_version": LAYOUT_SCHEMA_VERSION,
        "template": str(template),
        "records": str(records_path),
        "layout_config": str(config_path),
        "output": str(output_path),
        "record_count": len(records),
        "processed_record_count": 0,
        "slides_processed": 0,
        "warnings": [],
        "missing_fields": [],
        "missing_assets": [],
        "missing_required_fields": [],
        "failed_records": [],
        "stale_template_text": [],
        "capacity_warnings": [],
        "capacity": {"ok": True, "warnings": []},
        "expected_slide_count": 0,
        "slide_count": 0,
        "slide_count_status": {"expected": 0, "actual": 0, "ok": False},
        "reopen_ok": False,
        "reopen_status": {"status": "not_run", "ok": False},
        "reopen": {"status": "not_run", "ok": False},
    }

    repeat = config.get("repeat") or {}
    max_records = int(repeat["max_records"]) if "max_records" in repeat else len(records)
    validation_records = records[:max_records] if config.get("repeat") else records
    required_fields = config.get("required_fields") or []
    failed_indexes = validate_required_fields(validation_records, required_fields, report)

    presentation = Presentation(template)
    report["expected_slide_count"] = len(presentation.slides)

    if strict and failed_indexes:
        report["warnings"].append("Strict mode rejected records with missing required fields.")
        report["slide_count_status"] = {
            "expected": report["expected_slide_count"],
            "actual": 0,
            "ok": False,
        }
        report["ok"] = False
        return report

    effective_records = [
        (record_index, record)
        for record_index, record in enumerate(validation_records, start=1)
        if strict or record_index not in failed_indexes
    ]
    report["processed_record_count"] = len(effective_records)
    effective_record_values = [record for _, record in effective_records]

    prepare_repeated_slides(presentation, config, effective_record_values, report)
    report["expected_slide_count"] = len(presentation.slides)

    data_context = data if isinstance(data, dict) else {"records": records}
    paths_base = records_path.parent

    for slide_mapping in config.get("slides", []):
        slide_index = int(slide_mapping["slide_index"]) - 1
        if slide_index < 0 or slide_index >= len(presentation.slides):
            report["warnings"].append(f"Skipped slide mapping outside deck range: {slide_index + 1}")
            continue
        record = {}
        original_record_index = None
        if "record_index" in slide_mapping:
            original_record_index = int(slide_mapping["record_index"])
            if not strict and original_record_index in failed_indexes:
                clear_slide(presentation.slides[slide_index])
                report["warnings"].append(
                    f"Cleared slide {slide_index + 1} for failed record {original_record_index}."
                )
                continue
            record = records[original_record_index - 1]
        context = {"record": record, "globals": globals_data, "data": data_context}
        apply_slide_mapping(presentation.slides[slide_index], slide_mapping, context, paths_base, workspace, report)
        report["slides_processed"] += 1

    if repeat:
        start_index = int(repeat.get("start_slide_index") or repeat.get("source_slide_index")) - 1
        record_count = len(effective_records)
        mappings = {
            "placeholders": repeat.get("placeholders") or {},
            "texts": repeat.get("texts") or repeat.get("mappings", {}).get("texts", []),
            "tables": repeat.get("tables") or repeat.get("mappings", {}).get("tables", []),
            "images": repeat.get("images") or repeat.get("mappings", {}).get("images", []),
            "clear_shapes": repeat.get("clear_shapes") or [],
        }
        for offset, (_, record) in enumerate(effective_records[:record_count]):
            context = {"record": record, "globals": globals_data, "data": data_context}
            apply_slide_mapping(
                presentation.slides[start_index + offset],
                mappings,
                context,
                paths_base,
                workspace,
                report,
            )
            report["slides_processed"] += 1

    scan_shared_preflight(presentation, report)
    report["capacity"] = {
        "ok": not report["capacity_warnings"],
        "warnings": list(report["capacity_warnings"]),
    }
    report["slide_count_status"] = {
        "expected": report["expected_slide_count"],
        "actual": len(presentation.slides),
        "ok": report["expected_slide_count"] == len(presentation.slides),
    }
    report["ok"] = not (
        report["missing_required_fields"]
        or report["missing_assets"]
        or report["stale_template_text"]
    )
    save_presentation_atomically(presentation, output_path, report)
    report["reopen"] = dict(report["reopen_status"])
    report["ok"] = bool(report["ok"] and report["reopen_ok"])
    return report


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Fill a PPTX template from generic records.json and layout-config.json."
    )
    parser.add_argument("--template", required=True, help="Template PPTX path")
    parser.add_argument("--records", required=True, help="records.json path")
    parser.add_argument("--layout-config", "--config", dest="layout_config", required=True, help="layout-config.json path")
    parser.add_argument("--output", required=True, help="Output PPTX path")
    parser.add_argument("--workspace", help="Workspace for converted/cropped images")
    parser.add_argument("--report", help="Optional fill report JSON path")
    parser.add_argument("--strict", action="store_true", help="Exit non-zero when required fields are missing.")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    workspace = Path(args.workspace) if args.workspace else Path(args.output).parent / "_ppt_batch_workspace"
    try:
        report = fill_presentation(
            Path(args.template),
            Path(args.records),
            Path(args.layout_config),
            Path(args.output),
            workspace,
            strict=args.strict,
        )
    except Exception as exc:
        failure = {
            "ok": False,
            "error_type": exc.__class__.__name__,
            "error": str(exc),
        }
        if args.report:
            write_json(Path(args.report), failure)
        print(json.dumps(failure, ensure_ascii=False, indent=2), file=sys.stderr)
        return 2

    if args.report:
        write_json(Path(args.report), report)
    print(args.output)
    if args.strict and not report.get("ok", False):
        return 3
    if not report.get("reopen_ok", False):
        return 2
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
