from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_PIXEL = 12700


def emu_to_px(value: int) -> float:
    return round(value / EMU_PER_PIXEL, 2)


def extract_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return shape.text.strip()


def build_cover_config(slide) -> dict[str, Any]:
    text_candidates = []
    for index, shape in enumerate(slide.shapes, start=1):
        text = extract_text(shape)
        if not text:
            continue
        text_candidates.append(
            {
                "shape_index": index,
                "text": text,
                "top": shape.top,
                "height": shape.height,
            }
        )

    if not text_candidates:
        raise ValueError("No text boxes found on cover slide.")

    title_shape = max(text_candidates, key=lambda item: item["height"])
    country_shape = next((item for item in text_candidates if "国家" in item["text"]), None)
    if country_shape is None:
        remaining = [item for item in text_candidates if item["shape_index"] != title_shape["shape_index"]]
        country_shape = min(remaining, key=lambda item: item["top"]) if remaining else title_shape

    return {
        "slide_index": 1,
        "title_shape_index": title_shape["shape_index"],
        "country_shape_index": country_shape["shape_index"],
    }


def infer_value_key(label: str) -> str:
    mapping = {
        "企业": "name",
        "公司": "name",
        "国家": "country",
        "网站": "website",
        "官网": "website",
        "采购产品": "products",
        "采购需求": "products",
        "采购品类": "products",
        "简介": "bio",
        "企业简介": "bio",
        "买家简介": "bio",
    }
    return mapping.get(label, f"field_{label}")


def choose_title_shape(slide, table_shape_index: int) -> int:
    candidates = []
    for index, shape in enumerate(slide.shapes, start=1):
        if index == table_shape_index:
            continue
        text = extract_text(shape)
        if not text:
            continue
        if shape.top > 100 * EMU_PER_PIXEL:
            continue
        candidates.append((index, shape.height, shape.top))

    if not candidates:
        raise ValueError("No content title candidate found.")
    candidates.sort(key=lambda item: (-item[1], item[2]))
    return candidates[0][0]


def choose_footer_shape(slide, title_shape_index: int) -> int | None:
    candidates = []
    for index, shape in enumerate(slide.shapes, start=1):
        if index == title_shape_index:
            continue
        text = extract_text(shape)
        if not text:
            continue
        if shape.top < 400 * EMU_PER_PIXEL:
            continue
        candidates.append((index, shape.top))
    if not candidates:
        return None
    candidates.sort(key=lambda item: item[1])
    return candidates[0][0]


def detect_image_slots(slide, table_shape=None, slide_width: int | None = None, slide_height: int | None = None) -> dict[str, Any]:
    pictures = []
    for index, shape in enumerate(slide.shapes, start=1):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        pictures.append(
            {
                "shape_index": index,
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "area": shape.width * shape.height,
            }
        )

    logo_candidates = [item for item in pictures if item["left"] < 250 * EMU_PER_PIXEL and item["top"] > 80 * EMU_PER_PIXEL]
    site_candidates = [item for item in pictures if item["left"] > 450 * EMU_PER_PIXEL and item["top"] > 80 * EMU_PER_PIXEL]

    logo = None
    if logo_candidates:
        logo_pick = min(logo_candidates, key=lambda item: (item["left"], item["top"]))
        replaceable = emu_to_px(logo_pick["height"]) >= 20 and emu_to_px(logo_pick["width"]) >= 60
        if replaceable:
            logo = {
                "mode": "replace",
                "target_left": emu_to_px(logo_pick["left"]),
                "target_top": emu_to_px(logo_pick["top"]),
            }
        else:
            logo = {
                "mode": "add",
                "left": emu_to_px(logo_pick["left"]),
                "top": emu_to_px(max(logo_pick["top"] - 4 * EMU_PER_PIXEL, 0)),
                "width": 145.0,
                "height": 35.0,
                "clear_region": {
                    "left": emu_to_px(max(logo_pick["left"] - 10 * EMU_PER_PIXEL, 0)),
                    "top": emu_to_px(max(logo_pick["top"] - 20 * EMU_PER_PIXEL, 0)),
                    "right": emu_to_px(logo_pick["left"] + 140 * EMU_PER_PIXEL),
                    "bottom": emu_to_px(logo_pick["top"] + 35 * EMU_PER_PIXEL),
                },
            }

    site = None
    if site_candidates:
        site_pick = max(site_candidates, key=lambda item: item["area"])
        site = {
            "mode": "replace",
            "target_left": emu_to_px(site_pick["left"]),
            "target_top": emu_to_px(site_pick["top"]),
            "fill": True,
        }

    # Many manually designed buyer boards leave image regions blank instead of
    # inserting picture placeholders. Derive stable insertion boxes from the
    # table bounds so those templates can still receive fetched assets.
    if table_shape is not None and (logo is None or site is None):
        table_left = table_shape.left
        table_top = table_shape.top
        table_right = table_shape.left + table_shape.width
        table_bottom = table_shape.top + table_shape.height
        if logo is None and table_top > 100 * EMU_PER_PIXEL:
            logo = {
                "mode": "add",
                "left": emu_to_px(table_left),
                "top": emu_to_px(max(table_top - 44 * EMU_PER_PIXEL, 0)),
                "width": 172.0,
                "height": 34.0,
            }
        if site is None and slide_width and slide_height and table_right < slide_width - 240 * EMU_PER_PIXEL:
            site_left = table_right + 12 * EMU_PER_PIXEL
            site_top = table_top
            site_right = slide_width - 36 * EMU_PER_PIXEL
            site_bottom = min(table_bottom, slide_height - 80 * EMU_PER_PIXEL)
            if site_right > site_left and site_bottom > site_top:
                site = {
                    "mode": "add",
                    "left": emu_to_px(site_left),
                    "top": emu_to_px(site_top),
                    "width": emu_to_px(site_right - site_left),
                    "height": emu_to_px(site_bottom - site_top),
                    "fill": True,
                }

    return {"logo": logo, "site": site}


def build_content_config(presentation: Presentation) -> tuple[dict[str, Any], list[dict[str, Any]]]:
    content_slides = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        table_shape = next((shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.TABLE), None)
        if table_shape is None:
            continue
        content_slides.append((slide_index, slide, table_shape))

    if not content_slides:
        raise ValueError("No content slide with a table was found.")

    first_slide_index, first_slide, first_table = content_slides[0]
    table_shape_index = list(first_slide.shapes).index(first_table) + 1
    title_shape_index = choose_title_shape(first_slide, table_shape_index)
    footer_shape_index = choose_footer_shape(first_slide, title_shape_index)

    fields = []
    for row_index in range(len(first_table.table.rows)):
        label = first_table.table.cell(row_index, 0).text.strip()
        fields.append(
            {
                "row": row_index,
                "label": label,
                "value_key": infer_value_key(label),
                "value_column": 1,
                "label_column": 0,
                "label_color": "2A49F4",
                "value_color": "2A49F4",
            }
        )

    image_slots = []
    for slide_index, slide, table in content_slides:
        image_slots.append(
            {
                "slide_offset": slide_index - first_slide_index,
                **detect_image_slots(slide, table, presentation.slide_width, presentation.slide_height),
            }
        )

    content = {
        "source_slide_index": first_slide_index,
        "start_slide_index": first_slide_index,
        "template_slide_count": len(content_slides),
        "title_shape_index": title_shape_index,
        "table_shape_index": table_shape_index,
        "fields": fields,
        "preserve_title": True,
        "preserve_footer": True,
        "dynamic_row_height": True,
    }
    if footer_shape_index is not None:
        content["footer_shape_index"] = footer_shape_index
    return content, image_slots


def main() -> int:
    parser = argparse.ArgumentParser(description="Generate a starter layout-config.json from a buyer-board reference PPT.")
    parser.add_argument("--template", required=True, help="Reference PPTX path")
    parser.add_argument("--output", required=True, help="Output layout-config.json path")
    parser.add_argument("--cover-title", default="请替换封面标题", help="Default cover title written into config")
    parser.add_argument("--cover-country", default="国家：请替换", help="Default cover country line written into config")
    parser.add_argument("--content-title", default="请替换内容页标题", help="Default content title written into config")
    args = parser.parse_args()

    presentation = Presentation(args.template)
    cover = build_cover_config(presentation.slides[0])
    content, image_slots = build_content_config(presentation)

    config = {
        "version": 1,
        "defaults": {
            "cover_title": args.cover_title,
            "cover_country": args.cover_country,
            "content_title": args.content_title,
        },
        "cover": cover,
        "content": content,
        "images": {
            "layout_version": 2,
            "slides": image_slots,
        },
        "notes": [
            "This file is a scaffold generated from a reference PPT.",
            "Verify image slots, especially logo add/replace mode, before mass production.",
        ],
    }

    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
