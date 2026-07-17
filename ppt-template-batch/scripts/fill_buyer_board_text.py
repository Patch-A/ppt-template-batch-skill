from __future__ import annotations

import argparse
import copy
import json
import math
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches


EMU_PER_INCH = 914400
EMU_PER_POINT = 12700


def load_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8-sig"))


def get_shape(slide, shape_index: int):
    return slide.shapes[shape_index - 1]


def get_first_run(paragraph):
    if paragraph.runs:
        return paragraph.runs[0]
    return None


def copy_font(source_run, target_run) -> None:
    if source_run is None:
        return

    target_run.font.name = source_run.font.name
    target_run.font.size = source_run.font.size
    target_run.font.bold = source_run.font.bold
    target_run.font.italic = source_run.font.italic
    target_run.font.underline = source_run.font.underline

    try:
        target_run.font.color.rgb = source_run.font.color.rgb
    except Exception:
        try:
            target_run.font.color.theme_color = source_run.font.color.theme_color
        except Exception:
            pass


def safe_replace_text_frame(text_frame, value: str) -> bool:
    """Replace text inside the original run so template effects stay intact."""
    if not text_frame.paragraphs:
        return False

    first_run = get_first_run(text_frame.paragraphs[0])
    if first_run is None:
        first_run = text_frame.paragraphs[0].add_run()

    for paragraph_index, paragraph in enumerate(text_frame.paragraphs):
        for run_index, run in enumerate(paragraph.runs):
            if paragraph_index == 0 and run_index == 0:
                continue
            run.text = ""
    first_run.text = value
    return True


def safe_replace_shape_text(shape, value: str) -> bool:
    """Replace only the text of a title/footer/fixed text shape."""
    if not getattr(shape, "has_text_frame", False):
        return False
    return safe_replace_text_frame(shape.text_frame, value)


def replace_text(text_frame, value: str) -> None:
    if safe_replace_text_frame(text_frame, value):
        return
    text_frame.clear()
    text_frame.paragraphs[0].add_run().text = value


def set_shape_text(slide, shape_index: int, value: str) -> None:
    shape = get_shape(slide, shape_index)
    if not safe_replace_shape_text(shape, value):
        raise ValueError(f"Shape {shape_index} does not have a usable text frame")


def set_table_cell(table, row: int, col: int, value: str) -> None:
    cell = table.cell(row, col)
    replace_text(cell.text_frame, value)


def apply_override_color(shape_or_cell, rgb_hex: str) -> None:
    text_frame = shape_or_cell.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor.from_string(rgb_hex)


def cell_font_size_pt(cell, default: float = 16.0) -> float:
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                return float(run.font.size.pt)
    return default


def weighted_text_length(text: str) -> float:
    total = 0.0
    for char in text:
        if char == "\n":
            continue
        if "\u4e00" <= char <= "\u9fff" or char in "\uff0c\u3002\u3001\uff1b\uff1a\uff08\uff09":
            total += 1.0
        elif char.isspace():
            total += 0.3
        else:
            total += 0.55
    return total


def estimate_text_lines(cell, column_width: int, value: str) -> int:
    font_size = cell_font_size_pt(cell)
    usable_width = max(
        Inches(0.5),
        column_width - int(cell.margin_left or 0) - int(cell.margin_right or 0),
    )
    capacity = max(8.0, (usable_width / EMU_PER_INCH) * 72.0 / (font_size * 0.95))
    logical_lines = str(value or "").splitlines() or [""]
    return max(1, sum(max(1, math.ceil(weighted_text_length(line) / capacity)) for line in logical_lines))


def content_row_height(cell, column_width: int, value: str, value_key: str) -> int:
    lines = estimate_text_lines(cell, column_width, value)
    font_size = cell_font_size_pt(cell)
    vertical_margin_pt = (int(cell.margin_top or 0) + int(cell.margin_bottom or 0)) / EMU_PER_POINT
    height_inches = (lines * font_size * 1.25 + vertical_margin_pt + 2.0) / 72.0
    minimum = 0.42 if value_key == "products" else 0.78
    maximum = 1.45 if value_key == "products" else 2.6
    return Inches(max(minimum, min(maximum, height_inches)))


def adjust_dynamic_row_heights(table_shape, fields: list[dict[str, Any]], buyer: dict[str, Any], enabled: bool) -> None:
    if not enabled:
        return
    table = table_shape.table
    for field in fields:
        value_key = str(field.get("value_key", ""))
        if value_key not in {"products", "bio"}:
            continue
        row_index = int(field["row"])
        value_col = int(field.get("value_column", 1))
        value = str(buyer.get(value_key, "") or "")
        cell = table.cell(row_index, value_col)
        cell.text_frame.word_wrap = True
        table.rows[row_index].height = content_row_height(
            cell,
            table.columns[value_col].width,
            value,
            value_key,
        )
    table_shape.height = sum(int(row.height) for row in table.rows)


def duplicate_slide(presentation: Presentation, source_index: int):
    source_slide = presentation.slides[source_index]
    new_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    for shape in source_slide.shapes:
        new_slide.shapes._spTree.insert_element_before(copy.deepcopy(shape.element), "p:extLst")
    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype or rel.rId in new_slide.part.rels:
            continue
        try:
            new_slide.part.rels.add_relationship(rel.reltype, rel._target, rel.rId)
        except Exception:
            pass
    return new_slide


def move_slide(presentation: Presentation, old_index: int, new_index: int) -> None:
    slide_ids = presentation.slides._sldIdLst
    slide_id = list(slide_ids)[old_index]
    slide_ids.remove(slide_id)
    slide_ids.insert(new_index, slide_id)


def delete_slide(presentation: Presentation, slide_index: int) -> None:
    slide_ids = presentation.slides._sldIdLst
    slide_id = list(slide_ids)[slide_index]
    rel_id = slide_id.rId
    slide_ids.remove(slide_id)
    presentation.part.drop_rel(rel_id)


def resize_content_slides(presentation: Presentation, config: dict[str, Any], buyer_count: int) -> None:
    content_config = config["content"]
    start_index = int(content_config["start_slide_index"]) - 1
    source_index = int(content_config.get("source_slide_index", content_config["start_slide_index"])) - 1
    template_count = int(content_config.get("template_slide_count") or len(presentation.slides) - start_index)
    if template_count < 1:
        raise ValueError("Buyer-board template must contain at least one content slide.")

    if buyer_count > template_count:
        for offset in range(buyer_count - template_count):
            duplicate_slide(presentation, source_index)
            move_slide(presentation, len(presentation.slides) - 1, start_index + template_count + offset)
    elif template_count > buyer_count:
        for _ in range(template_count - buyer_count):
            delete_slide(presentation, start_index + buyer_count)


def fill_cover(presentation: Presentation, config: dict[str, Any], cover_title: str, cover_country: str) -> None:
    cover_config = config["cover"]
    slide = presentation.slides[cover_config["slide_index"] - 1]
    set_shape_text(slide, cover_config["title_shape_index"], cover_title)
    set_shape_text(slide, cover_config["country_shape_index"], cover_country)


def fill_content_slides(
    presentation: Presentation,
    buyers: list[dict[str, Any]],
    config: dict[str, Any],
    content_title: str,
) -> None:
    content_config = config["content"]
    start_slide_index = content_config["start_slide_index"]
    table_shape_index = content_config["table_shape_index"]
    title_shape_index = content_config["title_shape_index"]
    fields = content_config["fields"]
    preserve_title = content_config.get("preserve_title", True)

    available_content_slides = len(presentation.slides) - start_slide_index + 1
    if len(buyers) > available_content_slides:
        raise ValueError(
            f"Template only has {available_content_slides} content slides, but buyers.json contains {len(buyers)} buyers."
        )

    for offset, buyer in enumerate(buyers):
        slide = presentation.slides[start_slide_index - 1 + offset]
        if not preserve_title:
            set_shape_text(slide, title_shape_index, content_title)

        table_shape = get_shape(slide, table_shape_index)
        table = table_shape.table
        for field in fields:
            row = int(field["row"])
            label_col = int(field.get("label_column", 0))
            value_col = int(field.get("value_column", 1))
            label = str(field.get("label", ""))
            value_key = field["value_key"]
            value = str(buyer.get(value_key, "") or "")

            if label:
                set_table_cell(table, row, label_col, label)
            set_table_cell(table, row, value_col, value)

            if content_config.get("allow_style_overrides", False):
                if "label_color" in field:
                    apply_override_color(table.cell(row, label_col), field["label_color"])
                if "value_color" in field:
                    apply_override_color(table.cell(row, value_col), field["value_color"])

        adjust_dynamic_row_heights(
            table_shape,
            fields,
            buyer,
            bool(content_config.get("dynamic_row_height", True)),
        )


def main() -> int:
    parser = argparse.ArgumentParser(description="Fill buyer-board text content using buyers.json and layout-config.json.")
    parser.add_argument("template", help="Template PPTX path")
    parser.add_argument("buyers", help="Buyers JSON path")
    parser.add_argument("layout_config", help="Layout config JSON path")
    parser.add_argument("output", help="Output PPTX path")
    parser.add_argument("--cover-title", dest="cover_title", help="Override cover title")
    parser.add_argument("--cover-country", dest="cover_country", help="Override cover country line")
    parser.add_argument("--content-title", dest="content_title", help="Override content-page title")
    args = parser.parse_args()

    buyers = load_json(Path(args.buyers))
    config = load_json(Path(args.layout_config))
    presentation = Presentation(args.template)
    resize_content_slides(presentation, config, len(buyers))

    defaults = config.get("defaults", {})
    cover_title = args.cover_title or defaults.get("cover_title", "")
    cover_country = args.cover_country or defaults.get("cover_country", "")
    content_title = args.content_title or defaults.get("content_title", "")

    fill_cover(presentation, config, cover_title, cover_country)
    fill_content_slides(presentation, buyers, config, content_title)

    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
