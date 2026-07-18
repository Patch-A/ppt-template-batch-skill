from __future__ import annotations

import argparse
import copy
import json
import math
from pathlib import Path
from typing import Any

from pptx import Presentation


DEFAULT_MAPPING: dict[str, Any] = {
    "title_shape": 5,
    "buyers_per_slide": 6,
    "slots": [
        {"summary_shape": 15, "products_shape": 23},
        {"summary_group": 16, "summary_child": 2, "products_group": 16, "products_child": 3},
        {"summary_group": 17, "summary_child": 2, "products_group": 17, "products_child": 3},
        {"summary_shape": 19, "products_shape": 26},
        {"summary_shape": 20, "products_shape": 25},
        {"summary_shape": 22, "products_shape": 24},
    ],
}
MAX_BUYERS_PER_SLIDE = 6
EMU_PER_INCH = 914400


def load_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8-sig"))


def ensure_runs(paragraph, minimum: int) -> None:
    while len(paragraph.runs) < minimum:
        paragraph.add_run()


def get_text_frame(slide, shape_index: int, child_index: int | None = None):
    shape = slide.shapes[shape_index - 1]
    if child_index is not None:
        shape = shape.shapes[child_index - 1]
    return shape.text_frame


def set_single_run_text(text_frame, value: str) -> None:
    clear_text_frame(text_frame)
    paragraph = text_frame.paragraphs[0]
    ensure_runs(paragraph, 1)
    paragraph.runs[0].text = value


def clear_text_frame(text_frame) -> None:
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = ""


def split_summary(name: str, summary: str) -> tuple[str, str]:
    if summary.startswith(name):
        return name, summary[len(name) :]
    return name, summary


def set_summary_text(text_frame, buyer: dict[str, Any]) -> None:
    name = str(buyer.get("name", "") or "")
    summary = str(buyer.get("summary") or buyer.get("intro") or "")
    paragraph = text_frame.paragraphs[0]
    clear_text_frame(text_frame)
    ensure_runs(paragraph, 2)
    name_text, body_text = split_summary(name, summary)
    paragraph.runs[0].text = name_text
    paragraph.runs[1].text = body_text
    for run in paragraph.runs[2:]:
        run.text = ""


def _font_size_pt(text_frame, default: float = 14.0) -> float:
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                return float(run.font.size.pt)
    return default


def _weighted_text_length(value: str) -> float:
    total = 0.0
    for char in value:
        if char == "\n":
            continue
        if "\u4e00" <= char <= "\u9fff":
            total += 1.0
        elif char.isspace():
            total += 0.3
        else:
            total += 0.55
    return total


def text_capacity(text_frame) -> int:
    """Estimate a text frame's capacity from its geometry and existing font."""
    shape = getattr(text_frame, "_parent", None)
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    if not width or not height:
        return max(len(getattr(text_frame, "text", "") or ""), 1)

    font_size = max(_font_size_pt(text_frame), 1.0)
    width_pt = max(width / EMU_PER_INCH * 72.0, 1.0)
    height_pt = max(height / EMU_PER_INCH * 72.0, font_size)
    lines = max(1, math.floor(height_pt / (font_size * 1.2)))
    chars_per_line = max(1.0, width_pt / font_size * 1.6)
    return max(1, math.floor(chars_per_line * lines))


def _capacity_warning(
    text_frame,
    value: str,
    page_number: int,
    slot_number: int,
    field: str,
) -> dict[str, Any] | None:
    if not value:
        return None
    capacity = text_capacity(text_frame)
    length = _weighted_text_length(value)
    if length <= capacity:
        return None
    return {
        "page": page_number,
        "slot": slot_number,
        "field": field,
        "length": len(value),
        "capacity": capacity,
    }


def products_text(buyer: dict[str, Any]) -> str:
    value = str(buyer.get("products") or buyer.get("categories") or "").strip()
    if value.startswith("采购品类："):
        return value
    return f"采购品类：{value}" if value else ""


def get_slot_frames(slide, slot: dict[str, Any]):
    if "summary_group" in slot:
        summary_frame = get_text_frame(slide, int(slot["summary_group"]), int(slot["summary_child"]))
    else:
        summary_frame = get_text_frame(slide, int(slot["summary_shape"]))

    if "products_group" in slot:
        products_frame = get_text_frame(slide, int(slot["products_group"]), int(slot["products_child"]))
    else:
        products_frame = get_text_frame(slide, int(slot["products_shape"]))
    return summary_frame, products_frame


def fill_slot(
    slide,
    slot: dict[str, Any],
    buyer: dict[str, Any],
    page_number: int = 1,
    slot_number: int = 1,
) -> list[dict[str, Any]]:
    summary_frame, products_frame = get_slot_frames(slide, slot)
    name = str(buyer.get("name", "") or "")
    summary_value = str(buyer.get("summary") or buyer.get("intro") or "")
    if summary_value and not summary_value.startswith(name):
        summary_value = f"{name}{summary_value}"
    product_value = products_text(buyer)
    set_summary_text(summary_frame, buyer)
    set_single_run_text(products_frame, product_value)

    warnings = []
    for field, frame, value in (
        ("summary", summary_frame, summary_value),
        ("products", products_frame, product_value),
    ):
        warning = _capacity_warning(frame, value, page_number, slot_number, field)
        if warning:
            warnings.append(warning)
    return warnings


def clear_slot(slide, slot: dict[str, Any]) -> None:
    summary_frame, products_frame = get_slot_frames(slide, slot)
    set_single_run_text(summary_frame, "")
    set_single_run_text(products_frame, "")


def fill_slide(
    slide,
    page: dict[str, Any],
    mapping: dict[str, Any],
    page_number: int = 1,
) -> dict[str, list[dict[str, Any]]]:
    buyers = list(page.get("buyers") or [])
    slots = list(mapping.get("slots") or [])
    buyers_per_slide = min(
        MAX_BUYERS_PER_SLIDE,
        int(mapping.get("buyers_per_slide", len(slots))),
        len(slots),
    )
    if len(buyers) > buyers_per_slide:
        raise ValueError(f"Each briefing slide supports at most {buyers_per_slide} buyers.")

    title_frame = get_text_frame(slide, int(mapping["title_shape"]))
    set_single_run_text(title_frame, str(page.get("title", "") or ""))

    report: dict[str, list[dict[str, Any]]] = {
        "missing_buyers": [],
        "overlong_text": [],
    }

    for slot_index, slot in enumerate(slots):
        if slot_index >= buyers_per_slide:
            clear_slot(slide, slot)
            continue
        buyer = buyers[slot_index] if slot_index < len(buyers) else {}
        if not str(buyer.get("name", "") or "").strip():
            clear_slot(slide, slot)
            report["missing_buyers"].append({"page": page_number, "slot": slot_index + 1})
            continue
        report["overlong_text"].extend(
            fill_slot(slide, slot, buyer, page_number, slot_index + 1)
        )
    return report


def duplicate_slide(presentation: Presentation, source_index: int):
    source_slide = presentation.slides[source_index]
    blank_layout = presentation.slide_layouts[6]
    new_slide = presentation.slides.add_slide(blank_layout)
    for shape in source_slide.shapes:
        new_slide.shapes._spTree.insert_element_before(copy.deepcopy(shape.element), "p:extLst")
    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        if rel.rId in new_slide.part.rels:
            continue
        try:
            new_slide.part.rels.add_relationship(rel.reltype, rel._target, rel.rId)
        except Exception:
            pass
    return new_slide


def load_pages(path: Path) -> tuple[list[dict[str, Any]], dict[str, Any] | None]:
    payload = load_json(path)
    if isinstance(payload, list):
        return payload, None
    if isinstance(payload, dict):
        pages = payload.get("pages", [])
        if not isinstance(pages, list):
            raise ValueError("Buyer briefing JSON field 'pages' must be a list.")
        embedded_mapping = payload.get("mapping")
        if embedded_mapping is not None and not isinstance(embedded_mapping, dict):
            raise ValueError("Buyer briefing JSON field 'mapping' must be an object.")
        return pages, embedded_mapping
    raise ValueError("Buyer briefing JSON must be a list or an object containing 'pages'.")


def fill_presentation(
    template_path: str | Path,
    pages_path: str | Path,
    output_path: str | Path,
    report_path: str | Path | None = None,
    mapping: dict[str, Any] | None = None,
) -> dict[str, Any]:
    pages, embedded_mapping = load_pages(Path(pages_path))
    active_mapping = mapping or embedded_mapping or DEFAULT_MAPPING
    slots = list(active_mapping.get("slots") or [])
    buyers_per_slide = min(
        MAX_BUYERS_PER_SLIDE,
        int(active_mapping.get("buyers_per_slide", len(slots))),
        len(slots),
    )
    for page_number, page in enumerate(pages, start=1):
        buyers = list(page.get("buyers") or [])
        if len(buyers) > buyers_per_slide:
            raise ValueError(
                f"Page {page_number} supports at most {buyers_per_slide} buyers; "
                f"received {len(buyers)}."
            )

    prs = Presentation(str(template_path))
    if pages and not prs.slides:
        raise ValueError("The briefing template has no slides.")
    while len(pages) > len(prs.slides):
        duplicate_slide(prs, len(prs.slides) - 1)

    report: dict[str, Any] = {
        "ok": True,
        "pipeline_mode": "buyer_briefing",
        "page_count": len(pages),
        "output": str(output_path),
        "missing_buyers": [],
        "overlong_text": [],
        "warnings": [],
    }
    for index, page in enumerate(pages):
        page_report = fill_slide(prs.slides[index], page, active_mapping, index + 1)
        report["missing_buyers"].extend(page_report["missing_buyers"])
        report["overlong_text"].extend(page_report["overlong_text"])
    report["warnings"] = [
        {"type": "missing_buyer", **item} for item in report["missing_buyers"]
    ] + [
        {"type": "overlong_text", **item} for item in report["overlong_text"]
    ]

    target = Path(output_path)
    target.parent.mkdir(parents=True, exist_ok=True)
    prs.save(target)
    if report_path:
        report_target = Path(report_path)
        report_target.parent.mkdir(parents=True, exist_ok=True)
        report_target.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    return report


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Fill buyer-briefing PPT pages while preserving template run-level text styles."
    )
    parser.add_argument("template", help="Template PPTX path")
    parser.add_argument("pages_json", help="JSON with pages: title + up to 6 buyers per page")
    parser.add_argument("output", help="Output PPTX path")
    parser.add_argument("--layout-config", help="Optional briefing layout mapping JSON")
    parser.add_argument("--report", help="Optional export report JSON path")
    args = parser.parse_args()

    mapping = None
    if args.layout_config:
        mapping = load_json(Path(args.layout_config))
    fill_presentation(args.template, args.pages_json, args.output, args.report, mapping)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
