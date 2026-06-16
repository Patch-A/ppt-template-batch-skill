from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt


WHITE = RGBColor(255, 255, 255)
BODY_BLUE = RGBColor(42, 73, 244)


def write_shape(shape, text: str, size: float, bold: bool, color: RGBColor, font_name: str = "微软雅黑") -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def write_cell(cell, text: str, size: float = 16, font_name: str = "微软雅黑") -> None:
    text_frame = cell.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = BODY_BLUE
    run.font.name = font_name


def load_buyers(path: Path) -> list[dict]:
    return json.loads(path.read_text(encoding="utf-8"))


def fill_deck(template_path: Path, buyers_path: Path, output_path: Path, cover_title: str, cover_country: str, content_title: str) -> None:
    prs = Presentation(str(template_path))
    buyers = load_buyers(buyers_path)

    write_shape(prs.slides[0].shapes[3], cover_title, size=54, bold=True, color=WHITE)
    write_shape(prs.slides[0].shapes[5], cover_country, size=32, bold=True, color=WHITE)

    for slide, buyer in zip(list(prs.slides)[1:], buyers):
        write_shape(slide.shapes[3], content_title, size=32, bold=True, color=WHITE)
        table = slide.shapes[4].table
        rows = [
            ("企业", buyer["name"]),
            ("国家", buyer["country"]),
            ("网站", buyer["website"]),
            ("采购产品", buyer["products"]),
            ("简介", buyer["bio"]),
        ]
        for row_idx, (label, value) in enumerate(rows):
            write_cell(table.cell(row_idx, 0), label)
            write_cell(table.cell(row_idx, 1), value)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def main() -> int:
    if len(sys.argv) != 7:
        print(
            "Usage: python fill_buyer_board_text.py <template.pptx> <buyers.json> <output.pptx> <cover_title> <cover_country> <content_title>",
            file=sys.stderr,
        )
        return 1

    template_path = Path(sys.argv[1])
    buyers_path = Path(sys.argv[2])
    output_path = Path(sys.argv[3])
    cover_title = sys.argv[4]
    cover_country = sys.argv[5]
    content_title = sys.argv[6]

    fill_deck(template_path, buyers_path, output_path, cover_title, cover_country, content_title)
    print(output_path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
